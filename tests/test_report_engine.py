"""Tests for the Report Engine: persistence, PPTX chart rendering, auto-generate."""

from __future__ import annotations

import io

import pandas as pd
import pytest

from services.report_engine import (
    ChartDefinition,
    ChartType,
    Insight,
    KPIMetric,
    Recommendation,
    ReportCompositionService,
    ReportSection,
    ReportSectionType,
    ReportTemplate,
)

# ── Fixtures ────────────────────────────────────────────────────────────


@pytest.fixture
def sample_kpis():
    return [
        KPIMetric(label="Revenue", value="1.2M", unit="$", trend="up", trend_value="+15%"),
        KPIMetric(label="Users", value="45K", unit="", trend="up", trend_value="+8%"),
    ]


@pytest.fixture
def sample_chart():
    return ChartDefinition(
        title="Monthly Revenue",
        chart_type=ChartType.BAR,
        data=[
            {"x": "Jan", "y": 100},
            {"x": "Feb", "y": 150},
            {"x": "Mar", "y": 200},
            {"x": "Apr", "y": 180},
        ],
        x_axis="month",
        y_axis="revenue",
    )


@pytest.fixture
def sample_insights():
    return [
        Insight(
            title="Revenue Growth",
            description="Revenue increased 15% QoQ",
            severity="positive",
            impact="High",
        ),
    ]


@pytest.fixture
def sample_recommendations():
    return [
        Recommendation(
            title="Expand Marketing",
            description="Increase marketing spend in Q3",
            priority="high",
        ),
    ]


@pytest.fixture
def report_with_data(sample_kpis, sample_chart, sample_insights, sample_recommendations):
    """Create a report populated with KPIs, charts, insights, and recommendations."""
    report = ReportCompositionService.create_report(
        title="Test Report",
        template=ReportTemplate.EXECUTIVE,
        org_name="Test Org",
        author="Test Author",
        industry="retail",
    )
    # Add KPIs to the key metrics section
    for section in report.sections:
        if section.section_type == ReportSectionType.KEY_METRICS:
            ReportCompositionService.add_kpis(report.report_id, section.order, sample_kpis)
            break
    # Add chart to the chart section
    for section in report.sections:
        if section.section_type == ReportSectionType.CHART:
            ReportCompositionService.add_chart(report.report_id, section.order, sample_chart)
            break
    # Add insights to the insights section
    for section in report.sections:
        if section.section_type == ReportSectionType.INSIGHTS:
            ReportCompositionService.add_insights(report.report_id, section.order, sample_insights)
            break
    # Add recommendations
    for section in report.sections:
        if section.section_type == ReportSectionType.RECOMMENDATIONS:
            ReportCompositionService.add_recommendations(
                report.report_id, section.order, sample_recommendations
            )
            break
    return ReportCompositionService.get_report(report.report_id)


# ── 1. Report CRUD Tests ────────────────────────────────────────────────


class TestReportCRUD:
    """Test report creation, retrieval, listing, and deletion."""

    def test_create_report(self):
        report = ReportCompositionService.create_report(
            title="My Report",
            template=ReportTemplate.EXECUTIVE,
            org_name="ACME",
            author="John",
        )
        assert report.report_id
        assert report.title == "My Report"
        assert report.organization_name == "ACME"
        assert len(report.sections) > 0

    def test_get_report(self):
        created = ReportCompositionService.create_report(title="Get Test")
        fetched = ReportCompositionService.get_report(created.report_id)
        assert fetched is not None
        assert fetched.report_id == created.report_id

    def test_get_nonexistent_report(self):
        assert ReportCompositionService.get_report("nonexistent-id-12345") is None

    def test_list_reports(self):
        r1 = ReportCompositionService.create_report(title="List Test 1")
        r2 = ReportCompositionService.create_report(title="List Test 2")
        reports = ReportCompositionService.list_reports()
        report_ids = [r["report_id"] for r in reports]
        assert r1.report_id in report_ids
        assert r2.report_id in report_ids

    def test_delete_report(self):
        report = ReportCompositionService.create_report(title="Delete Me")
        assert ReportCompositionService.delete_report(report.report_id) is True
        assert ReportCompositionService.get_report(report.report_id) is None

    def test_delete_nonexistent_report(self):
        assert ReportCompositionService.delete_report("nonexistent-id-67890") is False


# ── 2. Section Management Tests ─────────────────────────────────────────


class TestSectionManagement:
    """Test adding, updating, and removing sections."""

    def test_add_section(self):
        report = ReportCompositionService.create_report(title="Section Test")
        initial_count = len(report.sections)
        new_section = ReportSection(
            section_type=ReportSectionType.CUSTOM,
            title="Custom Section",
            content="Custom content",
        )
        result = ReportCompositionService.add_section(report.report_id, new_section)
        assert result is not None
        assert len(result.sections) == initial_count + 1
        assert result.sections[-1].title == "Custom Section"

    def test_remove_section(self):
        report = ReportCompositionService.create_report(title="Remove Section Test")
        initial_count = len(report.sections)
        result = ReportCompositionService.remove_section(report.report_id, 0)
        assert result is not None
        assert len(result.sections) == initial_count - 1

    def test_add_kpis_to_section(self, sample_kpis):
        report = ReportCompositionService.create_report(title="KPI Test")
        kpi_section = next(
            s for s in report.sections if s.section_type == ReportSectionType.KEY_METRICS
        )
        result = ReportCompositionService.add_kpis(report.report_id, kpi_section.order, sample_kpis)
        assert result is not None
        updated_section = next(s for s in result.sections if s.order == kpi_section.order)
        assert len(updated_section.kpis) >= 2

    def test_add_chart_to_section(self, sample_chart):
        report = ReportCompositionService.create_report(title="Chart Test")
        chart_section = next(
            s for s in report.sections if s.section_type == ReportSectionType.CHART
        )
        result = ReportCompositionService.add_chart(
            report.report_id, chart_section.order, sample_chart
        )
        assert result is not None
        updated_section = next(s for s in result.sections if s.order == chart_section.order)
        assert len(updated_section.charts) >= 1
        assert updated_section.charts[0].title == "Monthly Revenue"


# ── 3. PPTX Export with Real Charts ─────────────────────────────────────


class TestPPTXExport:
    """Test that PPTX export contains real chart objects."""

    def test_pptx_contains_real_charts(self, report_with_data):
        """PPTX export must contain actual chart objects, not text placeholders."""
        pptx_bytes = ReportCompositionService.export_to_pptx(report_with_data)
        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 0

        # Re-open the PPTX and verify charts are present
        from pptx import Presentation as PptxPresentation

        prs = PptxPresentation(io.BytesIO(pptx_bytes))
        chart_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    chart_count += 1
        assert chart_count > 0, "PPTX must contain at least one real chart object"

    def test_pptx_has_multiple_slides(self, report_with_data):
        """PPTX should have multiple slides (title, KPIs, charts, etc.)."""
        pptx_bytes = ReportCompositionService.export_to_pptx(report_with_data)
        from pptx import Presentation as PptxPresentation

        prs = PptxPresentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) >= 3, "PPTX should have at least 3 slides"

    def test_pptx_chart_has_correct_data(self, report_with_data):
        """The chart in the PPTX should have the correct data from the report."""
        pptx_bytes = ReportCompositionService.export_to_pptx(report_with_data)
        from pptx import Presentation as PptxPresentation

        prs = PptxPresentation(io.BytesIO(pptx_bytes))
        found_chart = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    # Verify the chart has series data
                    for series in chart.series:
                        values = list(series.values)
                        assert len(values) > 0, "Chart series should have values"
                        found_chart = True
        assert found_chart, "Should find at least one chart with data"


# ── 4. Auto-Generate from Dataset ───────────────────────────────────────


class TestAutoGenerate:
    """Test the auto-generate pipeline that creates reports from datasets."""

    def test_auto_generate_with_sales_data(self):
        """Auto-generating a report from sales data should produce charts and KPIs."""
        from services.auto.orchestrator import AutoEngineOrchestrator

        df = pd.DataFrame(
            {
                "date": pd.date_range("2024-01-01", periods=30, freq="D"),
                "product": ["A", "B", "C"] * 10,
                "revenue": [100, 200, 150] * 10,
                "quantity": [10, 20, 15] * 10,
                "region": ["North", "South", "East", "West", "Central"] * 6,
            }
        )

        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(df, dataset_name="sales", industry="retail")
        dashboard_spec = result["dashboard"]

        # Create a report and populate it
        report = ReportCompositionService.create_report(
            title="Auto-Generated Sales Report",
            template=ReportTemplate.EXECUTIVE,
            industry="retail",
        )
        ReportCompositionService.populate_from_dashboard_spec(report.report_id, dashboard_spec)

        # Verify the report has content
        populated = ReportCompositionService.get_report(report.report_id)
        assert populated is not None
        chart_count = sum(len(s.charts) for s in populated.sections)
        kpi_count = sum(len(s.kpis) for s in populated.sections)
        insight_count = sum(len(s.insights) for s in populated.sections)

        assert chart_count > 0, "Auto-generated report should have charts"
        assert kpi_count > 0, "Auto-generated report should have KPIs"
        assert insight_count > 0, "Auto-generated report should have insights"

    def test_auto_generate_charts_have_data(self):
        """Charts from auto-generate should have actual data, not empty arrays."""
        from services.auto.orchestrator import AutoEngineOrchestrator

        df = pd.DataFrame(
            {
                "category": ["A", "B", "C", "D"] * 10,
                "value": range(40),
                "date": pd.date_range("2024-01-01", periods=40, freq="D"),
            }
        )

        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(df, dataset_name="test", industry="unknown")
        dashboard_spec = result["dashboard"]

        report = ReportCompositionService.create_report(
            title="Chart Data Test",
            template=ReportTemplate.EXECUTIVE,
        )
        ReportCompositionService.populate_from_dashboard_spec(report.report_id, dashboard_spec)

        populated = ReportCompositionService.get_report(report.report_id)
        assert populated is not None

        for section in populated.sections:
            for chart in section.charts:
                assert chart.data is not None, f"Chart '{chart.title}' should have data"
                assert len(chart.data) > 0, f"Chart '{chart.title}' should have non-empty data"

    def test_auto_generate_pptx_has_charts(self):
        """Full pipeline: auto-generate → report → PPTX with real charts."""
        from pptx import Presentation as PptxPresentation

        from services.auto.orchestrator import AutoEngineOrchestrator

        df = pd.DataFrame(
            {
                "product": ["A", "B", "C"] * 10,
                "revenue": [100, 200, 150] * 10,
                "date": pd.date_range("2024-01-01", periods=30, freq="D"),
            }
        )

        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(df, dataset_name="sales", industry="retail")
        dashboard_spec = result["dashboard"]

        report = ReportCompositionService.create_report(
            title="Full Pipeline PPTX Test",
            template=ReportTemplate.EXECUTIVE,
        )
        ReportCompositionService.populate_from_dashboard_spec(report.report_id, dashboard_spec)

        populated = ReportCompositionService.get_report(report.report_id)
        assert populated is not None

        pptx_bytes = ReportCompositionService.export_to_pptx(populated)
        prs = PptxPresentation(io.BytesIO(pptx_bytes))

        chart_count = sum(1 for slide in prs.slides for shape in slide.shapes if shape.has_chart)
        assert chart_count > 0, "Auto-generated PPTX should contain real charts"


# ── 5. Export Format Tests ──────────────────────────────────────────────


class TestExportFormats:
    """Test all export formats produce valid output."""

    def test_export_html(self, report_with_data):
        html = ReportCompositionService.export_to_html(report_with_data)
        assert isinstance(html, str)
        assert "<html" in html.lower() or "<div" in html.lower()
        assert report_with_data.title in html

    def test_export_pdf(self, report_with_data):
        pdf = ReportCompositionService.export_to_pdf(report_with_data)
        assert isinstance(pdf, bytes)
        assert len(pdf) > 0
        assert pdf.startswith(b"%PDF")

    def test_export_json(self, report_with_data):
        import json

        json_dict = ReportCompositionService.export_to_dict(report_with_data)
        json_str = json.dumps(json_dict)
        assert report_with_data.title in json_str
        assert "sections" in json_dict

    def test_export_pptx(self, report_with_data):
        pptx = ReportCompositionService.export_to_pptx(report_with_data)
        assert isinstance(pptx, bytes)
        assert len(pptx) > 0
        # PPTX files start with PK (zip magic bytes)
        assert pptx[:2] == b"PK"


# ── 6. Presentation Generator Tests ─────────────────────────────────────


class TestPresentationGenerator:
    """Test the PresentationGenerator produces valid slide definitions."""

    def test_presentation_has_slides(self, report_with_data):
        from services.report_engine import PresentationGenerator

        slides = PresentationGenerator.from_report(report_with_data)
        assert len(slides) > 0
        assert slides[0]["layout"] == "title"

    def test_presentation_chart_slides_have_data(self, report_with_data):
        """Chart slides should include chart_data for frontend rendering."""
        from services.report_engine import PresentationGenerator

        slides = PresentationGenerator.from_report(report_with_data)
        chart_slides = [s for s in slides if s["layout"] == "chart"]
        assert len(chart_slides) > 0
        for slide in chart_slides:
            assert "chart_data" in slide, "Chart slide must include chart_data"
            assert "chart_type" in slide, "Chart slide must include chart_type"
            assert len(slide["chart_data"]) > 0, "Chart slide must have non-empty data"
            assert "x_axis" in slide, "Chart slide must include x_axis"
            assert "y_axis" in slide, "Chart slide must include y_axis"

    def test_presentation_kpi_slides(self, report_with_data):
        from services.report_engine import PresentationGenerator

        slides = PresentationGenerator.from_report(report_with_data)
        kpi_slides = [s for s in slides if s["layout"] == "kpi"]
        assert len(kpi_slides) > 0
        assert len(kpi_slides[0]["kpis"]) > 0
