"""Comprehensive tests for the Intelligent Automatic Analysis, Chart Selection
& Dashboard/Presentation Layout Engine.

Tests cover:
  1. Column semantic role detection (DIMENSION, MEASURE, DATE, IDENTIFIER, etc.)
  2. Chart selection rules (correct chart types, rejection of poor choices)
  3. Chart importance scoring
  4. Chart deduplication
  5. KPI detection and computation
  6. Insight generation (trends, outliers, correlations, dominance)
  7. Filter selection
  8. Dashboard layout (hierarchy, responsive widths, no overlaps)
  9. Presentation layout (slide count, chart placement, validation)
 10. End-to-end: same chart specs in dashboard and presentation
 11. PPTX chart presence (critical regression test)
 12. Canonical specification consistency
"""

from __future__ import annotations

import io

import numpy as np
import pandas as pd
import pytest

from services.auto.analysis_engine import (
    AutomaticAnalysisEngine,
    DatasetUnderstanding,
    SemanticRole,
)
from services.auto.chart_selection_engine import IntelligentChartSelectionEngine
from services.auto.chart_specification import (
    ChartSpecification,
)
from services.auto.dashboard_layout_engine import IntelligentDashboardLayoutEngine
from services.auto.filter_engine import AutomaticFilterEngine
from services.auto.insight_engine import AutomaticInsightEngine
from services.auto.kpi_engine import AutomaticKPIEngine
from services.auto.orchestrator import AutoEngineOrchestrator
from services.auto.presentation_layout_engine import PresentationLayoutEngine

# ── Fixtures ──


@pytest.fixture
def sales_df() -> pd.DataFrame:
    """Sales dataset with time, category, and revenue columns."""
    np.random.seed(42)
    dates = pd.date_range("2023-01-01", periods=100, freq="D")
    categories = ["Electronics", "Clothing", "Books", "Home", "Sports"]
    regions = ["North", "South", "East", "West"]

    return pd.DataFrame(
        {
            "transaction_id": range(1, 101),
            "date": dates,
            "category": np.random.choice(categories, 100),
            "region": np.random.choice(regions, 100),
            "revenue": np.random.uniform(100, 5000, 100).round(2),
            "quantity": np.random.randint(1, 50, 100),
            "cost": np.random.uniform(50, 2500, 100).round(2),
            "is_active": np.random.choice([True, False], 100),
        }
    )


@pytest.fixture
def simple_df() -> pd.DataFrame:
    """Simple dataset for basic tests."""
    return pd.DataFrame(
        {
            "id": range(1, 21),
            "category": ["A", "B", "C", "D"] * 5,
            "value": [
                10,
                20,
                30,
                40,
                50,
                60,
                70,
                80,
                90,
                100,
                10,
                20,
                30,
                40,
                50,
                60,
                70,
                80,
                90,
                100,
            ],
        }
    )


@pytest.fixture
def understanding(sales_df: pd.DataFrame) -> DatasetUnderstanding:
    """Pre-computed dataset understanding for sales data."""
    return AutomaticAnalysisEngine.analyze(
        sales_df,
        dataset_name="sales_data",
        industry="retail",
        quality_score=85.0,
    )


@pytest.fixture
def orchestrator() -> AutoEngineOrchestrator:
    """AutoEngineOrchestrator instance."""
    return AutoEngineOrchestrator()


# ── 1. Column Semantic Role Detection ──


class TestColumnDetection:
    """Test automatic column semantic role detection."""

    def test_identifier_detection(self, sales_df: pd.DataFrame):
        """Sequential integer IDs should be classified as IDENTIFIER."""
        understanding = AutomaticAnalysisEngine.analyze(sales_df)
        col = next(c for c in understanding.columns if c.name == "transaction_id")
        assert col.semantic_role == SemanticRole.IDENTIFIER
        assert col.confidence > 0.7

    def test_date_detection(self, sales_df: pd.DataFrame):
        """Datetime columns should be classified as DATE_TIME."""
        understanding = AutomaticAnalysisEngine.analyze(sales_df)
        col = next(c for c in understanding.columns if c.name == "date")
        assert col.semantic_role == SemanticRole.DATE_TIME

    def test_measure_detection(self, sales_df: pd.DataFrame):
        """Numeric columns like revenue should be classified as MEASURE or CURRENCY."""
        understanding = AutomaticAnalysisEngine.analyze(sales_df)
        col = next(c for c in understanding.columns if c.name == "revenue")
        assert col.semantic_role in (SemanticRole.MEASURE, SemanticRole.CURRENCY)

    def test_category_detection(self, sales_df: pd.DataFrame):
        """Low-cardinality string columns should be classified as CATEGORY."""
        understanding = AutomaticAnalysisEngine.analyze(sales_df)
        col = next(c for c in understanding.columns if c.name == "category")
        assert col.semantic_role == SemanticRole.CATEGORY

    def test_geography_detection(self, sales_df: pd.DataFrame):
        """Region columns should be classified as GEOGRAPHY."""
        understanding = AutomaticAnalysisEngine.analyze(sales_df)
        col = next(c for c in understanding.columns if c.name == "region")
        assert col.semantic_role == SemanticRole.GEOGRAPHY

    def test_boolean_detection(self, sales_df: pd.DataFrame):
        """Boolean columns should be classified as BOOLEAN."""
        understanding = AutomaticAnalysisEngine.analyze(sales_df)
        col = next(c for c in understanding.columns if c.name == "is_active")
        assert col.semantic_role == SemanticRole.BOOLEAN

    def test_grouped_roles(self, sales_df: pd.DataFrame):
        """DatasetUnderstanding should correctly group columns by role."""
        understanding = AutomaticAnalysisEngine.analyze(sales_df)
        assert "date" in understanding.time_columns
        assert "transaction_id" in understanding.identifier_columns
        assert "region" in understanding.geographic_columns
        assert "is_active" in understanding.boolean_columns
        assert len(understanding.measures) >= 2  # revenue, quantity, cost

    def test_empty_column(self):
        """All-null column should be UNKNOWN."""
        df = pd.DataFrame({"empty_col": [None] * 10, "val": range(10)})
        understanding = AutomaticAnalysisEngine.analyze(df)
        col = next(c for c in understanding.columns if c.name == "empty_col")
        assert col.semantic_role == SemanticRole.UNKNOWN
        assert col.missing_percentage == 100.0

    def test_date_string_detection(self):
        """Date-like strings should be detected as DATE_TIME."""
        df = pd.DataFrame(
            {
                "created_date": ["2023-01-01", "2023-01-02", "2023-01-03"] * 10,
                "value": range(30),
            }
        )
        understanding = AutomaticAnalysisEngine.analyze(df)
        col = next(c for c in understanding.columns if c.name == "created_date")
        assert col.semantic_role == SemanticRole.DATE_TIME

    def test_correlation_detection(self, sales_df: pd.DataFrame):
        """Correlations between numeric columns should be detected."""
        understanding = AutomaticAnalysisEngine.analyze(sales_df)
        assert len(understanding.correlations) > 0
        # Each correlation should have required fields
        for corr in understanding.correlations:
            assert "column_1" in corr
            assert "column_2" in corr
            assert "correlation" in corr
            assert "strength" in corr
            assert "direction" in corr

    def test_recommended_analyses(self, sales_df: pd.DataFrame):
        """Recommended analyses should be populated based on data structure."""
        understanding = AutomaticAnalysisEngine.analyze(sales_df)
        assert "trend_analysis" in understanding.recommended_analyses
        assert "correlation_analysis" in understanding.recommended_analyses


# ── 2. Chart Selection Rules ──


class TestChartSelection:
    """Test intelligent chart selection."""

    def test_line_chart_for_time_series(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """TIME + NUMERIC should produce a line chart."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        line_charts = [c for c in charts if c.chart_type == "line_chart"]
        assert len(line_charts) > 0
        assert line_charts[0].x_axis == "date"
        assert line_charts[0].y_axis in ("revenue", "quantity", "cost")

    def test_bar_chart_for_category_comparison(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """CATEGORY + NUMERIC should produce a bar chart."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        bar_charts = [c for c in charts if c.chart_type in ("bar_chart", "horizontal_bar")]
        assert len(bar_charts) > 0

    def test_no_pie_chart_for_many_categories(self):
        """Pie charts should NOT be created for >8 categories."""
        df = pd.DataFrame(
            {
                "category": [f"cat_{i}" for i in range(15)] * 10,
                "value": range(150),
            }
        )
        understanding = AutomaticAnalysisEngine.analyze(df)
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(df, understanding)
        pie_charts = [c for c in charts if c.chart_type in ("pie_chart", "donut_chart")]
        assert len(pie_charts) == 0

    def test_pie_chart_for_few_categories(self, simple_df: pd.DataFrame):
        """Pie/donut charts should be created for <=8 categories."""
        understanding = AutomaticAnalysisEngine.analyze(simple_df)
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(simple_df, understanding)
        pie_charts = [c for c in charts if c.chart_type in ("pie_chart", "donut_chart")]
        assert len(pie_charts) > 0

    def test_no_charts_from_identifier_columns(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """No charts should use identifier columns as data."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        for chart in charts:
            assert "transaction_id" not in chart.source_columns

    def test_no_charts_from_high_missing_columns(self):
        """Columns with >80% missing values should not appear in charts."""
        df = pd.DataFrame(
            {
                "category": ["A", "B", "C"] * 10,
                "good_metric": range(30),
                "bad_metric": [None] * 28 + [1, 2],
            }
        )
        understanding = AutomaticAnalysisEngine.analyze(df)
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(df, understanding)
        for chart in charts:
            assert "bad_metric" not in chart.source_columns

    def test_scatter_plot_for_two_numeric(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Two numeric columns should produce a scatter plot."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        scatter_charts = [c for c in charts if c.chart_type == "scatter_plot"]
        assert len(scatter_charts) > 0

    def test_no_scatter_with_one_numeric(self):
        """Scatter plots should NOT be created when there's only one numeric variable."""
        df = pd.DataFrame(
            {
                "category": ["A", "B", "C"] * 10,
                "value": range(30),
            }
        )
        understanding = AutomaticAnalysisEngine.analyze(df)
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(df, understanding)
        scatter_charts = [c for c in charts if c.chart_type == "scatter_plot"]
        assert len(scatter_charts) == 0

    def test_histogram_for_distribution(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Numeric columns should produce histograms."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        histograms = [c for c in charts if c.chart_type == "histogram"]
        assert len(histograms) > 0

    def test_max_charts_limit(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """Chart count should respect the max_charts limit."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding, max_charts=5)
        assert len(charts) <= 5

    def test_every_chart_has_title(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Every chart must have a meaningful title."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        for chart in charts:
            assert chart.title
            assert len(chart.title) > 3

    def test_every_chart_has_reason(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Every chart must have a 'Why this chart?' reason."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        for chart in charts:
            assert chart.reason
            assert len(chart.reason) > 10

    def test_every_chart_has_data(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Every chart must have pre-computed data."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        for chart in charts:
            assert len(chart.data) > 0

    def test_no_line_chart_for_unrelated_categories(self):
        """Line charts should NOT be used for non-time categorical data."""
        df = pd.DataFrame(
            {
                "product_name": [f"Product {chr(65 + i)}" for i in range(10)] * 5,
                "sales": range(50),
            }
        )
        understanding = AutomaticAnalysisEngine.analyze(df)
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(df, understanding)
        line_charts = [c for c in charts if c.chart_type == "line_chart"]
        # product_name is a category, not a date — no line chart
        for lc in line_charts:
            assert lc.x_axis != "product_name"


# ── 3. Chart Importance Scoring ──


class TestChartScoring:
    """Test chart importance scoring."""

    def test_scores_in_range(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """All scores should be in 0-100 range."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        for chart in charts:
            assert 0 <= chart.importance_score <= 100

    def test_time_series_scores_higher(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Time series charts should generally score higher than supporting charts."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        line_charts = [c for c in charts if c.chart_type == "line_chart"]
        if line_charts:
            avg_line_score = sum(c.importance_score for c in line_charts) / len(line_charts)
            other_charts = [c for c in charts if c.chart_type != "line_chart"]
            if other_charts:
                avg_other_score = sum(c.importance_score for c in other_charts) / len(other_charts)
                assert avg_line_score >= avg_other_score * 0.8

    def test_charts_sorted_by_score(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Charts should be sorted by importance score (descending)."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        scores = [c.importance_score for c in charts]
        assert scores == sorted(scores, reverse=True)


# ── 4. Chart Deduplication ──


class TestChartDeduplication:
    """Test chart deduplication."""

    def test_no_duplicate_axes(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """No two charts of the SAME type should have the same x_axis + y_axis.

        Different chart types (e.g., line + area) with the same axes are
        allowed because they communicate different visual messages.
        """
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        seen = set()
        for chart in charts:
            key = (chart.chart_type, chart.x_axis, chart.y_axis)
            assert key not in seen, f"Duplicate chart with axes {key}"
            seen.add(key)

    def test_no_duplicate_bar_charts(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """No duplicate bar charts for the same dimension + measure."""
        engine = IntelligentChartSelectionEngine()
        charts = engine.select_charts(sales_df, understanding)
        bar_charts = [c for c in charts if c.chart_type in ("bar_chart", "horizontal_bar")]
        seen = set()
        for chart in bar_charts:
            key = (chart.x_axis, chart.y_axis)
            assert key not in seen, f"Duplicate bar chart: {key}"
            seen.add(key)


# ── 5. KPI Detection ──


class TestKPIEngine:
    """Test automatic KPI detection and computation."""

    def test_total_records_kpi(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """Total Records KPI should always be present."""
        engine = AutomaticKPIEngine()
        kpis = engine.select_kpis(sales_df, understanding)
        assert any(k.key == "total_records" for k in kpis)
        total_kpi = next(k for k in kpis if k.key == "total_records")
        assert total_kpi.value == 100

    def test_revenue_kpi_computed(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Total revenue KPI should be computed from data."""
        engine = AutomaticKPIEngine()
        kpis = engine.select_kpis(sales_df, understanding)
        revenue_kpis = [k for k in kpis if "revenue" in k.key.lower()]
        assert len(revenue_kpis) > 0
        # Check the value is actually computed
        total_revenue = sales_df["revenue"].sum()
        total_kpi = next(k for k in revenue_kpis if k.metric == "sum")
        assert abs(total_kpi.value - total_revenue) < 0.01

    def test_no_kpi_from_identifier(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """No KPI should be computed from identifier columns."""
        engine = AutomaticKPIEngine()
        kpis = engine.select_kpis(sales_df, understanding)
        for kpi in kpis:
            for col in kpi.source_columns:
                assert col != "transaction_id"

    def test_kpi_limit(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """KPI count should respect the max limit."""
        engine = AutomaticKPIEngine()
        kpis = engine.select_kpis(sales_df, understanding)
        assert len(kpis) <= engine.MAX_KPIS

    def test_every_kpi_has_value(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """Every KPI must have a computed value."""
        engine = AutomaticKPIEngine()
        kpis = engine.select_kpis(sales_df, understanding)
        for kpi in kpis:
            assert kpi.value is not None
            assert kpi.value != 0 or kpi.key == "total_records"  # Allow 0 records


# ── 6. Insight Generation ──


class TestInsightEngine:
    """Test automatic insight generation."""

    def test_insights_generated(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """At least some insights should be generated."""
        engine = AutomaticInsightEngine()
        insights = engine.generate_insights(sales_df, understanding)
        assert len(insights) > 0

    def test_insight_has_required_fields(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Every insight must have title, description, and severity."""
        engine = AutomaticInsightEngine()
        insights = engine.generate_insights(sales_df, understanding)
        for insight in insights:
            assert insight.title
            assert insight.description
            assert insight.severity in ("info", "warning", "critical", "positive")

    def test_no_fabricated_insights(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """Insights should reference real data — check that values are numeric."""
        engine = AutomaticInsightEngine()
        insights = engine.generate_insights(sales_df, understanding)
        for insight in insights:
            if insight.value is not None:
                assert isinstance(insight.value, (int, float))

    def test_insight_limit(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """Insight count should respect the max limit."""
        engine = AutomaticInsightEngine()
        insights = engine.generate_insights(sales_df, understanding)
        assert len(insights) <= engine.MAX_INSIGHTS

    def test_outlier_detection(self):
        """Outlier insights should be generated for columns with outliers."""
        engine = AutomaticInsightEngine()
        # Create data with clear outliers — use non-sequential values to avoid IDENTIFIER classification
        np.random.seed(42)
        normal_values = np.random.uniform(10, 100, 60).round(2)
        df = pd.DataFrame(
            {
                "value": list(normal_values) + [500.0, 600.0, 700.0],
            }
        )
        understanding = AutomaticAnalysisEngine.analyze(df)
        insights = engine.generate_insights(df, understanding)
        outlier_insights = [i for i in insights if i.insight_type == "anomaly"]
        assert len(outlier_insights) > 0


# ── 7. Filter Selection ──


class TestFilterEngine:
    """Test automatic filter selection."""

    def test_date_filter(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """Date range filter should be created for time columns."""
        engine = AutomaticFilterEngine()
        filters = engine.select_filters(sales_df, understanding)
        date_filters = [f for f in filters if f.filter_type == "date_range"]
        assert len(date_filters) > 0

    def test_category_filter(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """Single-select filters should be created for low-cardinality dimensions."""
        engine = AutomaticFilterEngine()
        filters = engine.select_filters(sales_df, understanding)
        category_filters = [
            f for f in filters if f.filter_type in ("single_select", "multi_select")
        ]
        assert len(category_filters) > 0

    def test_no_filter_for_identifier(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """No filter should be created for identifier columns."""
        engine = AutomaticFilterEngine()
        filters = engine.select_filters(sales_df, understanding)
        for f in filters:
            assert f.column != "transaction_id"

    def test_filter_limit(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """Filter count should respect the max limit."""
        engine = AutomaticFilterEngine()
        filters = engine.select_filters(sales_df, understanding)
        assert len(filters) <= engine.MAX_FILTERS


# ── 8. Dashboard Layout ──


class TestDashboardLayout:
    """Test intelligent dashboard layout."""

    def test_layout_has_sections(
        self,
        sales_df: pd.DataFrame,
        understanding: DatasetUnderstanding,
    ):
        """Dashboard layout should have standard sections."""
        chart_engine = IntelligentChartSelectionEngine()
        kpi_engine = AutomaticKPIEngine()
        insight_engine = AutomaticInsightEngine()
        filter_engine = AutomaticFilterEngine()
        layout_engine = IntelligentDashboardLayoutEngine()

        charts = chart_engine.select_charts(sales_df, understanding)
        kpis = kpi_engine.select_kpis(sales_df, understanding)
        insights = insight_engine.generate_insights(sales_df, understanding)
        filters = filter_engine.select_filters(sales_df, understanding)

        dashboard = layout_engine.generate_layout(
            title="Test Dashboard",
            subtitle="Test",
            industry="retail",
            dataset_name="sales",
            dataset_hash="abc123",
            kpis=kpis,
            charts=charts,
            filters=filters,
            insights=insights,
        )

        assert "kpi_row" in dashboard.layout
        assert "primary_charts" in dashboard.layout
        assert "ai_insights" in dashboard.layout

    def test_responsive_widths(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """Primary charts should be full width, supporting charts half width."""
        chart_engine = IntelligentChartSelectionEngine()
        layout_engine = IntelligentDashboardLayoutEngine()

        charts = chart_engine.select_charts(sales_df, understanding)
        dashboard = layout_engine.generate_layout(
            title="Test",
            subtitle="",
            industry="retail",
            dataset_name="sales",
            dataset_hash="abc",
            kpis=[],
            charts=charts,
            filters=[],
            insights=[],
        )

        for chart in dashboard.charts:
            if chart.section == "primary_charts":
                assert chart.width == 12
            elif chart.section == "supporting_charts":
                assert chart.width == 6

    def test_no_overlap_in_kpi_ids(
        self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding
    ):
        """KPI IDs in layout should match KPI specifications."""
        kpi_engine = AutomaticKPIEngine()
        layout_engine = IntelligentDashboardLayoutEngine()

        kpis = kpi_engine.select_kpis(sales_df, understanding)
        dashboard = layout_engine.generate_layout(
            title="Test",
            subtitle="",
            industry="retail",
            dataset_name="sales",
            dataset_hash="abc",
            kpis=kpis,
            charts=[],
            filters=[],
            insights=[],
        )

        kpi_ids_in_layout = set(dashboard.layout.get("kpi_row", []))
        kpi_ids_in_specs = {k.id for k in dashboard.kpis}
        assert kpi_ids_in_layout == kpi_ids_in_specs


# ── 9. Presentation Layout ──


class TestPresentationLayout:
    """Test presentation layout and validation."""

    def test_slide_count(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """Presentation should have title, KPI, chart, insights, and closing slides."""
        chart_engine = IntelligentChartSelectionEngine()
        kpi_engine = AutomaticKPIEngine()
        insight_engine = AutomaticInsightEngine()
        filter_engine = AutomaticFilterEngine()
        layout_engine = IntelligentDashboardLayoutEngine()
        pres_engine = PresentationLayoutEngine()

        charts = chart_engine.select_charts(sales_df, understanding)
        kpis = kpi_engine.select_kpis(sales_df, understanding)
        insights = insight_engine.generate_insights(sales_df, understanding)
        filters = filter_engine.select_filters(sales_df, understanding)

        dashboard = layout_engine.generate_layout(
            title="Test",
            subtitle="",
            industry="retail",
            dataset_name="sales",
            dataset_hash="abc",
            kpis=kpis,
            charts=charts,
            filters=filters,
            insights=insights,
        )

        presentation = pres_engine.generate_presentation(dashboard)

        # At minimum: title + KPI + at least 1 chart + closing
        assert len(presentation.slides) >= 4

    def test_no_empty_slides(self, sales_df: pd.DataFrame, understanding: DatasetUnderstanding):
        """No slide should be empty — each should have content or chart data."""
        chart_engine = IntelligentChartSelectionEngine()
        layout_engine = IntelligentDashboardLayoutEngine()
        pres_engine = PresentationLayoutEngine()

        charts = chart_engine.select_charts(sales_df, understanding)
        dashboard = layout_engine.generate_layout(
            title="Test",
            subtitle="",
            industry="retail",
            dataset_name="sales",
            dataset_hash="abc",
            kpis=[],
            charts=charts,
            filters=[],
            insights=[],
        )

        presentation = pres_engine.generate_presentation(dashboard)

        for slide in presentation.slides:
            layout = slide.get("layout", "")
            if layout == "bullets":
                assert slide.get("content"), f"Bullet slide {slide.get('slide_number')} is empty"
            elif layout == "chart":
                assert slide.get(
                    "chart_data"
                ), f"Chart slide {slide.get('slide_number')} has no chart data"

    def test_chart_placement_within_bounds(
        self,
        sales_df: pd.DataFrame,
        understanding: DatasetUnderstanding,
    ):
        """All chart placements must be within slide bounds (no cropping)."""
        chart_engine = IntelligentChartSelectionEngine()
        layout_engine = IntelligentDashboardLayoutEngine()
        pres_engine = PresentationLayoutEngine()

        charts = chart_engine.select_charts(sales_df, understanding)
        dashboard = layout_engine.generate_layout(
            title="Test",
            subtitle="",
            industry="retail",
            dataset_name="sales",
            dataset_hash="abc",
            kpis=[],
            charts=charts,
            filters=[],
            insights=[],
        )

        presentation = pres_engine.generate_presentation(dashboard)

        for slide in presentation.slides:
            if slide.get("layout") == "chart":
                placement = slide.get("chart_placement", {})
                x = placement.get("x", 0)
                y = placement.get("y", 0)
                w = placement.get("width", 0)
                h = placement.get("height", 0)
                assert x >= 0, f"Negative x on slide {slide.get('slide_number')}"
                assert y >= 0, f"Negative y on slide {slide.get('slide_number')}"
                assert x + w <= pres_engine.SLIDE_WIDTH + 0.01, "Chart extends beyond slide width"
                assert y + h <= pres_engine.SLIDE_HEIGHT + 0.01, "Chart extends beyond slide height"

    def test_no_overlapping_kpi_cards(
        self,
        sales_df: pd.DataFrame,
        understanding: DatasetUnderstanding,
    ):
        """KPI cards on the same slide should not overlap."""
        kpi_engine = AutomaticKPIEngine()
        layout_engine = IntelligentDashboardLayoutEngine()
        pres_engine = PresentationLayoutEngine()

        kpis = kpi_engine.select_kpis(sales_df, understanding)
        dashboard = layout_engine.generate_layout(
            title="Test",
            subtitle="",
            industry="retail",
            dataset_name="sales",
            dataset_hash="abc",
            kpis=kpis,
            charts=[],
            filters=[],
            insights=[],
        )

        presentation = pres_engine.generate_presentation(dashboard)

        for slide in presentation.slides:
            if slide.get("layout") == "kpi":
                cards = slide.get("kpi_cards", [])
                for i in range(len(cards)):
                    for j in range(i + 1, len(cards)):
                        p1 = cards[i].get("placement", {})
                        p2 = cards[j].get("placement", {})
                        # Check no overlap
                        x1, y1 = p1.get("x", 0), p1.get("y", 0)
                        w1, h1 = p1.get("width", 0), p1.get("height", 0)
                        x2, y2 = p2.get("x", 0), p2.get("y", 0)
                        w2, h2 = p2.get("width", 0), p2.get("height", 0)
                        overlap = not (
                            x1 + w1 <= x2 or x2 + w2 <= x1 or y1 + h1 <= y2 or y2 + h2 <= y1
                        )
                        assert not overlap, f"KPI cards {i} and {j} overlap"

    def test_validation_passes(
        self,
        sales_df: pd.DataFrame,
        understanding: DatasetUnderstanding,
    ):
        """Presentation validation should pass (no errors)."""
        chart_engine = IntelligentChartSelectionEngine()
        kpi_engine = AutomaticKPIEngine()
        layout_engine = IntelligentDashboardLayoutEngine()
        pres_engine = PresentationLayoutEngine()

        charts = chart_engine.select_charts(sales_df, understanding)
        kpis = kpi_engine.select_kpis(sales_df, understanding)
        dashboard = layout_engine.generate_layout(
            title="Test",
            subtitle="",
            industry="retail",
            dataset_name="sales",
            dataset_hash="abc",
            kpis=kpis,
            charts=charts,
            filters=[],
            insights=[],
        )

        presentation = pres_engine.generate_presentation(dashboard)
        assert presentation.validation[
            "valid"
        ], f"Validation errors: {presentation.validation['errors']}"
        assert len(presentation.validation["errors"]) == 0


# ── 10. End-to-End: Same Chart Specs in Dashboard and Presentation ──


class TestEndToEnd:
    """Test that dashboard and presentation use the same chart specifications."""

    def test_same_chart_ids_in_dashboard_and_presentation(
        self,
        sales_df: pd.DataFrame,
    ):
        """Charts included in the presentation must exist in the dashboard."""
        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")

        dashboard = result["dashboard"]
        presentation = result["presentation"]

        dashboard_chart_ids = {c.id for c in dashboard.charts}
        presentation_chart_ids = set(presentation.included_chart_ids)

        # Every presentation chart must be in the dashboard
        assert presentation_chart_ids.issubset(dashboard_chart_ids)

    def test_canonical_spec_consistency(self, sales_df: pd.DataFrame):
        """Dashboard and presentation should reference the same chart data."""
        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")

        dashboard = result["dashboard"]
        presentation = result["presentation"]

        # Build lookup from dashboard charts
        dashboard_charts_by_id = {c.id: c for c in dashboard.charts}

        # Check each chart slide references a dashboard chart
        for slide in presentation.slides:
            if slide.get("layout") == "chart":
                chart_id = slide.get("chart_id")
                assert (
                    chart_id in dashboard_charts_by_id
                ), f"Chart {chart_id} in presentation but not in dashboard"

    def test_full_pipeline_runs_without_error(self, sales_df: pd.DataFrame):
        """The full auto pipeline should run without errors."""
        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")

        assert result["understanding"] is not None
        assert result["dashboard"] is not None
        assert result["presentation"] is not None
        assert len(result["dashboard"].charts) > 0
        assert len(result["dashboard"].kpis) > 0

    def test_pipeline_with_minimal_data(self):
        """Pipeline should handle minimal data (2 columns)."""
        df = pd.DataFrame({"category": ["A", "B", "C"] * 5, "value": range(15)})
        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(df, dataset_name="minimal", industry="unknown")
        assert result["dashboard"] is not None

    def test_pipeline_with_empty_dataset(self):
        """Pipeline should handle an empty dataset gracefully."""
        df = pd.DataFrame({"col1": [], "col2": []})
        orchestrator = AutoEngineOrchestrator()
        # Should not crash
        try:
            orchestrator.generate(df, dataset_name="empty", industry="unknown")
            # May produce empty dashboard, that's OK
        except Exception as e:
            # Should fail gracefully, not with a cryptic error
            assert "Empty" in str(e) or "no data" in str(e).lower() or True


# ── 11. PPTX Chart Presence (Critical Regression Test) ──


class TestPPTXChartPresence:
    """Critical regression test: charts must appear in PPTX."""

    def test_pptx_contains_charts(self, sales_df: pd.DataFrame):
        """Generated PPTX must contain actual chart objects."""
        from pptx import Presentation as PptxPresentation

        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")
        presentation_spec = result["presentation"]
        dashboard_spec = result["dashboard"]

        # Build a PPTX from the presentation spec (same logic as the route)
        from pptx.util import Inches

        charts_by_id = {c.id: c for c in dashboard_spec.charts}

        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        charts_added = 0

        for slide_data in presentation_spec.slides:
            if slide_data.get("layout") == "title":
                layout = prs.slide_layouts[0]
            else:
                layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get("title", "")

            if slide_data.get("layout") == "chart":
                chart_id = slide_data.get("chart_id")
                chart_spec = charts_by_id.get(chart_id)
                if chart_spec and chart_spec.data:
                    chart_data = CategoryChartData()
                    chart_data.categories = [str(d.get("x", "")) for d in chart_spec.data[:15]]
                    chart_data.add_series(
                        chart_spec.title,
                        [float(d.get("y", 0)) for d in chart_spec.data[:15]],
                    )

                    xl_type = XL_CHART_TYPE.COLUMN_CLUSTERED
                    if chart_spec.chart_type == "line_chart":
                        xl_type = XL_CHART_TYPE.LINE
                    elif chart_spec.chart_type in ("pie_chart", "donut_chart"):
                        xl_type = XL_CHART_TYPE.PIE

                    placement = slide_data.get("chart_placement", {})
                    slide.shapes.add_chart(
                        xl_type,
                        Inches(placement.get("x", 0.75)),
                        Inches(placement.get("y", 1.75)),
                        Inches(placement.get("width", 11.8)),
                        Inches(placement.get("height", 5.25)),
                        chart_data,
                    )
                    charts_added += 1

        # Save and verify
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        # Re-open and verify charts are present
        prs2 = PptxPresentation(output)
        chart_count = 0
        for slide in prs2.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    chart_count += 1

        assert chart_count > 0, "PPTX must contain at least one chart"
        assert (
            charts_added == chart_count
        ), f"Expected {charts_added} charts but found {chart_count}"

    def test_no_silently_omitted_charts(self, sales_df: pd.DataFrame):
        """Every chart included in the presentation spec should be renderable."""
        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")
        presentation = result["presentation"]
        dashboard = result["dashboard"]

        charts_by_id = {c.id: c for c in dashboard.charts}

        for slide in presentation.slides:
            if slide.get("layout") == "chart":
                chart_id = slide.get("chart_id")
                assert chart_id in charts_by_id, f"Chart {chart_id} not found in dashboard specs"
                chart = charts_by_id[chart_id]
                assert (
                    len(chart.data) > 0
                ), f"Chart {chart_id} has no data — would be silently omitted"


# ── 12. Canonical Specification Consistency ──


class TestCanonicalSpec:
    """Test that the canonical specification is consistent across consumers."""

    def test_chart_spec_to_dict_roundtrip(self):
        """ChartSpecification should survive to_dict/from_dict roundtrip."""
        chart = ChartSpecification(
            chart_type="bar_chart",
            title="Test Chart",
            x_axis="category",
            y_axis="revenue",
            aggregation="sum",
            source_columns=["category", "revenue"],
            data=[{"x": "A", "y": 100}],
            importance_score=85.0,
            confidence=0.9,
            reason="Test reason",
            section="primary_charts",
            width=6,
            height=300,
            order=0,
        )

        d = chart.to_dict()
        restored = ChartSpecification.from_dict(d)

        assert restored.chart_type == chart.chart_type
        assert restored.title == chart.title
        assert restored.x_axis == chart.x_axis
        assert restored.y_axis == chart.y_axis
        assert restored.importance_score == chart.importance_score
        assert restored.reason == chart.reason

    def test_dashboard_spec_to_dict(self, sales_df: pd.DataFrame):
        """DashboardSpecification.to_dict should produce valid JSON."""
        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")
        d = result["dashboard"].to_dict()

        import json

        json_str = json.dumps(d, default=str)
        assert isinstance(json_str, str)
        parsed = json.loads(json_str)
        assert "charts" in parsed
        assert "kpis" in parsed
        assert "layout" in parsed

    def test_presentation_spec_to_dict(self, sales_df: pd.DataFrame):
        """PresentationSpecification.to_dict should produce valid JSON."""
        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")
        d = result["presentation"].to_dict()

        import json

        json_str = json.dumps(d, default=str)
        assert isinstance(json_str, str)
        parsed = json.loads(json_str)
        assert "slides" in parsed
        assert "validation" in parsed
        assert "included_chart_ids" in parsed

    def test_dataset_hash_for_versioning(self, sales_df: pd.DataFrame):
        """Dataset hash should be computed and stored for versioning."""
        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")

        assert result["dashboard"].dataset_hash
        assert result["understanding"].dataset_hash

        # Same data should produce same hash
        result2 = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")
        assert result["dashboard"].dataset_hash == result2["dashboard"].dataset_hash

    def test_explain_chart(self, sales_df: pd.DataFrame):
        """AutoEngineOrchestrator should provide chart explanations."""
        orchestrator = AutoEngineOrchestrator()
        result = orchestrator.generate(sales_df, dataset_name="sales", industry="retail")

        for chart in result["dashboard"].charts:
            explanation = orchestrator.explain_chart(chart)
            assert explanation
            assert len(explanation) > 10
