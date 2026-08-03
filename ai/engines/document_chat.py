"""AI Document Chat Engine — chat with uploaded documents.

Supports PDF, Word, Excel, CSV, PowerPoint, and text files.
The AI answers using only the uploaded document unless explicitly instructed otherwise.
"""

import os
import tempfile

from sqlalchemy.orm import Session as DbSession

from ai.config import AI_DOC_ALLOWED_TYPES, AI_DOC_MAX_SIZE_MB
from ai.gateway import AIGateway
from ai.models import AIDocument
from ai.security import AISecurityLayer


class DocumentChatEngine:
    """Enables chatting with uploaded documents."""

    def __init__(self, db: DbSession):
        self.db = db
        self.gateway = AIGateway(db)
        self.security = AISecurityLayer(db)

    def upload_document(
        self, filename: str, file_content: bytes, file_type: str, user_id: int | None = None,
        organization_id: int | None = None,
    ) -> dict:
        """Upload and index a document for chat.

        Returns:
            Dict with document_id, filename, file_type, page_count, is_indexed.
        """
        # Validate file size
        if len(file_content) > AI_DOC_MAX_SIZE_MB * 1024 * 1024:
            raise ValueError(f"File exceeds maximum size of {AI_DOC_MAX_SIZE_MB}MB")

        # Validate file type
        if file_type.lower() not in AI_DOC_ALLOWED_TYPES:
            raise ValueError(
                f"File type '{file_type}' not allowed. Allowed: {', '.join(AI_DOC_ALLOWED_TYPES)}"
            )

        # Save file
        upload_dir = os.path.join(tempfile.gettempdir(), "ai_documents")
        os.makedirs(upload_dir, exist_ok=True)
        file_path = os.path.join(upload_dir, f"{filename}")
        with open(file_path, "wb") as f:
            f.write(file_content)

        # Extract text based on file type
        extracted_text, metadata = self._extract_text(file_path, file_type)

        # Save to database
        doc = AIDocument(
            filename=filename,
            file_type=file_type,
            file_size=len(file_content),
            file_path=file_path,
            extracted_text=extracted_text,
            metadata=metadata,
            user_id=user_id,
            organization_id=organization_id,
            is_indexed=True,
        )
        self.db.add(doc)
        self.db.commit()
        self.db.refresh(doc)

        return {
            "document_id": doc.id,
            "filename": filename,
            "file_type": file_type,
            "page_count": metadata.get("page_count") if metadata else None,
            "is_indexed": True,
        }

    def chat(self, document_id: int, question: str, user_id: int | None = None) -> dict:
        """Chat with a document.

        Returns:
            Dict with answer, citations, confidence_score.
        """
        # Get document
        doc = self.db.query(AIDocument).filter(AIDocument.id == document_id).first()
        if not doc:
            return {"answer": "Document not found.", "citations": None, "confidence_score": None}

        if not doc.extracted_text:
            return {
                "answer": "Document content could not be extracted.",
                "citations": None,
                "confidence_score": None,
            }

        # Validate question
        question = self.security.validate_input(question)

        # Build context with document content
        # Truncate text to avoid token limits
        doc_text = doc.extracted_text[:8000]

        result = self.gateway.chat(
            user_message=(
                f"Document: {doc.filename}\n\n"
                f"Document content:\n{doc_text}\n\n"
                f"Question: {question}\n\n"
                f"Answer the question using ONLY the document content above. "
                f"If the answer is not in the document, say 'The document does not contain this information.'"
            ),
            assistant_type="data_copilot",
            user_id=user_id,
            context={"document_id": document_id, "filename": doc.filename},
        )

        return {
            "answer": result["response"],
            "citations": [{"source": doc.filename, "type": doc.file_type}],
            "confidence_score": result.get("confidence_score"),
        }

    def _extract_text(self, file_path: str, file_type: str) -> tuple[str, dict]:
        """Extract text from a document based on its type."""
        metadata = {}

        if file_type == "csv":
            return self._extract_csv(file_path, metadata)
        elif file_type == "xlsx":
            return self._extract_excel(file_path, metadata)
        elif file_type == "pdf":
            return self._extract_pdf(file_path, metadata)
        elif file_type == "docx":
            return self._extract_docx(file_path, metadata)
        elif file_type == "pptx":
            return self._extract_pptx(file_path, metadata)
        elif file_type == "txt":
            return self._extract_txt(file_path, metadata)
        else:
            return "", metadata

    def _extract_csv(self, file_path: str, metadata: dict) -> tuple[str, dict]:
        """Extract text from CSV."""
        import pandas as pd

        try:
            df = pd.read_csv(file_path)
            metadata["row_count"] = len(df)
            metadata["column_count"] = len(df.columns)
            metadata["columns"] = list(df.columns)
            # Convert to string representation
            text = f"CSV file with {len(df)} rows and {len(df.columns)} columns.\n"
            text += f"Columns: {', '.join(df.columns)}\n\n"
            text += df.head(50).to_string()
            return text, metadata
        except Exception as e:
            return f"Error reading CSV: {e}", metadata

    def _extract_excel(self, file_path: str, metadata: dict) -> tuple[str, dict]:
        """Extract text from Excel."""
        import pandas as pd

        try:
            xls = pd.ExcelFile(file_path)
            metadata["sheet_names"] = xls.sheet_names
            metadata["page_count"] = len(xls.sheet_names)
            text = (
                f"Excel file with {len(xls.sheet_names)} sheets: {', '.join(xls.sheet_names)}\n\n"
            )
            for sheet in xls.sheet_names[:5]:  # Limit to first 5 sheets
                df = pd.read_excel(file_path, sheet_name=sheet)
                text += f"Sheet: {sheet} ({len(df)} rows, {len(df.columns)} columns)\n"
                text += df.head(20).to_string() + "\n\n"
            return text, metadata
        except Exception as e:
            return f"Error reading Excel: {e}", metadata

    def _extract_pdf(self, file_path: str, metadata: dict) -> tuple[str, dict]:
        """Extract text from PDF."""
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(file_path)
            metadata["page_count"] = len(reader.pages)
            text = ""
            for page in reader.pages[:50]:  # Limit to first 50 pages
                text += page.extract_text() + "\n"
            return text, metadata
        except ImportError:
            return "PDF text extraction requires PyPDF2. Install with: pip install PyPDF2", metadata
        except ModuleNotFoundError:
            return "PDF text extraction requires PyPDF2. Install with: pip install PyPDF2", metadata
        except Exception as e:
            return f"Error reading PDF: {e}", metadata

    def _extract_docx(self, file_path: str, metadata: dict) -> tuple[str, dict]:
        """Extract text from Word document."""
        try:
            from docx import Document

            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            metadata["paragraph_count"] = len(paragraphs)
            text = "\n".join(paragraphs)
            return text, metadata
        except ImportError:
            return (
                "Word document extraction requires python-docx. Install with: pip install python-docx",
                metadata,
            )
        except ModuleNotFoundError:
            return (
                "Word document extraction requires python-docx. Install with: pip install python-docx",
                metadata,
            )
        except Exception as e:
            return f"Error reading Word document: {e}", metadata

    def _extract_pptx(self, file_path: str, metadata: dict) -> tuple[str, dict]:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            metadata["page_count"] = len(prs.slides)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                text += "\n"
            return text, metadata
        except ImportError:
            return (
                "PowerPoint extraction requires python-pptx. Install with: pip install python-pptx",
                metadata,
            )
        except ModuleNotFoundError:
            return (
                "PowerPoint extraction requires python-pptx. Install with: pip install python-pptx",
                metadata,
            )
        except Exception as e:
            return f"Error reading PowerPoint: {e}", metadata

    def _extract_txt(self, file_path: str, metadata: dict) -> tuple[str, dict]:
        """Extract text from plain text file."""
        try:
            with open(file_path, encoding="utf-8", errors="replace") as f:
                text = f.read()
            metadata["character_count"] = len(text)
            return text, metadata
        except Exception as e:
            return f"Error reading text file: {e}", metadata

    def list_documents(self, user_id: int | None = None, organization_id: int | None = None) -> list[dict]:
        """List uploaded documents."""
        query = self.db.query(AIDocument)
        if organization_id is not None:
            query = query.filter(AIDocument.organization_id == organization_id)
        if user_id:
            query = query.filter(AIDocument.user_id == user_id)
        docs = query.order_by(AIDocument.created_at.desc()).all()
        return [
            {
                "id": d.id,
                "filename": d.filename,
                "file_type": d.file_type,
                "file_size": d.file_size,
                "is_indexed": d.is_indexed,
                "created_at": str(d.created_at) if d.created_at else None,
            }
            for d in docs
        ]
