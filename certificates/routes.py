"""Certificate Intelligence API routes.

Builds on the existing Smart Data Capture pipeline to provide
certificate-specific endpoints for:
  - Multi-file batch upload (up to CERTIFICATE_MAX_BATCH_SIZE files)
  - Certificate search with filters
  - Certificate dashboard analytics
  - CSV/XLSX export of approved certificate data
  - Verification status management
  - Report and presentation generation from certificate data

All endpoints enforce organization isolation and authentication.
"""

from __future__ import annotations

import csv
import io
import logging

from fastapi import (
    APIRouter,
    Depends,
    File,
    Form,
    HTTPException,
    Query,
    UploadFile,
    status,
)
from fastapi.responses import StreamingResponse
from sqlalchemy import or_, select
from sqlalchemy.orm import Session as DbSession

import config
from audit.service import log_audit_event
from capture.models import (
    CaptureDocument,
    CaptureField,
    CertificateVerification,
)
from capture.service import CaptureError, CaptureService
from shared.database import get_db
from shared.dependencies import get_current_user
from shared.tenant import get_current_organization_id

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/certificates", tags=["certificate-intelligence"])

# Certificate document types (from the capture registry)
CERTIFICATE_DOC_TYPES = {
    "academic_certificate",
    "degree_certificate",
    "diploma",
    "professional_certificate",
    "training_certificate",
    "certificate_of_completion",
    "certificate_of_attendance",
    "membership_certificate",
    "license_certification",
}


def _get_max_batch_size() -> int:
    """Return the configured maximum batch size for certificate uploads."""
    return getattr(config, "CERTIFICATE_MAX_BATCH_SIZE", 50)


def _serialize_certificate(doc: CaptureDocument, fields: list[CaptureField] | None = None) -> dict:
    """Serialize a capture document as a certificate record."""
    result = {
        "id": doc.id,
        "batch_id": doc.batch_id,
        "filename": doc.filename,
        "file_type": doc.file_type,
        "status": doc.status,
        "error_message": doc.error_message,
        "document_type": doc.document_type,
        "document_type_label": doc.document_type_label,
        "classification_confidence": doc.classification_confidence,
        "overall_confidence": doc.overall_confidence,
        "needs_type_confirmation": doc.needs_type_confirmation,
        "verification_status": doc.verification_status,
        "verification_method": doc.verification_method,
        "verified_at": doc.verified_at.isoformat() if doc.verified_at else None,
        "duplicate_of_id": doc.duplicate_of_id,
        "created_at": doc.created_at.isoformat() if doc.created_at else None,
        "processed_at": doc.processed_at.isoformat() if doc.processed_at else None,
        "approved_at": doc.approved_at.isoformat() if doc.approved_at else None,
    }
    if fields is not None:
        result["fields"] = [
            {
                "id": f.id,
                "field_name": f.field_name,
                "field_label": f.field_label,
                "data_type": f.data_type,
                "value": f.value,
                "raw_value": f.raw_value,
                "confidence_score": f.confidence_score,
                "is_low_confidence": f.is_low_confidence,
                "was_corrected": f.was_corrected,
                "is_valid": f.is_valid,
                "validation_message": f.validation_message,
            }
            for f in fields
        ]
    return result


def _is_certificate_type(doc_type: str | None) -> bool:
    """Check if a document type is a certificate type."""
    return doc_type in CERTIFICATE_DOC_TYPES


# ═══════════════════════════════════════════════════════════════
# Batch upload (multi-file)
# ═══════════════════════════════════════════════════════════════


@router.post("/upload", status_code=status.HTTP_201_CREATED)
async def upload_certificates(
    files: list[UploadFile] = File(...),
    batch_name: str | None = Form(None),
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Upload multiple certificate files in a single batch.

    Accepts up to CERTIFICATE_MAX_BATCH_SIZE files (default 50).
    Each file is processed through the existing capture pipeline:
    upload → preprocess → OCR → classify → extract → validate.
    """
    org_id = get_current_organization_id(current_user, db)
    max_batch = _get_max_batch_size()

    logger.info(
        "Certificate upload started: org_id=%s user_id=%s file_count=%d batch_name=%s",
        org_id, current_user["id"], len(files), batch_name,
    )

    if len(files) > max_batch:
        logger.warning(
            "Certificate upload rejected — too many files: org_id=%s count=%d max=%d",
            org_id, len(files), max_batch,
        )
        raise HTTPException(
            status_code=413,
            detail=f"Too many files. Maximum {max_batch} certificates per batch. "
            f"Received {len(files)}.",
        )

    svc = CaptureService(db)
    batch = svc.create_batch(
        org_id, current_user["id"], batch_name or "Certificate Batch", "certificates"
    )
    logger.info("Batch created: batch_id=%s org_id=%s", batch.id, org_id)

    results: list[dict] = []
    succeeded = 0
    failed = 0
    review_required = 0

    for file in files:
        try:
            content = await file.read()
            doc = svc.upload_document(
                org_id,
                current_user["id"],
                file.filename or "certificate",
                content,
                source="web",
                batch_id=batch.id,
            )
            logger.info(
                "Document stored: doc_id=%s filename=%s file_type=%s size=%d org_id=%s",
                doc.id, doc.filename, doc.file_type, len(content), org_id,
            )
            # Enqueue background job for processing instead of synchronous
            job_id = None
            try:
                from jobs.handlers import register_builtin_handlers
                from jobs.service import JobService, get_registered_types

                if "ocr_document" not in get_registered_types():
                    register_builtin_handlers()

                job_svc = JobService(db)
                job = job_svc.create_job(
                    organization_id=org_id,
                    user_id=current_user["id"],
                    job_type="ocr_document",
                    name=f"Certificate OCR: {doc.filename}",
                    description=f"Processing certificate '{doc.filename}'",
                    payload={"document_id": doc.id, "organization_id": org_id},
                )
                job_id = job.id
                db.commit()
                logger.info(
                    "Job enqueued: job_id=%s doc_id=%s type=ocr_document org_id=%s",
                    job_id, doc.id, org_id,
                )
            except Exception as e:
                logger.warning(
                    "Job system unavailable for doc_id=%s, using thread fallback: %s",
                    doc.id, e,
                )
                import threading

                def _process_doc(doc_id: int) -> None:
                    from shared.database import get_engine, get_session_factory

                    engine = get_engine()
                    factory = get_session_factory(engine)
                    session = factory()
                    try:
                        CaptureService(session).process_document(doc_id)
                    finally:
                        session.close()

                threading.Thread(target=_process_doc, args=(doc.id,), daemon=True).start()

            succeeded += 1
            result = _serialize_certificate(doc)
            result["job_id"] = job_id
            results.append(result)
        except CaptureError as e:
            failed += 1
            results.append(
                {
                    "filename": file.filename,
                    "status": "failed",
                    "error_message": str(e),
                }
            )
            logger.warning("Certificate upload failed for %s: %s", file.filename, e)
        except Exception as e:
            failed += 1
            results.append(
                {
                    "filename": file.filename,
                    "status": "failed",
                    "error_message": f"Processing error: {e}",
                }
            )
            logger.exception("Certificate processing failed for %s", file.filename)

    batch.total_documents = len(files)
    batch.processed_documents = succeeded
    batch.failed_documents = failed
    db.commit()

    log_audit_event(
        db=db,
        action="certificate.batch_upload",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="certificate_batch",
        resource_id=batch.id,
        new_values={
            "batch_name": batch.name,
            "total": len(files),
            "succeeded": succeeded,
            "failed": failed,
            "review_required": review_required,
        },
    )
    db.commit()

    logger.info(
        "Certificate upload complete: batch_id=%s org_id=%s total=%d succeeded=%d failed=%d",
        batch.id, org_id, len(files), succeeded, failed,
    )

    return {
        "batch_id": batch.id,
        "batch_name": batch.name,
        "total": len(files),
        "succeeded": succeeded,
        "failed": failed,
        "review_required": review_required,
        "certificates": results,
    }


# ═══════════════════════════════════════════════════════════════
# Search
# ═══════════════════════════════════════════════════════════════


@router.get("/search")
async def search_certificates(
    q: str | None = Query(None, description="Search query (name, certificate number, institution)"),
    certificate_type: str | None = Query(None, description="Filter by certificate type"),
    verification_status: str | None = Query(None, description="Filter by verification status"),
    review_status: str | None = Query(
        None, description="Filter by review status (approved, ready_for_review, etc.)"
    ),
    institution: str | None = Query(None, description="Filter by institution name"),
    year: int | None = Query(None, description="Filter by award/issue year"),
    limit: int = Query(50, ge=1, le=200),
    offset: int = Query(0, ge=0),
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Search certificates with filters. Only returns certificate-type documents."""
    org_id = get_current_organization_id(current_user, db)

    query = (
        select(CaptureDocument)
        .where(
            CaptureDocument.organization_id == org_id,
            CaptureDocument.document_type.in_(CERTIFICATE_DOC_TYPES),
        )
        .order_by(CaptureDocument.id.desc())
    )

    if certificate_type:
        query = query.where(CaptureDocument.document_type == certificate_type)
    if verification_status:
        query = query.where(CaptureDocument.verification_status == verification_status)
    if review_status:
        query = query.where(CaptureDocument.status == review_status)

    # Text search: match filename or OCR text
    if q:
        pattern = f"%{q}%"
        query = query.where(
            or_(
                CaptureDocument.filename.ilike(pattern),
                CaptureDocument.raw_ocr_text.ilike(pattern),
            )
        )

    # Total count
    total_query = select(CaptureDocument.id).where(
        CaptureDocument.organization_id == org_id,
        CaptureDocument.document_type.in_(CERTIFICATE_DOC_TYPES),
    )
    if certificate_type:
        total_query = total_query.where(CaptureDocument.document_type == certificate_type)
    if verification_status:
        total_query = total_query.where(CaptureDocument.verification_status == verification_status)
    if review_status:
        total_query = total_query.where(CaptureDocument.status == review_status)
    if q:
        total_query = total_query.where(
            or_(
                CaptureDocument.filename.ilike(pattern),
                CaptureDocument.raw_ocr_text.ilike(pattern),
            )
        )
    total = len(list(db.execute(total_query).scalars()))

    # Pagination
    query = query.limit(limit).offset(offset)
    docs = list(db.execute(query).scalars().all())

    # Filter by institution/year if specified (requires field lookup)
    certificates = []
    for doc in docs:
        fields = (
            db.execute(select(CaptureField).where(CaptureField.document_id == doc.id))
            .scalars()
            .all()
        )
        field_dict = {f.field_name: f.value for f in fields}

        # Institution filter
        if institution:
            inst_val = (field_dict.get("institution") or "").lower()
            if institution.lower() not in inst_val:
                continue

        # Year filter
        if year:
            date_val = (
                field_dict.get("date_awarded")
                or field_dict.get("date_issued")
                or field_dict.get("graduation_date")
                or ""
            )
            if str(year) not in date_val:
                continue

        certificates.append(_serialize_certificate(doc, fields))

    return {
        "certificates": certificates,
        "total": total if not (institution or year) else len(certificates),
        "limit": limit,
        "offset": offset,
    }


# ═══════════════════════════════════════════════════════════════
# Dashboard analytics
# ═══════════════════════════════════════════════════════════════


@router.get("/dashboard")
async def certificate_dashboard(
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Certificate intelligence dashboard with real data from the organization."""
    org_id = get_current_organization_id(current_user, db)

    docs = list(
        db.execute(
            select(CaptureDocument).where(
                CaptureDocument.organization_id == org_id,
                CaptureDocument.document_type.in_(CERTIFICATE_DOC_TYPES),
            )
        )
        .scalars()
        .all()
    )

    total = len(docs)
    by_status: dict[str, int] = {}
    by_type: dict[str, int] = {}
    by_verification: dict[str, int] = {}
    by_institution: dict[str, int] = {}

    for doc in docs:
        by_status[doc.status] = by_status.get(doc.status, 0) + 1
        if doc.document_type:
            by_type[doc.document_type] = by_type.get(doc.document_type, 0) + 1
        by_verification[doc.verification_status] = (
            by_verification.get(doc.verification_status, 0) + 1
        )

    # Get institution distribution from fields
    for doc in docs:
        fields = (
            db.execute(
                select(CaptureField).where(
                    CaptureField.document_id == doc.id,
                    CaptureField.field_name == "institution",
                )
            )
            .scalars()
            .all()
        )
        for f in fields:
            if f.value:
                inst = f.value.strip()
                by_institution[inst] = by_institution.get(inst, 0) + 1

    # Get year distribution from date fields
    by_year: dict[str, int] = {}
    for doc in docs:
        fields = (
            db.execute(
                select(CaptureField).where(
                    CaptureField.document_id == doc.id,
                    CaptureField.field_name.in_(["date_awarded", "date_issued", "graduation_date"]),
                )
            )
            .scalars()
            .all()
        )
        for f in fields:
            if f.value and len(f.value) >= 4:
                # Extract year from date string
                year_str = f.value[-4:] if f.value[-4:].isdigit() else None
                if year_str:
                    by_year[year_str] = by_year.get(year_str, 0) + 1

    return {
        "total": total,
        "processed": by_status.get("ready_for_review", 0) + by_status.get("approved", 0),
        "processing": by_status.get("preprocessing", 0)
        + by_status.get("extracting", 0)
        + by_status.get("classifying", 0)
        + by_status.get("validating", 0),
        "review_required": by_status.get("ready_for_review", 0),
        "approved": by_status.get("approved", 0),
        "rejected": by_status.get("rejected", 0),
        "failed": by_status.get("failed", 0),
        "duplicates": sum(1 for d in docs if d.duplicate_of_id is not None),
        "verified": by_verification.get("verified", 0),
        "not_verified": by_verification.get("not_verified", 0),
        "verification_pending": by_verification.get("verification_pending", 0),
        "verification_failed": by_verification.get("verification_failed", 0),
        "by_type": by_type,
        "by_status": by_status,
        "by_verification": by_verification,
        "by_institution": dict(
            sorted(by_institution.items(), key=lambda kv: kv[1], reverse=True)[:20]
        ),
        "by_year": dict(sorted(by_year.items())),
    }


# ═══════════════════════════════════════════════════════════════
# Export (CSV / XLSX)
# ═══════════════════════════════════════════════════════════════


def _get_certificate_field_dict(db: DbSession, doc_ids: list[int]) -> dict[int, dict[str, str]]:
    """Get a dict of document_id -> {field_name: value} for the given docs."""
    if not doc_ids:
        return {}
    fields = (
        db.execute(select(CaptureField).where(CaptureField.document_id.in_(doc_ids)))
        .scalars()
        .all()
    )
    result: dict[int, dict[str, str]] = {}
    for f in fields:
        result.setdefault(f.document_id, {})[f.field_name] = f.value or ""
    return result


EXPORT_COLUMNS = [
    "full_name",
    "certificate_type",
    "qualification",
    "degree",
    "programme",
    "course",
    "institution",
    "certificate_number",
    "license_number",
    "date_awarded",
    "date_issued",
    "graduation_date",
    "expiry_date",
    "grade",
    "gpa",
    "department",
    "country",
    "verification_status",
    "review_status",
    "classification_confidence",
    "overall_confidence",
]


@router.get("/export/csv")
async def export_certificates_csv(
    certificate_type: str | None = Query(None),
    review_status: str | None = Query(None),
    verification_status: str | None = Query(None),
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Export approved certificate data as CSV."""
    org_id = get_current_organization_id(current_user, db)

    query = (
        select(CaptureDocument)
        .where(
            CaptureDocument.organization_id == org_id,
            CaptureDocument.document_type.in_(CERTIFICATE_DOC_TYPES),
        )
        .order_by(CaptureDocument.id.desc())
    )
    if certificate_type:
        query = query.where(CaptureDocument.document_type == certificate_type)
    if review_status:
        query = query.where(CaptureDocument.status == review_status)
    if verification_status:
        query = query.where(CaptureDocument.verification_status == verification_status)

    docs = list(db.execute(query).scalars().all())
    field_dicts = _get_certificate_field_dict(db, [d.id for d in docs])

    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(["Filename", "Certificate Type", "Certificate Type Label"] + EXPORT_COLUMNS)

    for doc in docs:
        fields = field_dicts.get(doc.id, {})
        row = [doc.filename, doc.document_type or "", doc.document_type_label or ""]
        for col in EXPORT_COLUMNS:
            if col == "certificate_type":
                row.append(doc.document_type or "")
            elif col == "verification_status":
                row.append(doc.verification_status)
            elif col == "review_status":
                row.append(doc.status)
            elif col == "classification_confidence":
                row.append(
                    f"{doc.classification_confidence:.3f}" if doc.classification_confidence else ""
                )
            elif col == "overall_confidence":
                row.append(f"{doc.overall_confidence:.3f}" if doc.overall_confidence else "")
            else:
                row.append(fields.get(col, ""))
        writer.writerow(row)

    output.seek(0)
    log_audit_event(
        db=db,
        action="certificate.export_csv",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="certificate",
        new_values={"count": len(docs), "format": "csv"},
    )
    db.commit()

    return StreamingResponse(
        iter([output.getvalue()]),
        media_type="text/csv",
        headers={"Content-Disposition": "attachment; filename=certificates_export.csv"},
    )


@router.get("/export/xlsx")
async def export_certificates_xlsx(
    certificate_type: str | None = Query(None),
    review_status: str | None = Query(None),
    verification_status: str | None = Query(None),
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Export approved certificate data as XLSX."""
    try:
        import openpyxl
    except ImportError as e:
        raise HTTPException(
            status_code=500,
            detail="Excel export requires openpyxl. Install with: pip install openpyxl",
        ) from e

    org_id = get_current_organization_id(current_user, db)

    query = (
        select(CaptureDocument)
        .where(
            CaptureDocument.organization_id == org_id,
            CaptureDocument.document_type.in_(CERTIFICATE_DOC_TYPES),
        )
        .order_by(CaptureDocument.id.desc())
    )
    if certificate_type:
        query = query.where(CaptureDocument.document_type == certificate_type)
    if review_status:
        query = query.where(CaptureDocument.status == review_status)
    if verification_status:
        query = query.where(CaptureDocument.verification_status == verification_status)

    docs = list(db.execute(query).scalars().all())
    field_dicts = _get_certificate_field_dict(db, [d.id for d in docs])

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Certificates"
    headers = ["Filename", "Certificate Type", "Certificate Type Label"] + EXPORT_COLUMNS
    ws.append(headers)

    for doc in docs:
        fields = field_dicts.get(doc.id, {})
        row = [doc.filename, doc.document_type or "", doc.document_type_label or ""]
        for col in EXPORT_COLUMNS:
            if col == "certificate_type":
                row.append(doc.document_type or "")
            elif col == "verification_status":
                row.append(doc.verification_status)
            elif col == "review_status":
                row.append(doc.status)
            elif col == "classification_confidence":
                row.append(doc.classification_confidence)
            elif col == "overall_confidence":
                row.append(doc.overall_confidence)
            else:
                row.append(fields.get(col, ""))
        ws.append(row)

    # Auto-width columns
    for col_idx, header in enumerate(headers, 1):
        max_len = len(str(header))
        for row in ws.iter_rows(min_row=2, min_col=col_idx, max_col=col_idx):
            for cell in row:
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
        ws.column_dimensions[ws.cell(row=1, column=col_idx).column_letter].width = min(
            max_len + 2, 50
        )

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)

    log_audit_event(
        db=db,
        action="certificate.export_xlsx",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="certificate",
        new_values={"count": len(docs), "format": "xlsx"},
    )
    db.commit()

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=certificates_export.xlsx"},
    )


# ═══════════════════════════════════════════════════════════════
# Verification
# ═══════════════════════════════════════════════════════════════


@router.post("/{document_id}/verify")
async def verify_certificate(
    document_id: int,
    payload: dict,
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Record a verification attempt for a certificate.

    Verification means an authoritative source confirmed the certificate.
    This does NOT mean the certificate is "genuine" — it means an external
    source was contacted and returned a verification result.
    """
    org_id = get_current_organization_id(current_user, db)

    doc = db.execute(
        select(CaptureDocument).where(
            CaptureDocument.id == document_id,
            CaptureDocument.organization_id == org_id,
        )
    ).scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="Certificate not found")
    if not _is_certificate_type(doc.document_type):
        raise HTTPException(status_code=400, detail="Document is not a certificate type")

    method = payload.get("method", "manual_check")
    verify_status = payload.get("status", "pending")
    source = payload.get("verification_source")
    ref_num = payload.get("reference_number")
    notes = payload.get("notes")
    verified_fields = payload.get("verified_fields")

    if verify_status not in ("pending", "verified", "failed", "inconclusive"):
        raise HTTPException(
            status_code=422,
            detail="Status must be one of: pending, verified, failed, inconclusive",
        )

    # Record the verification attempt
    verification = CertificateVerification(
        organization_id=org_id,
        document_id=document_id,
        method=method,
        status=verify_status,
        verified_by=current_user["id"],
        verification_source=source,
        reference_number=ref_num,
        notes=notes,
        verified_fields=verified_fields,
    )
    db.add(verification)

    # Update document verification status
    if verify_status == "verified":
        doc.verification_status = "verified"
        doc.verification_method = method
        doc.verified_at = __import__("datetime").datetime.now(__import__("datetime").timezone.utc)
        doc.verified_by = current_user["id"]
    elif verify_status == "failed":
        doc.verification_status = "verification_failed"
        doc.verification_method = method
    else:
        doc.verification_status = "verification_pending"
        doc.verification_method = method

    log_audit_event(
        db=db,
        action="certificate.verify",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="certificate",
        resource_id=document_id,
        new_values={
            "method": method,
            "status": verify_status,
            "source": source,
        },
    )
    db.commit()

    return {
        "document_id": document_id,
        "verification_status": doc.verification_status,
        "method": method,
        "status": verify_status,
        "verification_id": verification.id,
    }


@router.get("/{document_id}/verifications")
async def list_verifications(
    document_id: int,
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """List all verification attempts for a certificate."""
    org_id = get_current_organization_id(current_user, db)

    doc = db.execute(
        select(CaptureDocument).where(
            CaptureDocument.id == document_id,
            CaptureDocument.organization_id == org_id,
        )
    ).scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="Certificate not found")

    verifications = (
        db.execute(
            select(CertificateVerification)
            .where(
                CertificateVerification.document_id == document_id,
                CertificateVerification.organization_id == org_id,
            )
            .order_by(CertificateVerification.id.desc())
        )
        .scalars()
        .all()
    )

    return {
        "document_id": document_id,
        "current_status": doc.verification_status,
        "verifications": [
            {
                "id": v.id,
                "method": v.method,
                "status": v.status,
                "verification_source": v.verification_source,
                "reference_number": v.reference_number,
                "notes": v.notes,
                "verified_by": v.verified_by,
                "created_at": v.created_at.isoformat() if v.created_at else None,
            }
            for v in verifications
        ],
    }


# ═══════════════════════════════════════════════════════════════
# Certificate types
# ═══════════════════════════════════════════════════════════════


@router.get("/types")
async def list_certificate_types():
    """List all supported certificate types and their fields."""
    from capture.document_types import CERTIFICATE_TYPES

    return {
        "certificate_types": [
            {
                "key": t.key,
                "label": t.label,
                "industry": t.industry,
                "keywords": t.keywords,
                "fields": [
                    {
                        "name": f.name,
                        "label": f.label,
                        "data_type": f.data_type,
                        "required": f.required,
                        "keywords": f.keywords,
                    }
                    for f in t.fields
                ],
            }
            for t in CERTIFICATE_TYPES
        ],
    }


# ═══════════════════════════════════════════════════════════════
# Analytics dataset (for integration with existing analytics pipeline)
# ═══════════════════════════════════════════════════════════════


@router.post("/to-dataset")
async def certificates_to_dataset(
    payload: dict | None = None,
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Convert approved certificate data into a dataset for the analytics pipeline.

    This bridges certificate data into the existing Data-to-Decision workflow:
    certificates → structured dataset → analysis → dashboard → report → presentation.
    """
    org_id = get_current_organization_id(current_user, db)
    dataset_name = (payload or {}).get("dataset_name", "Certificate Data")

    # Get all approved certificates
    docs = list(
        db.execute(
            select(CaptureDocument).where(
                CaptureDocument.organization_id == org_id,
                CaptureDocument.document_type.in_(CERTIFICATE_DOC_TYPES),
                CaptureDocument.status == "approved",
            )
        )
        .scalars()
        .all()
    )

    if not docs:
        raise HTTPException(
            status_code=422,
            detail="No approved certificates found. Approve certificates first, then export to dataset.",
        )

    field_dicts = _get_certificate_field_dict(db, [d.id for d in docs])

    # Build CSV from certificate data
    output = io.StringIO()
    writer = csv.writer(output)

    # Collect all field names across all certificates
    all_fields: set[str] = set()
    for fd in field_dicts.values():
        all_fields.update(fd.keys())
    all_fields.update(["certificate_type", "verification_status"])

    headers = ["filename", "certificate_type_label"] + sorted(all_fields) + ["verification_status"]
    writer.writerow(headers)

    for doc in docs:
        fields = field_dicts.get(doc.id, {})
        row = [doc.filename, doc.document_type_label or ""]
        for field_name in sorted(all_fields):
            if field_name == "certificate_type":
                row.append(doc.document_type or "")
            elif field_name == "verification_status":
                row.append(doc.verification_status)
            else:
                row.append(fields.get(field_name, ""))
        row.append(doc.verification_status)
        writer.writerow(row)

    csv_content = output.getvalue()

    # Save as a dataset file for the analytics pipeline
    import os

    dataset_dir = os.path.join(getattr(config, "UPLOAD_DIR", "uploads"), "certificate_datasets")
    os.makedirs(dataset_dir, exist_ok=True)
    safe_name = dataset_name.replace(" ", "_").replace("/", "_")
    csv_path = os.path.join(dataset_dir, f"{safe_name}_{org_id}.csv")
    with open(csv_path, "w", newline="", encoding="utf-8") as f:
        f.write(csv_content)

    log_audit_event(
        db=db,
        action="certificate.to_dataset",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="certificate",
        new_values={
            "dataset_name": dataset_name,
            "csv_path": csv_path,
            "row_count": len(docs),
        },
    )
    db.commit()

    return {
        "dataset_name": dataset_name,
        "csv_path": csv_path,
        "row_count": len(docs),
        "field_count": len(headers),
        "message": f"Exported {len(docs)} approved certificates to dataset '{dataset_name}'. "
        f"Use the Data-to-Decision workflow to analyze this data.",
    }


# ═══════════════════════════════════════════════════════════════
# Certificate detail with full analysis
# ═══════════════════════════════════════════════════════════════


def _serialize_fields(fields: list[CaptureField]) -> list[dict]:
    """Serialize CaptureField objects to dicts for analysis."""
    return [
        {
            "field_name": f.field_name,
            "field_label": f.field_label,
            "data_type": f.data_type,
            "value": f.value,
            "raw_value": f.raw_value,
            "confidence_score": f.confidence_score,
            "is_low_confidence": f.is_low_confidence,
            "was_corrected": f.was_corrected,
            "is_valid": f.is_valid,
            "validation_message": f.validation_message,
        }
        for f in fields
    ]


@router.get("/{document_id}/status")
async def get_certificate_status(
    document_id: int,
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Lightweight status check for a single certificate.

    Returns only the document status, error message, and basic metadata.
    Used by the frontend for polling during processing.
    """
    org_id = get_current_organization_id(current_user, db)

    doc = db.execute(
        select(CaptureDocument).where(
            CaptureDocument.id == document_id,
            CaptureDocument.organization_id == org_id,
        )
    ).scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="Certificate not found")

    return {
        "id": doc.id,
        "status": doc.status,
        "error_message": doc.error_message,
        "document_type": doc.document_type,
        "document_type_label": doc.document_type_label,
        "overall_confidence": doc.overall_confidence,
        "classification_confidence": doc.classification_confidence,
    }


@router.get("/{document_id}/detail")
async def get_certificate_detail(
    document_id: int,
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Get a single certificate with full field data and intelligence analysis.

    Returns the document record, all extracted fields with confidence scores,
    and a comprehensive analysis including:
    - Completeness assessment (required vs optional fields)
    - Consistency checks (cross-field validation)
    - Academic performance summary
    - Anomaly detection
    - Actionable recommendations
    """
    from certificates.analysis import analyze_certificate

    org_id = get_current_organization_id(current_user, db)

    doc = db.execute(
        select(CaptureDocument).where(
            CaptureDocument.id == document_id,
            CaptureDocument.organization_id == org_id,
        )
    ).scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="Certificate not found")

    fields = (
        db.execute(
            select(CaptureField)
            .where(CaptureField.document_id == document_id)
            .order_by(CaptureField.id)
        )
        .scalars()
        .all()
    )

    doc_dict = _serialize_certificate(doc)
    field_list = _serialize_fields(fields)
    analysis = analyze_certificate(doc_dict, field_list)

    return {
        "certificate": doc_dict,
        "fields": field_list,
        "analysis": {
            "document_type": analysis.document_type,
            "document_type_label": analysis.document_type_label,
            "classification_confidence": analysis.classification_confidence,
            "overall_confidence": analysis.overall_confidence,
            "summary": analysis.summary,
            "verification_status": analysis.verification_status,
            "is_duplicate": analysis.is_duplicate,
            "duplicate_of_id": analysis.duplicate_of_id,
            "completeness": {
                "total_fields": analysis.completeness.total_fields,
                "required_fields": analysis.completeness.required_fields,
                "required_filled": analysis.completeness.required_filled,
                "optional_fields": analysis.completeness.optional_fields,
                "optional_filled": analysis.completeness.optional_filled,
                "completeness_pct": analysis.completeness.completeness_pct,
                "overall_pct": analysis.completeness.overall_pct,
                "missing_required": analysis.completeness.missing_required,
                "missing_optional": analysis.completeness.missing_optional,
            },
            "consistency_checks": [
                {
                    "check_name": c.check_name,
                    "description": c.description,
                    "passed": c.passed,
                    "severity": c.severity,
                    "detail": c.detail,
                }
                for c in analysis.consistency_checks
            ],
            "academic_performance": {
                "gpa": analysis.academic_performance.gpa,
                "grade": analysis.academic_performance.grade,
                "qualification": analysis.academic_performance.qualification,
                "programme": analysis.academic_performance.programme,
                "has_performance_data": analysis.academic_performance.has_performance_data,
                "summary": analysis.academic_performance.summary,
            },
            "anomalies": [
                {
                    "anomaly_type": a.anomaly_type,
                    "field_name": a.field_name,
                    "description": a.description,
                    "severity": a.severity,
                }
                for a in analysis.anomalies
            ],
            "recommendations": [
                {
                    "action": r.action,
                    "description": r.description,
                    "priority": r.priority,
                }
                for r in analysis.recommendations
            ],
            "field_analysis": [
                {
                    "field_name": fa.field_name,
                    "field_label": fa.field_label,
                    "value": fa.value,
                    "raw_value": fa.raw_value,
                    "confidence": fa.confidence,
                    "is_low_confidence": fa.is_low_confidence,
                    "is_present": fa.is_present,
                    "is_required": fa.is_required,
                    "is_valid": fa.is_valid,
                    "validation_message": fa.validation_message,
                    "was_corrected": fa.was_corrected,
                }
                for fa in analysis.fields
            ],
        },
    }


# ═══════════════════════════════════════════════════════════════
# Batch analytics with intelligence
# ═══════════════════════════════════════════════════════════════


@router.get("/batch/{batch_id}/analytics")
async def get_batch_analytics(
    batch_id: int,
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Get aggregate analytics for a batch of certificates.

    Returns batch-level intelligence including:
    - Completeness distribution (high/medium/low)
    - Average confidence and completeness
    - Anomaly summary
    - Institution and qualification breakdowns
    - Verification status distribution
    """
    from certificates.analysis import analyze_batch, batch_analytics

    org_id = get_current_organization_id(current_user, db)

    docs = list(
        db.execute(
            select(CaptureDocument)
            .where(
                CaptureDocument.organization_id == org_id,
                CaptureDocument.batch_id == batch_id,
                CaptureDocument.document_type.in_(CERTIFICATE_DOC_TYPES),
            )
            .order_by(CaptureDocument.id)
        )
        .scalars()
        .all()
    )

    if not docs:
        raise HTTPException(status_code=404, detail="No certificates found in this batch")

    doc_ids = [d.id for d in docs]
    all_fields = (
        db.execute(select(CaptureField).where(CaptureField.document_id.in_(doc_ids)))
        .scalars()
        .all()
    )
    fields_by_doc: dict[int, list[dict]] = {}
    for f in all_fields:
        fields_by_doc.setdefault(f.document_id, []).append({
            "field_name": f.field_name,
            "field_label": f.field_label,
            "data_type": f.data_type,
            "value": f.value,
            "raw_value": f.raw_value,
            "confidence_score": f.confidence_score,
            "is_low_confidence": f.is_low_confidence,
            "was_corrected": f.was_corrected,
            "is_valid": f.is_valid,
            "validation_message": f.validation_message,
        })

    doc_dicts = [_serialize_certificate(d) for d in docs]
    analyses = analyze_batch(doc_dicts, fields_by_doc)
    analytics = batch_analytics(analyses)

    return {
        "batch_id": batch_id,
        "total": analytics.total,
        "by_type": analytics.by_type,
        "by_verification": analytics.by_verification,
        "by_completeness_tier": analytics.by_completeness_tier,
        "avg_completeness": analytics.avg_completeness,
        "avg_confidence": analytics.avg_confidence,
        "total_anomalies": analytics.total_anomalies,
        "total_duplicates": analytics.total_duplicates,
        "common_anomalies": analytics.common_anomalies,
        "institutions": analytics.institutions,
        "qualifications": analytics.qualifications,
        "summary": analytics.summary,
        "certificates": [
            {
                "id": d.id,
                "filename": d.filename,
                "document_type": d.document_type,
                "document_type_label": d.document_type_label,
                "status": d.status,
                "verification_status": d.verification_status,
                "overall_confidence": d.overall_confidence,
                "duplicate_of_id": d.duplicate_of_id,
            }
            for d in docs
        ],
    }


# ═══════════════════════════════════════════════════════════════
# Field correction (human review)
# ═══════════════════════════════════════════════════════════════


@router.patch("/{document_id}/fields/{field_id}")
async def correct_certificate_field(
    document_id: int,
    field_id: int,
    payload: dict,
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Correct a single extracted field value during human review.

    The original OCR value is preserved in `raw_value`. The corrected value
    is stored in `value` with confidence set to 1.0 and `was_corrected` flag
    set. All corrections are logged in the audit trail.
    """
    org_id = get_current_organization_id(current_user, db)

    doc = db.execute(
        select(CaptureDocument).where(
            CaptureDocument.id == document_id,
            CaptureDocument.organization_id == org_id,
        )
    ).scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="Certificate not found")

    if not _is_certificate_type(doc.document_type):
        raise HTTPException(status_code=400, detail="Document is not a certificate type")

    field = db.execute(
        select(CaptureField).where(
            CaptureField.id == field_id,
            CaptureField.document_id == document_id,
        )
    ).scalar_one_or_none()
    if not field:
        raise HTTPException(status_code=404, detail="Field not found")

    new_value = payload.get("value")
    if new_value is None:
        raise HTTPException(status_code=422, detail="Field 'value' is required")

    old_value = field.value
    field.value = str(new_value)
    field.was_corrected = True
    field.confidence_score = 1.0
    field.is_low_confidence = False

    # Re-validate the corrected value
    from capture.document_types import get_document_type
    from capture.validators import validate_field

    doc_type_spec = get_document_type(doc.document_type) if doc.document_type else None
    enum_values = None
    if doc_type_spec:
        spec = next((f for f in doc_type_spec.fields if f.name == field.field_name), None)
        enum_values = spec.enum_values if spec else None
    is_valid, message = validate_field(field.field_name, str(new_value), field.data_type, enum_values)
    field.is_valid = is_valid
    field.validation_message = message

    # Record correction in the learning system
    from capture.models import CaptureCorrection

    correction = CaptureCorrection(
        document_id=document_id,
        field_id=field_id,
        field_name=field.field_name,
        old_value=old_value,
        new_value=str(new_value),
        corrected_by=current_user["id"],
    )
    db.add(correction)

    log_audit_event(
        db=db,
        action="certificate.field_corrected",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="certificate",
        resource_id=document_id,
        old_values={"field": field.field_name, "value": old_value},
        new_values={"field": field.field_name, "value": str(new_value)},
    )
    db.commit()

    return {
        "document_id": document_id,
        "field_id": field_id,
        "field_name": field.field_name,
        "old_value": old_value,
        "new_value": str(new_value),
        "confidence_score": field.confidence_score,
        "was_corrected": True,
        "is_valid": field.is_valid,
        "validation_message": field.validation_message,
    }


# ═══════════════════════════════════════════════════════════════
# Certificate report generation
# ═══════════════════════════════════════════════════════════════


@router.get("/report")
async def generate_certificate_report(
    certificate_type: str | None = Query(None),
    verification_status: str | None = Query(None),
    review_status: str | None = Query(None),
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Generate a structured JSON report from certificate data.

    The report includes:
    - Executive summary with key metrics
    - Breakdown by type, verification status, institution
    - Completeness and confidence statistics
    - Anomaly summary
    - Per-certificate analysis highlights

    This report can be consumed by the existing report engine for PDF
    generation or used standalone as a structured data export.
    """
    from certificates.analysis import analyze_batch, batch_analytics

    org_id = get_current_organization_id(current_user, db)

    query = (
        select(CaptureDocument)
        .where(
            CaptureDocument.organization_id == org_id,
            CaptureDocument.document_type.in_(CERTIFICATE_DOC_TYPES),
        )
        .order_by(CaptureDocument.id.desc())
    )
    if certificate_type:
        query = query.where(CaptureDocument.document_type == certificate_type)
    if verification_status:
        query = query.where(CaptureDocument.verification_status == verification_status)
    if review_status:
        query = query.where(CaptureDocument.status == review_status)

    docs = list(db.execute(query).scalars().all())
    if not docs:
        raise HTTPException(status_code=422, detail="No certificates found matching the criteria")

    doc_ids = [d.id for d in docs]
    all_fields = (
        db.execute(select(CaptureField).where(CaptureField.document_id.in_(doc_ids)))
        .scalars()
        .all()
    )
    fields_by_doc: dict[int, list[dict]] = {}
    for f in all_fields:
        fields_by_doc.setdefault(f.document_id, []).append({
            "field_name": f.field_name,
            "field_label": f.field_label,
            "data_type": f.data_type,
            "value": f.value,
            "raw_value": f.raw_value,
            "confidence_score": f.confidence_score,
            "is_low_confidence": f.is_low_confidence,
            "was_corrected": f.was_corrected,
            "is_valid": f.is_valid,
            "validation_message": f.validation_message,
        })

    doc_dicts = [_serialize_certificate(d) for d in docs]
    analyses = analyze_batch(doc_dicts, fields_by_doc)
    analytics = batch_analytics(analyses)

    report = {
        "title": "Certificate Intelligence Report",
        "organization_id": org_id,
        "generated_at": __import__("datetime").datetime.now(__import__("datetime").timezone.utc).isoformat(),
        "generated_by": current_user["id"],
        "filters": {
            "certificate_type": certificate_type,
            "verification_status": verification_status,
            "review_status": review_status,
        },
        "executive_summary": analytics.summary,
        "metrics": {
            "total_certificates": analytics.total,
            "avg_completeness": analytics.avg_completeness,
            "avg_confidence": analytics.avg_confidence,
            "total_anomalies": analytics.total_anomalies,
            "total_duplicates": analytics.total_duplicates,
        },
        "breakdowns": {
            "by_type": analytics.by_type,
            "by_verification": analytics.by_verification,
            "by_completeness_tier": analytics.by_completeness_tier,
            "by_institution": analytics.institutions,
            "by_qualification": analytics.qualifications,
        },
        "anomaly_summary": analytics.common_anomalies,
        "certificates": [
            {
                "id": d.id,
                "filename": d.filename,
                "document_type": d.document_type,
                "document_type_label": d.document_type_label,
                "status": d.status,
                "verification_status": d.verification_status,
                "overall_confidence": d.overall_confidence,
                "completeness_pct": analyses[i].completeness.completeness_pct,
                "anomaly_count": len(analyses[i].anomalies),
                "is_duplicate": analyses[i].is_duplicate,
                "summary": analyses[i].summary,
                "top_recommendation": (
                    {
                        "action": analyses[i].recommendations[0].action,
                        "description": analyses[i].recommendations[0].description,
                        "priority": analyses[i].recommendations[0].priority,
                    }
                    if analyses[i].recommendations
                    else None
                ),
            }
            for i, d in enumerate(docs)
        ],
    }

    log_audit_event(
        db=db,
        action="certificate.report_generated",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="certificate",
        new_values={"total": analytics.total, "filters": report["filters"]},
    )
    db.commit()

    return report


# ═══════════════════════════════════════════════════════════════
# Certificate PowerPoint generation
# ═══════════════════════════════════════════════════════════════


@router.get("/presentation")
async def generate_certificate_presentation(
    certificate_type: str | None = Query(None),
    verification_status: str | None = Query(None),
    review_status: str | None = Query(None),
    db: DbSession = Depends(get_db),
    current_user: dict = Depends(get_current_user),
):
    """Generate a PowerPoint presentation from certificate data.

    Creates a professional PPTX with:
    - Title slide
    - Executive summary slide
    - Distribution charts (by type, verification, completeness)
    - Institution and qualification breakdowns
    - Per-certificate highlights (top 10)

    Uses real data from the database — no fabricated content.
    """
    from certificates.analysis import analyze_batch, batch_analytics

    org_id = get_current_organization_id(current_user, db)

    query = (
        select(CaptureDocument)
        .where(
            CaptureDocument.organization_id == org_id,
            CaptureDocument.document_type.in_(CERTIFICATE_DOC_TYPES),
        )
        .order_by(CaptureDocument.id.desc())
    )
    if certificate_type:
        query = query.where(CaptureDocument.document_type == certificate_type)
    if verification_status:
        query = query.where(CaptureDocument.verification_status == verification_status)
    if review_status:
        query = query.where(CaptureDocument.status == review_status)

    docs = list(db.execute(query).scalars().all())
    if not docs:
        raise HTTPException(status_code=422, detail="No certificates found matching the criteria")

    doc_ids = [d.id for d in docs]
    all_fields = (
        db.execute(select(CaptureField).where(CaptureField.document_id.in_(doc_ids)))
        .scalars()
        .all()
    )
    fields_by_doc: dict[int, list[dict]] = {}
    for f in all_fields:
        fields_by_doc.setdefault(f.document_id, []).append({
            "field_name": f.field_name,
            "field_label": f.field_label,
            "data_type": f.data_type,
            "value": f.value,
            "raw_value": f.raw_value,
            "confidence_score": f.confidence_score,
            "is_low_confidence": f.is_low_confidence,
            "was_corrected": f.was_corrected,
            "is_valid": f.is_valid,
            "validation_message": f.validation_message,
        })

    doc_dicts = [_serialize_certificate(d) for d in docs]
    analyses = analyze_batch(doc_dicts, fields_by_doc)
    analytics = batch_analytics(analyses)

    # Build the PPTX
    try:
        from pptx import Presentation as PptxPresentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise HTTPException(
            status_code=500,
            detail="PowerPoint generation requires python-pptx. Install with: pip install python-pptx",
        ) from e

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)

    blank_layout = prs.slide_layouts[6]

    def _add_title_slide():
        slide = prs.slides.add_slide(blank_layout)
        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE
        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Certificate Intelligence Report"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = f"{analytics.total} certificates analyzed | Generated {__import__('datetime').datetime.now().strftime('%Y-%m-%d')}"
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xAE, 0xC7, 0xE0)

    def _add_summary_slide():
        slide = prs.slides.add_slide(blank_layout)
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Summary text
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = analytics.summary
        p2.font.size = Pt(18)
        p2.font.color.rgb = DARK_TEXT

        # Key metrics
        metrics_text = (
            f"Total Certificates: {analytics.total}\n"
            f"Average Completeness: {analytics.avg_completeness:.0f}%\n"
            f"Average Confidence: {analytics.avg_confidence:.0%}\n"
            f"Total Anomalies: {analytics.total_anomalies}\n"
            f"Duplicates: {analytics.total_duplicates}"
        )
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(6), Inches(3))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        for i, line in enumerate(metrics_text.split("\n")):
            p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT

        # Completeness tiers
        tier_text = (
            f"Completeness Distribution:\n"
            f"  High (≥80%): {analytics.by_completeness_tier.get('high', 0)}\n"
            f"  Medium (50-80%): {analytics.by_completeness_tier.get('medium', 0)}\n"
            f"  Low (<50%): {analytics.by_completeness_tier.get('low', 0)}"
        )
        txBox4 = slide.shapes.add_textbox(Inches(7), Inches(4), Inches(6), Inches(3))
        tf4 = txBox4.text_frame
        tf4.word_wrap = True
        for i, line in enumerate(tier_text.split("\n")):
            p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT

    def _add_chart_slide(title: str, data: dict[str, int], chart_type=XL_CHART_TYPE.BAR_CHART):
        slide = prs.slides.add_slide(blank_layout)
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        if not data:
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(1))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = "No data available"
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            return

        chart_data = CategoryChartData()
        chart_data.categories = list(data.keys())
        chart_data.add_series("Count", list(data.values()))

        slide.shapes.add_chart(
            chart_type,
            Inches(0.5),
            Inches(1.5),
            Inches(12),
            Inches(5.5),
            chart_data,
        )

    def _add_certificate_highlights_slide():
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Certificate Highlights"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Show top 10 certificates
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, (d, a) in enumerate(zip(docs[:10], analyses[:10], strict=False)):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            name = ""
            for fa in a.fields:
                if fa.field_name == "full_name" and fa.value:
                    name = fa.value
                    break
            p.text = (
                f"{i + 1}. {d.filename} — {a.document_type_label or 'Unknown'}"
                f" | Holder: {name or 'N/A'}"
                f" | Completeness: {a.completeness.completeness_pct:.0f}%"
                f" | Verification: {a.verification_status}"
            )
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT

    # Build slides
    _add_title_slide()
    _add_summary_slide()

    if analytics.by_type:
        _add_chart_slide("Certificates by Type", analytics.by_type, XL_CHART_TYPE.BAR_CHART)

    if analytics.by_verification:
        _add_chart_slide("Verification Status", analytics.by_verification, XL_CHART_TYPE.PIE_CHART)

    if analytics.institutions:
        # Top 10 institutions
        top_inst = dict(sorted(analytics.institutions.items(), key=lambda kv: kv[1], reverse=True)[:10])
        _add_chart_slide("Top Institutions", top_inst, XL_CHART_TYPE.BAR_CHART)

    if analytics.qualifications:
        top_qual = dict(sorted(analytics.qualifications.items(), key=lambda kv: kv[1], reverse=True)[:10])
        _add_chart_slide("Qualifications", top_qual, XL_CHART_TYPE.BAR_CHART)

    _add_certificate_highlights_slide()

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    log_audit_event(
        db=db,
        action="certificate.presentation_generated",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="certificate",
        new_values={"total": analytics.total},
    )
    db.commit()

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=certificate_intelligence.pptx"},
    )
