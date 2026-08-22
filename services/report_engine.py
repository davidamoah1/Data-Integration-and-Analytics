"""Reporting & Presentation Engine — Phase 8.

Builds executive-ready reports with:
  - Professional PDF export (cover page, TOC, sections, charts, tables)
  - PowerPoint-style presentations (slide generation from report sections)
  - Executive summaries (AI-generated or template-based)
  - Charts, tables, insights, and recommendations

Report structure:
  Report
    ├── Cover Page (title, org, date, author)
    ├── Executive Summary
    ├── Sections
    │     ├── Key Metrics (KPI cards)
    │     ├── Charts (visualizations)
    │     ├── Data Tables
    │     ├── Insights (AI-generated)
    │     └── Recommendations
    └── Appendix

Workflow: Dataset → Analysis → Insights → Report → Presentation
"""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass, field
from datetime import datetime, timezone
from enum import Enum
from typing import Any

logger = logging.getLogger(__name__)


# ── Enums ─────────────────────────────────────────────


class ReportSectionType(str, Enum):
    COVER = "cover"
    EXECUTIVE_SUMMARY = "executive_summary"
    KEY_METRICS = "key_metrics"
    CHART = "chart"
    TABLE = "table"
    INSIGHTS = "insights"
    RECOMMENDATIONS = "recommendations"
    METHODOLOGY = "methodology"
    APPENDIX = "appendix"
    CUSTOM = "custom"


class ReportTemplate(str, Enum):
    EXECUTIVE = "executive"
    ANALYTICAL = "analytical"
    RESEARCH = "research"
    OPERATIONAL = "operational"
    COMPLIANCE = "compliance"


class ExportFormat(str, Enum):
    PDF = "pdf"
    PPTX = "pptx"
    HTML = "html"
    JSON = "json"


class ChartType(str, Enum):
    BAR = "bar"
    LINE = "line"
    PIE = "pie"
    DONUT = "donut"
    AREA = "area"
    SCATTER = "scatter"
    HEATMAP = "heatmap"
    GAUGE = "gauge"


# ── Data Classes ──────────────────────────────────────


@dataclass
class KPIMetric:
    label: str
    value: str | int | float
    unit: str = ""
    trend: str = ""  # "up", "down", "flat"
    trend_value: str = ""
    target: str | int | float | None = None
    color: str = ""  # "green", "red", "amber", "blue"


@dataclass
class ChartDefinition:
    title: str
    chart_type: ChartType
    data: list[dict[str, Any]] = field(default_factory=list)
    x_axis: str = ""
    y_axis: str = ""
    series: list[dict[str, Any]] = field(default_factory=list)
    config: dict[str, Any] = field(default_factory=dict)

    @classmethod
    def from_canonical_spec(cls, spec: Any) -> ChartDefinition:
        """Create a ChartDefinition from a canonical ChartSpecification.

        This bridges the canonical chart spec (services.auto.chart_specification)
        to the report engine's ChartDefinition, ensuring reports use the
        SAME chart specifications as dashboards and presentations.
        """
        # Map canonical chart types to report ChartType enum
        type_map = {
            "bar_chart": ChartType.BAR,
            "horizontal_bar": ChartType.BAR,
            "line_chart": ChartType.LINE,
            "area_chart": ChartType.AREA,
            "pie_chart": ChartType.PIE,
            "donut_chart": ChartType.DONUT,
            "scatter_plot": ChartType.SCATTER,
            "histogram": ChartType.BAR,  # histogram renders as bar
            "box_plot": ChartType.BAR,  # box plot renders as bar variant
            "heatmap": ChartType.HEATMAP,
            "geo_map": ChartType.BAR,  # geo renders as bar variant in report
            "treemap": ChartType.BAR,  # treemap renders as bar variant in report
            "leaderboard": ChartType.BAR,
        }
        ct = type_map.get(spec.chart_type, ChartType.BAR)
        return cls(
            title=spec.title,
            chart_type=ct,
            data=spec.data[:100],  # cap for report
            x_axis=spec.x_axis or "",
            y_axis=spec.y_axis or "",
            series=spec.series,
            config={
                "aggregation": spec.aggregation,
                "source_columns": spec.source_columns,
                "importance_score": spec.importance_score,
                "reason": spec.reason,
                "source_analysis": spec.source_analysis,
                "chart_id": spec.id,
                "original_chart_type": spec.chart_type,
            },
        )


@dataclass
class TableDefinition:
    title: str
    columns: list[str] = field(default_factory=list)
    rows: list[list[Any]] = field(default_factory=list)
    summary: str = ""


@dataclass
class Insight:
    title: str
    description: str
    severity: str = "info"  # "info", "warning", "critical", "positive"
    metric: str = ""
    impact: str = ""


@dataclass
class Recommendation:
    title: str
    description: str
    priority: str = "medium"  # "high", "medium", "low"
    action: str = ""
    expected_impact: str = ""
    timeline: str = ""


@dataclass
class ReportSection:
    section_type: ReportSectionType
    title: str
    content: str = ""
    kpis: list[KPIMetric] = field(default_factory=list)
    charts: list[ChartDefinition] = field(default_factory=list)
    tables: list[TableDefinition] = field(default_factory=list)
    insights: list[Insight] = field(default_factory=list)
    recommendations: list[Recommendation] = field(default_factory=list)
    order: int = 0
    page_break: bool = False


@dataclass
class ReportComposition:
    report_id: str
    title: str
    subtitle: str = ""
    organization_name: str = ""
    author_name: str = ""
    template: ReportTemplate = ReportTemplate.EXECUTIVE
    industry: str = ""
    dataset_id: int | None = None
    analysis_id: int | None = None
    sections: list[ReportSection] = field(default_factory=list)
    created_at: str = ""
    tags: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "report_id": self.report_id,
            "title": self.title,
            "subtitle": self.subtitle,
            "organization_name": self.organization_name,
            "author_name": self.author_name,
            "template": self.template.value,
            "industry": self.industry,
            "dataset_id": self.dataset_id,
            "analysis_id": self.analysis_id,
            "sections": [
                {
                    "section_type": s.section_type.value,
                    "title": s.title,
                    "content": s.content,
                    "kpis": [
                        {
                            "label": k.label,
                            "value": k.value,
                            "unit": k.unit,
                            "trend": k.trend,
                            "trend_value": k.trend_value,
                            "target": k.target,
                            "color": k.color,
                        }
                        for k in s.kpis
                    ],
                    "charts": [
                        {
                            "title": c.title,
                            "chart_type": c.chart_type.value,
                            "data": c.data,
                            "x_axis": c.x_axis,
                            "y_axis": c.y_axis,
                            "series": c.series,
                            "config": c.config,
                        }
                        for c in s.charts
                    ],
                    "tables": [
                        {
                            "title": t.title,
                            "columns": t.columns,
                            "rows": t.rows,
                            "summary": t.summary,
                        }
                        for t in s.tables
                    ],
                    "insights": [
                        {
                            "title": i.title,
                            "description": i.description,
                            "severity": i.severity,
                            "metric": i.metric,
                            "impact": i.impact,
                        }
                        for i in s.insights
                    ],
                    "recommendations": [
                        {
                            "title": r.title,
                            "description": r.description,
                            "priority": r.priority,
                            "action": r.action,
                            "expected_impact": r.expected_impact,
                            "timeline": r.timeline,
                        }
                        for r in s.recommendations
                    ],
                    "order": s.order,
                    "page_break": s.page_break,
                }
                for s in sorted(self.sections, key=lambda x: x.order)
            ],
            "created_at": self.created_at,
            "tags": self.tags,
        }


# ── Report Templates ──────────────────────────────────


class ReportTemplateFactory:
    """Pre-built report templates by industry and use case."""

    @staticmethod
    def executive_template(
        title: str,
        org_name: str = "",
        author: str = "",
        industry: str = "",
    ) -> ReportComposition:
        return ReportComposition(
            report_id=f"rpt_exec_{datetime.now(timezone.utc).strftime('%Y%m%d%H%M%S')}",
            title=title,
            subtitle="Executive Summary Report",
            organization_name=org_name,
            author_name=author,
            template=ReportTemplate.EXECUTIVE,
            industry=industry,
            created_at=datetime.now(timezone.utc).isoformat(),
            sections=[
                ReportSection(
                    section_type=ReportSectionType.COVER,
                    title=title,
                    content=f"{org_name} — Executive Report",
                    order=0,
                    page_break=True,
                ),
                ReportSection(
                    section_type=ReportSectionType.EXECUTIVE_SUMMARY,
                    title="Executive Summary",
                    content="This report provides a comprehensive overview of key metrics, trends, and recommendations.",
                    order=1,
                ),
                ReportSection(
                    section_type=ReportSectionType.KEY_METRICS,
                    title="Key Performance Indicators",
                    order=2,
                ),
                ReportSection(
                    section_type=ReportSectionType.CHART,
                    title="Performance Trends",
                    order=3,
                ),
                ReportSection(
                    section_type=ReportSectionType.TABLE,
                    title="Detailed Breakdown",
                    order=4,
                ),
                ReportSection(
                    section_type=ReportSectionType.INSIGHTS,
                    title="Key Insights",
                    order=5,
                ),
                ReportSection(
                    section_type=ReportSectionType.RECOMMENDATIONS,
                    title="Strategic Recommendations",
                    order=6,
                    page_break=True,
                ),
            ],
        )

    @staticmethod
    def analytical_template(
        title: str,
        org_name: str = "",
        author: str = "",
        industry: str = "",
    ) -> ReportComposition:
        return ReportComposition(
            report_id=f"rpt_anal_{datetime.now(timezone.utc).strftime('%Y%m%d%H%M%S')}",
            title=title,
            subtitle="Analytical Report",
            organization_name=org_name,
            author_name=author,
            template=ReportTemplate.ANALYTICAL,
            industry=industry,
            created_at=datetime.now(timezone.utc).isoformat(),
            sections=[
                ReportSection(
                    section_type=ReportSectionType.COVER,
                    title=title,
                    content=f"{org_name} — Analytical Report",
                    order=0,
                    page_break=True,
                ),
                ReportSection(
                    section_type=ReportSectionType.EXECUTIVE_SUMMARY,
                    title="Overview",
                    order=1,
                ),
                ReportSection(
                    section_type=ReportSectionType.METHODOLOGY,
                    title="Methodology",
                    content="Description of data sources, analysis methods, and assumptions.",
                    order=2,
                ),
                ReportSection(
                    section_type=ReportSectionType.KEY_METRICS,
                    title="Descriptive Statistics",
                    order=3,
                ),
                ReportSection(
                    section_type=ReportSectionType.CHART,
                    title="Trend Analysis",
                    order=4,
                ),
                ReportSection(
                    section_type=ReportSectionType.TABLE,
                    title="Statistical Summary",
                    order=5,
                ),
                ReportSection(
                    section_type=ReportSectionType.INSIGHTS,
                    title="Findings",
                    order=6,
                ),
                ReportSection(
                    section_type=ReportSectionType.RECOMMENDATIONS,
                    title="Conclusions & Recommendations",
                    order=7,
                    page_break=True,
                ),
                ReportSection(
                    section_type=ReportSectionType.APPENDIX,
                    title="Appendix — Raw Data",
                    order=8,
                    page_break=True,
                ),
            ],
        )

    @staticmethod
    def research_template(
        title: str,
        org_name: str = "",
        author: str = "",
        industry: str = "",
    ) -> ReportComposition:
        return ReportComposition(
            report_id=f"rpt_res_{datetime.now(timezone.utc).strftime('%Y%m%d%H%M%S')}",
            title=title,
            subtitle="Research Report",
            organization_name=org_name,
            author_name=author,
            template=ReportTemplate.RESEARCH,
            industry=industry,
            created_at=datetime.now(timezone.utc).isoformat(),
            sections=[
                ReportSection(
                    section_type=ReportSectionType.COVER,
                    title=title,
                    content=f"{org_name} — Research Report",
                    order=0,
                    page_break=True,
                ),
                ReportSection(
                    section_type=ReportSectionType.EXECUTIVE_SUMMARY,
                    title="Abstract",
                    order=1,
                ),
                ReportSection(
                    section_type=ReportSectionType.METHODOLOGY,
                    title="Research Methodology",
                    content="Research design, data collection methods, and analytical framework.",
                    order=2,
                ),
                ReportSection(
                    section_type=ReportSectionType.KEY_METRICS,
                    title="Key Statistics",
                    order=3,
                ),
                ReportSection(
                    section_type=ReportSectionType.CHART,
                    title="Results Visualization",
                    order=4,
                ),
                ReportSection(
                    section_type=ReportSectionType.TABLE,
                    title="Statistical Results",
                    order=5,
                ),
                ReportSection(
                    section_type=ReportSectionType.INSIGHTS,
                    title="Discussion",
                    order=6,
                ),
                ReportSection(
                    section_type=ReportSectionType.RECOMMENDATIONS,
                    title="Conclusions & Future Work",
                    order=7,
                    page_break=True,
                ),
                ReportSection(
                    section_type=ReportSectionType.APPENDIX,
                    title="Appendix",
                    order=8,
                    page_break=True,
                ),
            ],
        )

    @classmethod
    def create_template(
        cls,
        template: ReportTemplate,
        title: str,
        org_name: str = "",
        author: str = "",
        industry: str = "",
    ) -> ReportComposition:
        if template == ReportTemplate.EXECUTIVE:
            return cls.executive_template(title, org_name, author, industry)
        elif template == ReportTemplate.ANALYTICAL:
            return cls.analytical_template(title, org_name, author, industry)
        elif template == ReportTemplate.RESEARCH:
            return cls.research_template(title, org_name, author, industry)
        else:
            return cls.executive_template(title, org_name, author, industry)


# ── Report Composition Service ────────────────────────


class ReportCompositionService:
    """Service for composing, managing, and exporting reports."""

    _store: dict[str, ReportComposition] = {}

    @classmethod
    def create_report(
        cls,
        title: str,
        template: ReportTemplate = ReportTemplate.EXECUTIVE,
        org_name: str = "",
        author: str = "",
        industry: str = "",
        dataset_id: int | None = None,
        analysis_id: int | None = None,
    ) -> ReportComposition:
        report = ReportTemplateFactory.create_template(template, title, org_name, author, industry)
        report.dataset_id = dataset_id
        report.analysis_id = analysis_id
        cls._store[report.report_id] = report
        logger.info(f"Created report '{title}' with template {template.value}")
        return report

    @classmethod
    def get_report(cls, report_id: str) -> ReportComposition | None:
        return cls._store.get(report_id)

    @classmethod
    def list_reports(cls) -> list[dict[str, Any]]:
        return [
            {
                "report_id": r.report_id,
                "title": r.title,
                "subtitle": r.subtitle,
                "template": r.template.value,
                "industry": r.industry,
                "section_count": len(r.sections),
                "created_at": r.created_at,
                "tags": r.tags,
            }
            for r in cls._store.values()
        ]

    @classmethod
    def delete_report(cls, report_id: str) -> bool:
        if report_id in cls._store:
            del cls._store[report_id]
            return True
        return False

    @classmethod
    def add_section(cls, report_id: str, section: ReportSection) -> ReportComposition | None:
        report = cls._store.get(report_id)
        if not report:
            return None
        section.order = len(report.sections)
        report.sections.append(section)
        return report

    @classmethod
    def remove_section(cls, report_id: str, section_order: int) -> ReportComposition | None:
        report = cls._store.get(report_id)
        if not report:
            return None
        report.sections = [s for s in report.sections if s.order != section_order]
        for i, s in enumerate(report.sections):
            s.order = i
        return report

    @classmethod
    def update_section(
        cls, report_id: str, section_order: int, updates: dict[str, Any]
    ) -> ReportComposition | None:
        report = cls._store.get(report_id)
        if not report:
            return None
        for s in report.sections:
            if s.order == section_order:
                if "title" in updates:
                    s.title = updates["title"]
                if "content" in updates:
                    s.content = updates["content"]
                if "kpis" in updates:
                    s.kpis = [KPIMetric(**k) for k in updates["kpis"]]
                if "charts" in updates:
                    s.charts = [
                        ChartDefinition(
                            title=c["title"],
                            chart_type=ChartType(c["chart_type"]),
                            data=c.get("data", []),
                            x_axis=c.get("x_axis", ""),
                            y_axis=c.get("y_axis", ""),
                            series=c.get("series", []),
                            config=c.get("config", {}),
                        )
                        for c in updates["charts"]
                    ]
                if "tables" in updates:
                    s.tables = [TableDefinition(**t) for t in updates["tables"]]
                if "insights" in updates:
                    s.insights = [Insight(**i) for i in updates["insights"]]
                if "recommendations" in updates:
                    s.recommendations = [Recommendation(**r) for r in updates["recommendations"]]
                break
        return report

    @classmethod
    def add_kpis(
        cls, report_id: str, section_order: int, kpis: list[KPIMetric]
    ) -> ReportComposition | None:
        report = cls._store.get(report_id)
        if not report:
            return None
        for s in report.sections:
            if s.order == section_order:
                s.kpis.extend(kpis)
                break
        return report

    @classmethod
    def add_chart(
        cls, report_id: str, section_order: int, chart: ChartDefinition
    ) -> ReportComposition | None:
        report = cls._store.get(report_id)
        if not report:
            return None
        for s in report.sections:
            if s.order == section_order:
                s.charts.append(chart)
                break
        return report

    @classmethod
    def add_insights(
        cls, report_id: str, section_order: int, insights: list[Insight]
    ) -> ReportComposition | None:
        report = cls._store.get(report_id)
        if not report:
            return None
        for s in report.sections:
            if s.order == section_order:
                s.insights.extend(insights)
                break
        return report

    @classmethod
    def add_recommendations(
        cls, report_id: str, section_order: int, recommendations: list[Recommendation]
    ) -> ReportComposition | None:
        report = cls._store.get(report_id)
        if not report:
            return None
        for s in report.sections:
            if s.order == section_order:
                s.recommendations.extend(recommendations)
                break
        return report

    @classmethod
    def generate_executive_summary(cls, report: ReportComposition) -> str:
        """Generate an executive summary from the report's sections."""
        kpi_count = sum(len(s.kpis) for s in report.sections)
        chart_count = sum(len(s.charts) for s in report.sections)
        table_count = sum(len(s.tables) for s in report.sections)
        insight_count = sum(len(s.insights) for s in report.sections)
        rec_count = sum(len(s.recommendations) for s in report.sections)

        summary_parts = [
            f"This {report.template.value} report for {report.organization_name or 'the organization'}",
            f"contains {len(report.sections)} sections with {kpi_count} key metrics,",
            f"{chart_count} charts, {table_count} data tables,",
            f"{insight_count} insights, and {rec_count} recommendations.",
        ]

        # Add top insights
        all_insights = [i for s in report.sections for i in s.insights]
        if all_insights:
            critical = [i for i in all_insights if i.severity == "critical"]
            positive = [i for i in all_insights if i.severity == "positive"]
            if critical:
                summary_parts.append(
                    f"\n\nCritical findings: {len(critical)} item(s) require immediate attention."
                )
            if positive:
                summary_parts.append(
                    f"\nPositive trends: {len(positive)} item(s) show strong performance."
                )

        # Add top recommendations
        all_recs = [r for s in report.sections for r in s.recommendations]
        high_priority = [r for r in all_recs if r.priority == "high"]
        if high_priority:
            summary_parts.append(
                f"\n\n{len(high_priority)} high-priority recommendation(s) identified for immediate action."
            )

        return " ".join(summary_parts)

    @classmethod
    def populate_from_dashboard_spec(
        cls,
        report_id: str,
        dashboard_spec: Any,
    ) -> ReportComposition | None:
        """Populate report sections from a canonical DashboardSpecification.

        This ensures the report uses the SAME chart specifications as the
        dashboard and presentation — no independent chart regeneration.

        Args:
            report_id: The report ID to populate.
            dashboard_spec: DashboardSpecification from the Visualization Intelligence Engine.
        """
        report = cls._store.get(report_id)
        if not report:
            return None

        # Find the chart section and key metrics section
        chart_section = next(
            (s for s in report.sections if s.section_type == ReportSectionType.CHART), None
        )
        kpi_section = next(
            (s for s in report.sections if s.section_type == ReportSectionType.KEY_METRICS), None
        )
        insight_section = next(
            (s for s in report.sections if s.section_type == ReportSectionType.INSIGHTS), None
        )
        rec_section = next(
            (s for s in report.sections if s.section_type == ReportSectionType.RECOMMENDATIONS), None
        )

        # Convert canonical charts to report ChartDefinitions
        if chart_section:
            chart_section.charts = [
                ChartDefinition.from_canonical_spec(c) for c in dashboard_spec.charts
            ]

        # Convert canonical KPIs to report KPIMetrics
        if kpi_section:
            kpi_section.kpis = [
                KPIMetric(
                    label=k.label,
                    value=k.value,
                    unit=k.unit,
                    trend=k.comparison_direction or "",
                    trend_value=str(k.comparison_value) if k.comparison_value is not None else "",
                )
                for k in dashboard_spec.kpis
            ]

        # Convert canonical insights to report Insights
        if insight_section:
            insight_section.insights = [
                Insight(
                    title=i.title,
                    description=i.description,
                    severity=i.severity,
                    metric=i.metric,
                    impact=i.recommendation,
                )
                for i in dashboard_spec.insights
            ]

        # Convert recommendations
        if rec_section and dashboard_spec.recommendations:
            rec_section.recommendations = [
                Recommendation(
                    title=rec[:80] if len(rec) > 80 else rec,
                    description=rec,
                    priority="medium",
                )
                for rec in dashboard_spec.recommendations
            ]

        return report

    @classmethod
    def export_to_dict(cls, report: ReportComposition) -> dict[str, Any]:
        """Export report as a dictionary (for JSON export)."""
        result = report.to_dict()
        result["executive_summary"] = cls.generate_executive_summary(report)
        return result

    @classmethod
    def export_to_html(cls, report: ReportComposition) -> str:
        """Export report as a styled HTML document."""
        exec_summary = cls.generate_executive_summary(report)
        html_parts = [
            "<!DOCTYPE html>",
            '<html lang="en">',
            "<head>",
            '<meta charset="UTF-8">',
            f"<title>{report.title}</title>",
            "<style>",
            "body { font-family: 'Helvetica Neue', Arial, sans-serif; margin: 40px; color: #1a1a2e; }",
            ".cover { text-align: center; padding: 80px 0; page-break-after: always; }",
            ".cover h1 { font-size: 32px; color: #0f3460; margin-bottom: 10px; }",
            ".cover .subtitle { font-size: 18px; color: #666; }",
            ".cover .meta { margin-top: 30px; font-size: 14px; color: #999; }",
            "h2 { color: #0f3460; border-bottom: 2px solid #0f3460; padding-bottom: 5px; }",
            ".kpi-grid { display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 15px; }",
            ".kpi-card { border: 1px solid #e0e0e0; border-radius: 8px; padding: 15px; text-align: center; }",
            ".kpi-card .label { font-size: 12px; color: #666; text-transform: uppercase; }",
            ".kpi-card .value { font-size: 28px; font-weight: bold; color: #0f3460; }",
            ".kpi-card .trend { font-size: 12px; margin-top: 5px; }",
            ".trend.up { color: #16a34a; } .trend.down { color: #dc2626; }",
            "table { width: 100%; border-collapse: collapse; margin: 15px 0; }",
            "th, td { border: 1px solid #ddd; padding: 8px 12px; text-align: left; }",
            "th { background: #0f3460; color: white; }",
            "tr:nth-child(even) { background: #f8f9fa; }",
            ".insight { border-left: 4px solid #0f3460; padding: 10px 15px; margin: 10px 0; background: #f8f9fa; }",
            ".insight.critical { border-color: #dc2626; }",
            ".insight.warning { border-color: #f59e0b; }",
            ".insight.positive { border-color: #16a34a; }",
            ".recommendation { border: 1px solid #e0e0e0; border-radius: 8px; padding: 15px; margin: 10px 0; }",
            ".recommendation .priority { display: inline-block; padding: 2px 8px; border-radius: 4px; font-size: 11px; font-weight: bold; }",
            ".priority.high { background: #fee2e2; color: #dc2626; }",
            ".priority.medium { background: #fef3c7; color: #f59e0b; }",
            ".priority.low { background: #dbeafe; color: #2563eb; }",
            ".page-break { page-break-before: always; }",
            ".chart-placeholder { border: 2px dashed #ccc; border-radius: 8px; padding: 40px; text-align: center; color: #999; }",
            "</style>",
            "</head>",
            "<body>",
        ]

        for section in sorted(report.sections, key=lambda x: x.order):
            if section.page_break:
                html_parts.append('<div class="page-break"></div>')

            if section.section_type == ReportSectionType.COVER:
                html_parts.extend(
                    [
                        '<div class="cover">',
                        f"<h1>{section.title}</h1>",
                        f'<div class="subtitle">{report.subtitle}</div>',
                        '<div class="meta">',
                        f"<p>Organization: {report.organization_name or 'N/A'}</p>",
                        f"<p>Author: {report.author_name or 'N/A'}</p>",
                        f"<p>Date: {datetime.now().strftime('%B %d, %Y')}</p>",
                        "</div>",
                        "</div>",
                    ]
                )
                continue

            html_parts.append(f"<h2>{section.title}</h2>")

            if section.section_type == ReportSectionType.EXECUTIVE_SUMMARY:
                html_parts.append(f"<p>{exec_summary}</p>")
                if section.content:
                    html_parts.append(f"<p>{section.content}</p>")

            elif section.section_type == ReportSectionType.KEY_METRICS:
                if section.kpis:
                    html_parts.append('<div class="kpi-grid">')
                    for kpi in section.kpis:
                        trend_class = (
                            "up" if kpi.trend == "up" else "down" if kpi.trend == "down" else ""
                        )
                        html_parts.extend(
                            [
                                '<div class="kpi-card">',
                                f'<div class="label">{kpi.label}</div>',
                                f'<div class="value">{kpi.value}{kpi.unit}</div>',
                                (
                                    f'<div class="trend {trend_class}">{kpi.trend_value}</div>'
                                    if kpi.trend_value
                                    else ""
                                ),
                                "</div>",
                            ]
                        )
                    html_parts.append("</div>")
                if section.content:
                    html_parts.append(f"<p>{section.content}</p>")

            elif section.section_type == ReportSectionType.CHART:
                for chart in section.charts:
                    html_parts.extend(
                        [
                            f"<h3>{chart.title}</h3>",
                            '<div class="chart-placeholder">',
                            f"[{chart.chart_type.value.upper()} CHART: {chart.title}]",
                            f"<br/><small>X: {chart.x_axis} | Y: {chart.y_axis}</small>",
                            "</div>",
                        ]
                    )

            elif section.section_type == ReportSectionType.TABLE:
                for table in section.tables:
                    html_parts.append(f"<h3>{table.title}</h3>")
                    if table.columns:
                        html_parts.append("<table><thead><tr>")
                        for col in table.columns:
                            html_parts.append(f"<th>{col}</th>")
                        html_parts.append("</tr></thead><tbody>")
                        for row in table.rows:
                            html_parts.append("<tr>")
                            for cell in row:
                                html_parts.append(f"<td>{cell}</td>")
                            html_parts.append("</tr>")
                        html_parts.append("</tbody></table>")
                    if table.summary:
                        html_parts.append(f"<p><em>{table.summary}</em></p>")

            elif section.section_type == ReportSectionType.INSIGHTS:
                for insight in section.insights:
                    html_parts.extend(
                        [
                            f'<div class="insight {insight.severity}">',
                            f"<strong>{insight.title}</strong>",
                            f"<p>{insight.description}</p>",
                            f"<small>Impact: {insight.impact}</small>" if insight.impact else "",
                            "</div>",
                        ]
                    )

            elif section.section_type == ReportSectionType.RECOMMENDATIONS:
                for rec in section.recommendations:
                    html_parts.extend(
                        [
                            '<div class="recommendation">',
                            f'<span class="priority {rec.priority}">{rec.priority.upper()}</span>',
                            f"<strong> {rec.title}</strong>",
                            f"<p>{rec.description}</p>",
                            f"<p><small>Action: {rec.action}</small></p>" if rec.action else "",
                            (
                                f"<p><small>Expected Impact: {rec.expected_impact}</small></p>"
                                if rec.expected_impact
                                else ""
                            ),
                            (
                                f"<p><small>Timeline: {rec.timeline}</small></p>"
                                if rec.timeline
                                else ""
                            ),
                            "</div>",
                        ]
                    )

            else:
                if section.content:
                    html_parts.append(f"<p>{section.content}</p>")

        html_parts.extend(["</body>", "</html>"])
        return "\n".join(html_parts)

    @classmethod
    def export_to_pdf(cls, report: ReportComposition) -> bytes:
        """Export report as a professional PDF document."""
        from fpdf import FPDF

        pdf = FPDF(format="A4")
        pdf.set_auto_page_break(auto=True, margin=15)

        page_width = pdf.w - pdf.l_margin - pdf.r_margin

        for section in sorted(report.sections, key=lambda x: x.order):
            if section.page_break and pdf.page_no() > 0 or pdf.page_no() == 0:
                pdf.add_page()

            # Cover page
            if section.section_type == ReportSectionType.COVER:
                pdf.ln(60)
                pdf.set_font("Helvetica", "B", 24)
                pdf.cell(0, 15, section.title, new_x="LMARGIN", new_y="NEXT", align="C")
                pdf.ln(5)
                pdf.set_font("Helvetica", "", 14)
                pdf.cell(0, 8, report.subtitle, new_x="LMARGIN", new_y="NEXT", align="C")
                pdf.ln(20)
                pdf.set_font("Helvetica", "", 11)
                pdf.cell(
                    0,
                    6,
                    f"Organization: {report.organization_name or 'N/A'}",
                    new_x="LMARGIN",
                    new_y="NEXT",
                    align="C",
                )
                pdf.cell(
                    0,
                    6,
                    f"Author: {report.author_name or 'N/A'}",
                    new_x="LMARGIN",
                    new_y="NEXT",
                    align="C",
                )
                pdf.cell(
                    0,
                    6,
                    f"Date: {datetime.now().strftime('%B %d, %Y')}",
                    new_x="LMARGIN",
                    new_y="NEXT",
                    align="C",
                )
                continue

            # Section header
            pdf.set_font("Helvetica", "B", 14)
            pdf.cell(0, 10, section.title, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)

            # Executive summary
            if section.section_type == ReportSectionType.EXECUTIVE_SUMMARY:
                exec_summary = cls.generate_executive_summary(report)
                pdf.set_font("Helvetica", "", 10)
                pdf.multi_cell(page_width, 5, exec_summary)
                pdf.ln(3)
                if section.content:
                    pdf.multi_cell(page_width, 5, section.content)

            # Key metrics
            elif section.section_type == ReportSectionType.KEY_METRICS:
                if section.kpis:
                    col_count = min(len(section.kpis), 3)
                    col_width = page_width / col_count
                    for i in range(0, len(section.kpis), col_count):
                        row = section.kpis[i : i + col_count]
                        for kpi in row:
                            pdf.set_font("Helvetica", "B", 9)
                            pdf.cell(col_width, 5, kpi.label[:20], border=1, align="C")
                        pdf.ln()
                        for kpi in row:
                            pdf.set_font("Helvetica", "B", 16)
                            val = f"{kpi.value}{kpi.unit}"
                            pdf.cell(col_width, 8, val, border=1, align="C")
                        pdf.ln()
                        for kpi in row:
                            pdf.set_font("Helvetica", "", 8)
                            trend = kpi.trend_value or kpi.trend or ""
                            pdf.cell(col_width, 4, trend, border=1, align="C")
                        pdf.ln(2)
                if section.content:
                    pdf.set_font("Helvetica", "", 10)
                    pdf.multi_cell(page_width, 5, section.content)

            # Charts
            elif section.section_type == ReportSectionType.CHART:
                for chart in section.charts:
                    pdf.set_font("Helvetica", "B", 11)
                    pdf.cell(0, 6, chart.title, new_x="LMARGIN", new_y="NEXT")
                    pdf.set_font("Helvetica", "", 9)
                    pdf.cell(
                        0,
                        5,
                        f"[{chart.chart_type.value.upper()} Chart - X: {chart.x_axis}, Y: {chart.y_axis}]",
                        new_x="LMARGIN",
                        new_y="NEXT",
                    )
                    pdf.ln(2)

            # Tables
            elif section.section_type == ReportSectionType.TABLE:
                for table in section.tables:
                    pdf.set_font("Helvetica", "B", 11)
                    pdf.cell(0, 6, table.title, new_x="LMARGIN", new_y="NEXT")
                    if table.columns:
                        pdf.set_font("Helvetica", "B", 9)
                        col_width = page_width / len(table.columns)
                        for col in table.columns:
                            pdf.cell(col_width, 6, str(col)[:18], border=1, align="C")
                        pdf.ln()
                        pdf.set_font("Helvetica", "", 8)
                        for row in table.rows[:30]:
                            for cell in row:
                                pdf.cell(col_width, 5, str(cell)[:18], border=1)
                            pdf.ln()
                    if table.summary:
                        pdf.set_font("Helvetica", "I", 9)
                        pdf.multi_cell(page_width, 4, table.summary)
                    pdf.ln(2)

            # Insights
            elif section.section_type == ReportSectionType.INSIGHTS:
                for insight in section.insights:
                    severity_marker = {
                        "critical": "[!]",
                        "warning": "[~]",
                        "positive": "[+]",
                        "info": "[i]",
                    }.get(insight.severity, "[i]")
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.cell(
                        0, 6, f"{severity_marker} {insight.title}", new_x="LMARGIN", new_y="NEXT"
                    )
                    pdf.set_font("Helvetica", "", 9)
                    pdf.multi_cell(page_width, 4, insight.description)
                    if insight.impact:
                        pdf.set_font("Helvetica", "I", 8)
                        pdf.cell(0, 4, f"Impact: {insight.impact}", new_x="LMARGIN", new_y="NEXT")
                    pdf.ln(2)

            # Recommendations
            elif section.section_type == ReportSectionType.RECOMMENDATIONS:
                for rec in section.recommendations:
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.cell(
                        0, 6, f"[{rec.priority.upper()}] {rec.title}", new_x="LMARGIN", new_y="NEXT"
                    )
                    pdf.set_font("Helvetica", "", 9)
                    pdf.multi_cell(page_width, 4, rec.description)
                    if rec.action:
                        pdf.set_font("Helvetica", "", 8)
                        pdf.multi_cell(page_width, 4, f"Action: {rec.action}")
                    if rec.expected_impact:
                        pdf.cell(
                            0,
                            4,
                            f"Expected Impact: {rec.expected_impact}",
                            new_x="LMARGIN",
                            new_y="NEXT",
                        )
                    if rec.timeline:
                        pdf.cell(0, 4, f"Timeline: {rec.timeline}", new_x="LMARGIN", new_y="NEXT")
                    pdf.ln(2)

            else:
                if section.content:
                    pdf.set_font("Helvetica", "", 10)
                    pdf.multi_cell(page_width, 5, section.content)

        return bytes(pdf.output())

    @classmethod
    def export_to_pptx(cls, report: ReportComposition) -> bytes:
        """Export report as a PowerPoint-style presentation."""
        try:
            from pptx import Presentation as PptxPresentation
            from pptx.dml.color import RGBColor
            from pptx.util import Inches, Pt
        except ImportError:
            logger.warning("python-pptx not installed, generating JSON fallback")
            return cls.export_to_dict(report).__str__().encode("utf-8")

        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Color scheme
        PRIMARY = RGBColor(0x0F, 0x34, 0x60)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        DARK = RGBColor(0x1A, 0x1A, 0x2E)
        GRAY = RGBColor(0x66, 0x66, 0x66)
        GREEN = RGBColor(0x16, 0xA3, 0x4A)
        RED = RGBColor(0xDC, 0x26, 0x26)
        RGBColor(0xF5, 0x9E, 0x0B)

        def add_title_slide(title: str, subtitle: str):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            # Background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = PRIMARY
            # Title
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = WHITE
            # Subtitle
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = WHITE

        def add_content_slide(title: str, bullets: list[str]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Title bar
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = PRIMARY
            # Bullets
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(16)
                p.font.color.rgb = DARK
                p.space_after = Pt(8)

        def add_kpi_slide(title: str, kpis: list[KPIMetric]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = PRIMARY

            col_count = min(len(kpis), 4)
            col_width = 12 / col_count
            for i, kpi in enumerate(kpis[:8]):
                col = i % col_count
                row = i // col_count
                x = Inches(0.5 + col * col_width)
                y = Inches(1.5 + row * 2.5)
                w = Inches(col_width - 0.3)
                h = Inches(2)
                txBox = slide.shapes.add_textbox(x, y, w, h)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = kpi.label
                p.font.size = Pt(12)
                p.font.color.rgb = GRAY
                p2 = tf.add_paragraph()
                p2.text = f"{kpi.value}{kpi.unit}"
                p2.font.size = Pt(32)
                p2.font.bold = True
                p2.font.color.rgb = PRIMARY
                if kpi.trend_value:
                    p3 = tf.add_paragraph()
                    p3.text = kpi.trend_value
                    p3.font.size = Pt(11)
                    p3.font.color.rgb = (
                        GREEN if kpi.trend == "up" else RED if kpi.trend == "down" else GRAY
                    )

        def add_table_slide(title: str, table: TableDefinition):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = PRIMARY

            rows = min(len(table.rows) + 1, 10)
            cols = min(len(table.columns), 6)
            if cols > 0 and rows > 1:
                tbl = slide.shapes.add_table(
                    rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(5)
                ).table
                for j, col_name in enumerate(table.columns[:cols]):
                    tbl.cell(0, j).text = col_name
                    tbl.cell(0, j).text_frame.paragraphs[0].font.bold = True
                for i, row in enumerate(table.rows[: rows - 1]):
                    for j, cell in enumerate(row[:cols]):
                        tbl.cell(i + 1, j).text = str(cell)[:30]

        # Generate slides from report sections
        for section in sorted(report.sections, key=lambda x: x.order):
            if section.section_type == ReportSectionType.COVER:
                add_title_slide(section.title, report.subtitle or report.organization_name)

            elif section.section_type == ReportSectionType.EXECUTIVE_SUMMARY:
                exec_summary = cls.generate_executive_summary(report)
                bullets = exec_summary.split(". ")
                add_content_slide("Executive Summary", [b.strip() for b in bullets if b.strip()])

            elif section.section_type == ReportSectionType.KEY_METRICS:
                if section.kpis:
                    add_kpi_slide(section.title, section.kpis)
                else:
                    add_content_slide(
                        section.title, ["Key metrics will appear here once data is loaded."]
                    )

            elif section.section_type == ReportSectionType.CHART:
                chart_bullets = [f"{c.title} ({c.chart_type.value})" for c in section.charts]
                if not chart_bullets:
                    chart_bullets = ["Charts will appear here once data is loaded."]
                add_content_slide(section.title, chart_bullets)

            elif section.section_type == ReportSectionType.TABLE:
                for table in section.tables:
                    add_table_slide(table.title, table)
                if not section.tables:
                    add_content_slide(
                        section.title, ["Data tables will appear here once data is loaded."]
                    )

            elif section.section_type == ReportSectionType.INSIGHTS:
                bullets = []
                for insight in section.insights:
                    prefix = {
                        "critical": "⚠️",
                        "warning": "⚡",
                        "positive": "✅",
                        "info": "ℹ️",
                    }.get(insight.severity, "•")
                    bullets.append(f"{prefix} {insight.title}: {insight.description}")
                if not bullets:
                    bullets = ["Insights will appear here after analysis."]
                add_content_slide(section.title, bullets)

            elif section.section_type == ReportSectionType.RECOMMENDATIONS:
                bullets = []
                for rec in section.recommendations:
                    bullets.append(f"[{rec.priority.upper()}] {rec.title}: {rec.description}")
                if not bullets:
                    bullets = ["Recommendations will appear here after analysis."]
                add_content_slide(section.title, bullets)

            else:
                if section.content:
                    add_content_slide(section.title, [section.content])

        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    @classmethod
    def export_report(cls, report_id: str, format: ExportFormat) -> tuple[bytes, str, str]:
        """Export a report in the specified format."""
        report = cls._store.get(report_id)
        if not report:
            raise ValueError(f"Report not found: {report_id}")

        if format == ExportFormat.JSON:
            data = cls.export_to_dict(report)
            return (
                __import__("json").dumps(data, indent=2, default=str).encode("utf-8"),
                "application/json",
                "json",
            )
        elif format == ExportFormat.HTML:
            html = cls.export_to_html(report)
            return html.encode("utf-8"), "text/html", "html"
        elif format == ExportFormat.PDF:
            pdf_bytes = cls.export_to_pdf(report)
            return pdf_bytes, "application/pdf", "pdf"
        elif format == ExportFormat.PPTX:
            pptx_bytes = cls.export_to_pptx(report)
            return (
                pptx_bytes,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "pptx",
            )
        else:
            raise ValueError(f"Unsupported format: {format}")


# ── Presentation Generator ────────────────────────────


class PresentationGenerator:
    """Generates PowerPoint-style presentations from report compositions."""

    @staticmethod
    def from_report(report: ReportComposition) -> list[dict[str, Any]]:
        """Generate slide definitions from a report."""
        slides: list[dict[str, Any]] = []
        svc = ReportCompositionService

        for section in sorted(report.sections, key=lambda x: x.order):
            if section.section_type == ReportSectionType.COVER:
                slides.append(
                    {
                        "slide_number": len(slides) + 1,
                        "layout": "title",
                        "title": section.title,
                        "subtitle": report.subtitle or report.organization_name,
                        "speaker_notes": f"Welcome. Today we'll cover {report.title}.",
                    }
                )

            elif section.section_type == ReportSectionType.EXECUTIVE_SUMMARY:
                exec_summary = svc.generate_executive_summary(report)
                slides.append(
                    {
                        "slide_number": len(slides) + 1,
                        "layout": "bullets",
                        "title": "Executive Summary",
                        "content": exec_summary,
                        "speaker_notes": "Provide a high-level overview of the main findings.",
                    }
                )

            elif section.section_type == ReportSectionType.KEY_METRICS and section.kpis:
                slides.append(
                    {
                        "slide_number": len(slides) + 1,
                        "layout": "kpi",
                        "title": section.title,
                        "kpis": [
                            {
                                "label": k.label,
                                "value": f"{k.value}{k.unit}",
                                "trend": k.trend,
                                "trend_value": k.trend_value,
                            }
                            for k in section.kpis
                        ],
                        "speaker_notes": "Walk through the most important metrics.",
                    }
                )

            elif section.section_type == ReportSectionType.CHART and section.charts:
                for chart in section.charts:
                    slides.append(
                        {
                            "slide_number": len(slides) + 1,
                            "layout": "chart",
                            "title": chart.title,
                            "chart_type": chart.chart_type.value,
                            "chart_data": chart.data,
                            "x_axis": chart.x_axis,
                            "y_axis": chart.y_axis,
                            "speaker_notes": f"Discuss the {chart.chart_type.value} chart showing {chart.title}.",
                        }
                    )

            elif section.section_type == ReportSectionType.TABLE and section.tables:
                for table in section.tables:
                    slides.append(
                        {
                            "slide_number": len(slides) + 1,
                            "layout": "table",
                            "title": table.title,
                            "columns": table.columns,
                            "rows": table.rows[:15],
                            "speaker_notes": f"Review the data in {table.title}.",
                        }
                    )

            elif section.section_type == ReportSectionType.INSIGHTS and section.insights:
                slides.append(
                    {
                        "slide_number": len(slides) + 1,
                        "layout": "bullets",
                        "title": section.title,
                        "content": "\n".join(
                            [
                                f"{'⚠️' if i.severity == 'critical' else '✅' if i.severity == 'positive' else '•'} {i.title}: {i.description}"
                                for i in section.insights
                            ]
                        ),
                        "speaker_notes": "Discuss each insight with supporting evidence.",
                    }
                )

            elif (
                section.section_type == ReportSectionType.RECOMMENDATIONS
                and section.recommendations
            ):
                slides.append(
                    {
                        "slide_number": len(slides) + 1,
                        "layout": "bullets",
                        "title": section.title,
                        "content": "\n".join(
                            [
                                f"[{r.priority.upper()}] {r.title}: {r.description}"
                                for r in section.recommendations
                            ]
                        ),
                        "speaker_notes": "Present clear, actionable recommendations.",
                    }
                )

        # Closing slide
        slides.append(
            {
                "slide_number": len(slides) + 1,
                "layout": "title",
                "title": "Thank You",
                "subtitle": "Questions & Discussion",
                "speaker_notes": "Open the floor for questions.",
            }
        )

        return slides
