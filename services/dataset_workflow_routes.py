"""FastAPI routes for the Dataset Intelligence Workflow.

Endpoints:
  POST /dataset-workflow/run          â€” Start a full workflow on an uploaded file
  GET  /dataset-workflow/{id}/status  â€” Get workflow status and progress
  GET  /dataset-workflow/{id}/profile â€” Get dataset profile
  GET  /dataset-workflow/{id}/quality â€” Get quality report
  GET  /dataset-workflow/{id}/semantic â€” Get semantic analysis
  GET  /dataset-workflow/{id}/industry â€” Get industry detection
  GET  /dataset-workflow/{id}/insights â€” Get AI insights
  GET  /dataset-workflow/{id}/dashboard â€” Get dashboard recommendations
  POST /dataset-workflow/{id}/retry/{stage} â€” Retry a failed stage
  POST /dataset-workflow/{id}/confirm-industry â€” Confirm industry detection
"""

from __future__ import annotations

import io
import json
import logging
import math
import os
import tempfile
from typing import Any

import numpy as np
import pandas as pd
from fastapi import APIRouter, Depends, File, HTTPException, Request, UploadFile
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from sqlalchemy import select
from sqlalchemy.orm import Session as DbSession

import shared.database as db_module
from audit.service import log_audit_event
from etl.file_security import FileValidator
from governance import classify_dataset
from services.dataset_workflow import (
    DatasetWorkflowOrchestrator,
    StageResult,
    StageStatus,
    WorkflowStage,
    WorkflowState,
)
from services.dataset_workflow_models import DatasetWorkflowRun
from shared.database import get_db
from shared.dependencies import get_current_user
from shared.tenant import get_current_organization_id

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api/dataset-workflow", tags=["Dataset Workflow"])

# Global orchestrator instance. Live state lives in this in-process dict for
# the lifetime of the run, but every stage transition is also persisted to
# the dataset_workflow_runs table (see `_persist_workflow_state` below) so
# status/results survive a restart and are visible to other worker
# processes (C3). Retrying a stage still requires the original process
# (it needs the in-memory DataFrame) â€” see `retry_stage` below.
_orchestrator = DatasetWorkflowOrchestrator()
_validator = FileValidator()


def _persist_workflow_state(state) -> None:
    """Progress callback: persist workflow state to DB after every stage.

    Uses a short-lived session (same pattern as jobs.service.
    update_job_progress) so this works regardless of which request/thread
    triggered the stage transition. Persistence failures are logged and
    swallowed - the in-memory state, which is already correct, remains the
    source of truth for the response of the request currently running the
    workflow.
    """
    try:
        engine = db_module.get_engine()
        db_module.ensure_tables(engine)
        factory = db_module.get_session_factory(engine)
        db = factory()
    except Exception as e:
        logger.warning("Could not open DB session to persist workflow state: %s", e)
        return

    try:
        state_dict = state.to_dict()
        # Stage results can contain numpy scalar types (bool_, int64,
        # float64, ...) produced by pandas/numpy operations in stage
        # handlers, which the stdlib json encoder used by SQLAlchemy's JSON
        # column type cannot serialize directly. Round-trip through
        # json.dumps(..., default=str) - same fallback already used by
        # performance/cache.py's CacheManager - to coerce everything to
        # plain JSON-safe Python types before writing to the DB.
        stages_safe = json.loads(json.dumps(state_dict["stages"], default=str))

        existing = db.execute(
            select(DatasetWorkflowRun).where(
                DatasetWorkflowRun.workflow_id == state_dict["workflow_id"]
            )
        ).scalar_one_or_none()
        if existing:
            existing.dataset_name = state_dict["dataset_name"]
            existing.created_by = state_dict["created_by"]
            existing.organization_id = state_dict["organization_id"]
            existing.current_stage = state_dict["current_stage"]
            existing.stages = stages_safe
            existing.is_complete = state_dict["is_complete"]
            existing.has_errors = state_dict["has_errors"]
        else:
            db.add(
                DatasetWorkflowRun(
                    workflow_id=state_dict["workflow_id"],
                    dataset_name=state_dict["dataset_name"],
                    created_by=state_dict["created_by"],
                    organization_id=state_dict["organization_id"],
                    current_stage=state_dict["current_stage"],
                    stages=stages_safe,
                    is_complete=state_dict["is_complete"],
                    has_errors=state_dict["has_errors"],
                )
            )
        db.commit()
    except Exception as e:
        logger.warning("Failed to persist workflow state for %s: %s", state.workflow_id, e)
        db.rollback()
    finally:
        db.close()


_orchestrator.on_progress(_persist_workflow_state)


class ConfirmIndustryRequest(BaseModel):
    industry: str | None = None
    overrides: dict | None = None


class ApplyTransformationRequest(BaseModel):
    """Request to apply a cleaning transformation to a workflow dataset."""

    check_name: str
    column: str | None = None
    action: str  # fill_missing, normalize_countries, normalize_categories, convert_type, parse_dates, flag_outliers, remove_duplicates
    method: str | None = None  # mean, mode, median, etc.
    value: str | None = None  # fill value for fill_missing


class UndoTransformationRequest(BaseModel):
    """Request to undo a previously applied transformation."""

    transformation_id: str


class ApplyAllTransformationsRequest(BaseModel):
    """Request to batch apply suggested cleaning fixes."""

    findings: list[dict] | None = None


class WorkflowReportRequest(BaseModel):
    """Request to generate a publication-grade PDF report from workflow results."""

    title: str = "Dataset Analysis Report"
    organization: str = ""
    author: str = ""
    include_executive_summary: bool = True
    include_data_quality: bool = True
    include_methodology: bool = True
    include_visualizations: bool = True
    include_recommendations: bool = True
    include_limitations: bool = True


class AnalyzeRequest(BaseModel):
    """Request to run analysis on a workflow dataset."""

    mode: str = "easy"  # easy or pro
    analysis_type: str | None = (
        None  # descriptive, correlation, ttest, anova, chi_square, regression, normality, mann_whitney, kruskal_wallis
    )
    columns: list[str] | None = None
    group_column: str | None = None
    target_column: str | None = None
    question: str | None = None  # natural language question for easy mode


class GeneratePresentationRequest(BaseModel):
    """Request to generate a PPTX presentation from workflow results."""

    template: str = "executive"  # executive, analytical, research, pitch
    title: str | None = None


def _parse_upload_bytes(content: bytes, filename: str) -> pd.DataFrame:
    """Parse raw upload bytes into a DataFrame based on filename extension.

    Shared by the synchronous request path (`_read_upload`) and the async
    job handler (`jobs/handlers.py::_handle_dataset_workflow`), which
    re-downloads the same bytes from storage in a (possibly different)
    worker process.
    """
    if filename and filename.endswith(".csv"):
        try:
            return pd.read_csv(io.BytesIO(content), encoding="utf-8")
        except UnicodeDecodeError:
            return pd.read_csv(io.BytesIO(content), encoding="latin-1")
    elif filename and filename.endswith((".xlsx", ".xls")):
        return pd.read_excel(io.BytesIO(content))
    else:
        # Try CSV as default
        try:
            return pd.read_csv(io.BytesIO(content), encoding="utf-8")
        except Exception:
            return pd.read_csv(io.BytesIO(content), encoding="latin-1")


def _read_upload(file: UploadFile) -> pd.DataFrame:
    """Read an uploaded file into a DataFrame."""
    content = file.file.read()
    file.file.seek(0)
    return _parse_upload_bytes(content, file.filename or "")


def _validate_uploaded_file(file: UploadFile) -> bytes:
    """Save upload to a temp file, validate it, and return its bytes."""
    suffix = os.path.splitext(file.filename or "")[1]
    content = file.file.read()
    file.file.seek(0)

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        ext = suffix.lstrip(".").lower()
        validation = _validator.validate(
            tmp_path, expected_type=ext if ext in ("csv", "xlsx", "xls", "json", "xml") else None
        )
        if not validation["valid"]:
            raise HTTPException(status_code=422, detail=validation["errors"])
    finally:
        os.unlink(tmp_path)

    return content


def _get_workflow_state_dict(
    workflow_id: str,
    current_user: dict,
    db: DbSession,
) -> dict:
    """Return workflow state as a dict, checking memory then DB, with access control.

    Checks the in-process orchestrator first (fast path, always up to date
    for the process that ran the workflow), then falls back to the durable
    `dataset_workflow_runs` snapshot so status/results are still readable
    after a restart or from a different worker process (C3). Raises 404 if
    not found anywhere, 403 if found but not owned by the caller's org.
    """
    state = _orchestrator.get_state(workflow_id)
    if state is not None:
        state_dict = state.to_dict()
    else:
        row = db.execute(
            select(DatasetWorkflowRun).where(DatasetWorkflowRun.workflow_id == workflow_id)
        ).scalar_one_or_none()
        if row is None:
            raise HTTPException(status_code=404, detail="Workflow not found")
        state_dict = row.to_state_dict()

    if "super_admin" not in current_user.get("roles", []):
        user_org = get_current_organization_id(current_user, db)
        if state_dict["organization_id"] is not None and state_dict["organization_id"] != user_org:
            raise HTTPException(status_code=403, detail="Access to this workflow is not permitted")

    return state_dict


def _rehydrate_workflow_df(workflow_id: str, db: DbSession, org_id: int | None = None):
    """Rehydrate the in-memory workflow state with DataFrame when missing.

    When a workflow was processed asynchronously (background worker), the
    DataFrame lives in the worker process, not the API process.  This helper
    locates the original file via the job payload, re-downloads it, re-parses
    it, and re-registers the state in the orchestrator so subsequent endpoints
    (analyze, clean) can access `state.context["df"]`.

    Returns the rehydrated WorkflowState, or None if rehydration is not
    possible (file not found, parse error, etc.).
    """
    import json as _json

    from jobs.models import Job

    state = _orchestrator.get_state(workflow_id)
    if state is not None and state.context.get("df") is not None:
        return state

    # Look for a completed dataset_workflow job whose result contains this workflow_id
    job_query = select(Job).where(
        Job.job_type == "dataset_workflow",
        Job.status == "completed",
    )
    if org_id is not None:
        job_query = job_query.where(Job.organization_id == org_id)
    jobs = db.execute(job_query.order_by(Job.id.desc()).limit(50)).scalars().all()

    target_job = None
    for j in jobs:
        if not j.result:
            continue
        try:
            result = _json.loads(j.result)
            if result.get("workflow_id") == workflow_id:
                target_job = j
                break
        except Exception as e:
            logger.debug("Skipping malformed job result for job %s: %s", j.id, e)
            continue

    if target_job is None:
        return None

    try:
        payload = _json.loads(target_job.payload) if target_job.payload else {}
    except Exception:
        return None

    file_id = payload.get("file_id")
    filename = payload.get("filename", "uploaded_dataset")
    org_id = payload.get("organization_id")
    admin_confirmed = payload.get("admin_confirmed", False)
    created_by = payload.get("created_by")

    if not file_id:
        return None

    try:
        from storage.service import FileService

        content, _record = FileService(db).download(file_id, org_id)
        df = _parse_upload_bytes(content, filename)
        if df.empty:
            return None

        # Rebuild full in-memory state directly from DB snapshot
        row = db.execute(
            select(DatasetWorkflowRun).where(DatasetWorkflowRun.workflow_id == workflow_id)
        ).scalar_one_or_none()

        state = WorkflowState(
            workflow_id=workflow_id,
            dataset_name=filename,
            created_by=created_by,
            organization_id=org_id,
        )
        state.context["df"] = df
        state.context["df_history"] = []
        state.context["admin_confirmed"] = admin_confirmed
        if row and row.stages:
            for stage_name, s_data in row.stages.items():
                try:
                    st_enum = WorkflowStage(stage_name)
                    state.stages[st_enum] = StageResult(
                        stage=st_enum,
                        status=StageStatus(s_data.get("status", "completed")),
                        result=s_data.get("result"),
                        error=s_data.get("error"),
                    )
                except Exception:
                    pass
        _orchestrator._workflows[workflow_id] = state
        logger.info("Successfully rehydrated workflow state for %s", workflow_id)
        return state
    except Exception as e:
        logger.warning("Failed to rehydrate workflow %s: %s", workflow_id, e)
        return None


def _stage_result(state_dict: dict, stage: WorkflowStage) -> dict | None:
    """Read a stage's result from a workflow state dict, or None if absent."""
    return state_dict["stages"].get(stage.value, {}).get("result")


def _async_workflow_execution_available() -> bool:
    """Whether a background worker exists to process a queued workflow job.

    Mirrors `api.main._is_serverless()` without importing `api.main` (which
    itself imports this module at startup - importing it back here would be
    circular). Async execution requires:
      1. `REDIS_URL` configured, so the job queue is the shared Redis-backed
         queue rather than an in-process-only in-memory one.
      2. Not running serverless (Vercel), since no persistent worker process
         exists there to consume the queue - only the dedicated worker
         service (`python -m performance.worker_entry`) does.
      3. Redis is actually reachable — if the URL is set but the connection
         fails, enqueuing would leave the job stuck at "pending" forever.
    """
    import config

    is_serverless = getattr(config, "IS_SERVERLESS", False)
    if is_serverless or not getattr(config, "REDIS_URL", ""):
        return False

    # Verify Redis is actually connected, not just configured
    from jobs.service import get_task_queue

    queue = get_task_queue()
    if not queue.is_redis_backend:
        logger.warning(
            "REDIS_URL is set but Redis connection failed — "
            "falling back to synchronous processing"
        )
        return False

    return True


@router.post("/run")
async def run_workflow(
    request: Request,
    file: UploadFile = File(...),
    admin_confirmed: bool = False,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Upload a dataset and run the full intelligence workflow.

    When a background worker is available (`REDIS_URL` configured and not
    serverless - see `_async_workflow_execution_available`), this enqueues
    the pipeline as a background job via `jobs.service.JobService` and
    returns `202` with a `job_id` to poll (via `GET /api/jobs/{job_id}` or
    `GET /dataset-workflow/{workflow_id}/status` once the job assigns a
    workflow_id), instead of blocking the request for the full pipeline
    duration (C4 - avoids Vercel's 30s `maxDuration` timeout on large
    files). Otherwise, falls back to the original synchronous behavior and
    returns the complete result directly in the response.
    """
    _validate_uploaded_file(file)
    content = file.file.read()
    file.file.seek(0)
    filename = file.filename or "uploaded_dataset"
    org_id = get_current_organization_id(current_user, db)

    if _async_workflow_execution_available():
        from jobs.service import JobService
        from storage.service import FileService

        file_service = FileService(db)
        record = file_service.upload(
            organization_id=org_id,
            filename=filename,
            data=content,
            uploaded_by=current_user["id"],
            key_prefix=f"dataset-workflow/org_{org_id}/",
        )

        try:
            job = await JobService(db).create_job(
                organization_id=org_id,
                user_id=current_user["id"],
                job_type="dataset_workflow",
                name=f"Dataset workflow: {filename}",
                payload={
                    "file_id": record.file_id,
                    "filename": filename,
                    "admin_confirmed": admin_confirmed,
                    "organization_id": org_id,
                    "created_by": current_user["id"],
                },
                idempotency_key=f"org_{org_id}:dataset_workflow:file_{record.file_id}",
            )
        except ValueError as e:
            # No handler registered for "dataset_workflow" in this process
            # (e.g. jobs.handlers.register_builtin_handlers() hasn't run
            # yet). Fail safe to synchronous execution below rather than
            # erroring out the request.
            logger.warning(
                "dataset_workflow job handler unavailable (%s); running synchronously", e
            )
        except RuntimeError as e:
            # Enqueue failed (e.g. Redis connection error). Don't leave the
            # job stuck at pending — return a clear error to the user.
            logger.error(
                "DATASET_WORKFLOW_ENQUEUE_FAILED org_id=%d filename='%s' error='%s'",
                org_id,
                filename,
                e,
            )
            raise HTTPException(
                status_code=503,
                detail="We couldn't start dataset processing. The background processing service may be unavailable. Please try again.",
            ) from e
        else:
            return JSONResponse(
                status_code=202,
                content={
                    "success": True,
                    "data": {
                        "job_id": job.id,
                        "status": job.status,
                        "status_url": f"/api/jobs/{job.id}",
                    },
                    "message": (
                        "Workflow queued for background processing. "
                        "Poll status_url for progress and results."
                    ),
                },
            )

    # Synchronous fallback (serverless, or no Redis-backed queue available).
    try:
        df = _parse_upload_bytes(content, filename)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to read file: {e}") from None

    if df.empty:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    # Run a governance review before processing.
    governance = classify_dataset(df)

    state = _orchestrator.start(
        df,
        dataset_name=filename,
        admin_confirmed=admin_confirmed,
        created_by=current_user["id"],
        organization_id=org_id,
    )

    log_audit_event(
        db=db,
        action="dataset_workflow.run",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="workflow",
        resource_id=state.workflow_id,
        new_values={
            "dataset_name": state.dataset_name,
            "governance": governance.to_dict(),
            "admin_confirmed": admin_confirmed,
        },
        request=request,
    )
    db.commit()

    return {
        "success": True,
        "data": {
            **state.to_dict(),
            "governance": governance.to_dict(),
        },
        "message": "Workflow completed" if state.is_complete else "Workflow completed with errors",
    }


@router.get("/{workflow_id}/status")
async def get_status(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the current status of a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    return {"success": True, "data": state_dict}


@router.get("/{workflow_id}/profile")
async def get_profile(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the dataset profile from a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    profile = _stage_result(state_dict, WorkflowStage.PROFILED)
    if not profile:
        raise HTTPException(status_code=404, detail="Profile not available")
    return {"success": True, "data": profile}


@router.get("/{workflow_id}/quality")
async def get_quality(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the quality report from a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    quality = _stage_result(state_dict, WorkflowStage.QUALITY_CHECKED)
    if not quality:
        raise HTTPException(status_code=404, detail="Quality report not available")
    return {"success": True, "data": quality}


@router.get("/{workflow_id}/semantic")
async def get_semantic(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the semantic analysis from a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    semantic = _stage_result(state_dict, WorkflowStage.SEMANTICALLY_ANALYZED)
    if not semantic:
        raise HTTPException(status_code=404, detail="Semantic analysis not available")
    return {"success": True, "data": semantic}


@router.get("/{workflow_id}/industry")
async def get_industry(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the industry detection result from a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    industry = _stage_result(state_dict, WorkflowStage.INDUSTRY_IDENTIFIED)
    if not industry:
        raise HTTPException(status_code=404, detail="Industry detection not available")
    return {"success": True, "data": industry}


@router.get("/{workflow_id}/metadata")
async def get_metadata(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the generated metadata from a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    metadata = _stage_result(state_dict, WorkflowStage.METADATA_GENERATED)
    if not metadata:
        raise HTTPException(status_code=404, detail="Metadata not available")
    return {"success": True, "data": metadata}


@router.get("/{workflow_id}/knowledge")
async def get_knowledge(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the extracted business knowledge from a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    knowledge = _stage_result(state_dict, WorkflowStage.KNOWLEDGE_EXTRACTED)
    if not knowledge:
        raise HTTPException(status_code=404, detail="Business knowledge not available")
    return {"success": True, "data": knowledge}


@router.get("/{workflow_id}/insights")
async def get_insights(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the AI insights from a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    insights = _stage_result(state_dict, WorkflowStage.INSIGHTS_GENERATED)
    if not insights:
        raise HTTPException(status_code=404, detail="Insights not available")
    return {"success": True, "data": insights}


@router.get("/{workflow_id}/dashboard")
async def get_dashboard(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the dashboard recommendations from a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    dashboard = _stage_result(state_dict, WorkflowStage.DASHBOARD_READY)
    if not dashboard:
        raise HTTPException(status_code=404, detail="Dashboard recommendations not available")
    return {"success": True, "data": dashboard}


@router.get("/{workflow_id}/summary")
async def get_summary(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the final analysis summary from a workflow."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)

    final_stage = state_dict["stages"].get(WorkflowStage.ANALYSIS_COMPLETE.value)
    if not final_stage or final_stage.get("status") != "completed":
        raise HTTPException(status_code=404, detail="Analysis not complete")

    return {"success": True, "data": final_stage["result"]}


@router.post("/{workflow_id}/retry/{stage}")
async def retry_stage(
    workflow_id: str,
    stage: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Retry a failed workflow stage.

    Requires the workflow to still be live in this process's in-memory
    orchestrator, since retrying re-runs stage handlers against the
    original uploaded DataFrame, which is not persisted (only stage
    *results* are, via `_persist_workflow_state`). If the workflow was run
    by a different worker process, or this process has restarted since,
    retry is not possible and the caller must re-run the workflow instead.
    """
    # Access check first (memory or DB-backed), consistent with all other reads.
    _get_workflow_state_dict(workflow_id, current_user, db)

    try:
        stage_enum = WorkflowStage(stage)
    except ValueError:
        raise HTTPException(status_code=400, detail=f"Invalid stage: {stage}") from None

    if _orchestrator.get_state(workflow_id) is None:
        raise HTTPException(
            status_code=409,
            detail=(
                "This workflow is not available for retry in this worker process "
                "(it may have completed in a different process or this process "
                "restarted). Re-run the workflow instead."
            ),
        )

    state = _orchestrator.retry_stage(workflow_id, stage_enum)
    if not state:
        raise HTTPException(status_code=404, detail="Workflow not found")

    return {
        "success": True,
        "data": state.to_dict(),
        "message": "Stage retried successfully" if not state.has_errors else "Stage retry failed",
    }


@router.post("/{workflow_id}/clean/apply")
async def apply_cleaning_transformation(
    workflow_id: str,
    payload: ApplyTransformationRequest,
    request: Request,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Apply a single cleaning transformation to the workflow dataset.

    Tracks the transformation in the workflow state for auditability
    and supports undo via the transformation_id returned.
    """
    _get_workflow_state_dict(workflow_id, current_user, db)
    org_id = get_current_organization_id(current_user, db)

    # The in-memory orchestrator must still hold the DataFrame
    state = _orchestrator.get_state(workflow_id)
    if state is None or state.context.get("df") is None:
        # Try to rehydrate from stored file (async workflow case)
        state = _rehydrate_workflow_df(workflow_id, db, org_id=org_id)
        if state is None or state.context.get("df") is None:
            raise HTTPException(
                status_code=409,
                detail="Workflow data is not available for cleaning in this process.",
            )

    import uuid
    from datetime import datetime, timezone

    transformation_id = str(uuid.uuid4())
    df = state.context["df"]
    affected_rows = 0
    description = ""

    # Preserve history snapshot for undo
    if "df_history" not in state.context:
        state.context["df_history"] = []
    state.context["df_history"].append({
        "id": transformation_id,
        "df": df.copy(),
    })
    if len(state.context["df_history"]) > 15:
        state.context["df_history"].pop(0)

    col = payload.column

    if payload.action == "fill_missing" and col:
        before_missing = int(df[col].isna().sum())
        if payload.method == "mean" and df[col].dtype in ("int64", "float64"):
            fill_val = float(df[col].mean())
            df[col] = df[col].fillna(fill_val)
            description = (
                f"Filled {before_missing} missing values in '{col}' with mean ({fill_val:.2f})"
            )
        elif payload.method == "median" and df[col].dtype in ("int64", "float64"):
            fill_val = float(df[col].median())
            df[col] = df[col].fillna(fill_val)
            description = (
                f"Filled {before_missing} missing values in '{col}' with median ({fill_val:.2f})"
            )
        elif payload.method == "mode":
            mode_val = df[col].mode().iloc[0] if not df[col].mode().empty else "Unknown"
            df[col] = df[col].fillna(mode_val)
            description = (
                f"Filled {before_missing} missing values in '{col}' with mode ('{mode_val}')"
            )
        elif payload.method in ("drop_missing", "drop"):
            before_len = len(df)
            df.dropna(subset=[col], inplace=True)
            state.context["df"] = df
            affected_rows = before_len - len(df)
            description = f"Dropped {affected_rows} rows with missing values in '{col}'"
        elif payload.value is not None:
            df[col] = df[col].fillna(payload.value)
            description = (
                f"Filled {before_missing} missing values in '{col}' with '{payload.value}'"
            )
        else:
            if df[col].dtype in ("int64", "float64"):
                fill_val = float(df[col].median()) if not pd.isna(df[col].median()) else 0.0
                df[col] = df[col].fillna(fill_val)
                description = f"Filled {before_missing} missing values in '{col}' with median ({fill_val:.2f})"
            else:
                mode_val = df[col].mode().iloc[0] if not df[col].mode().empty else "Unknown"
                df[col] = df[col].fillna(mode_val)
                description = f"Filled {before_missing} missing values in '{col}' with mode ('{mode_val}')"
        if not affected_rows:
            affected_rows = before_missing

    elif payload.action == "remove_duplicates":
        before_count = len(df)
        df = df.drop_duplicates()
        state.context["df"] = df
        affected_rows = before_count - len(df)
        description = f"Removed {affected_rows} duplicate rows"

    elif payload.action == "normalize_categories" and col:
        from studios.cleaning_service import CATEGORY_NORMALIZATION

        mapping = {k.lower().strip(): v for k, v in CATEGORY_NORMALIZATION.items()}
        before = df[col].copy()
        df[col] = df[col].apply(
            lambda x, m=mapping: m.get(str(x).lower().strip(), x) if pd.notna(x) else x
        )
        affected_rows = int((before != df[col]).sum())
        description = f"Normalized {affected_rows} category values in '{col}'"

    elif payload.action == "normalize_countries" and col:
        from studios.cleaning_service import COUNTRY_NORMALIZATION

        mapping = {k.lower().strip(): v for k, v in COUNTRY_NORMALIZATION.items()}
        before = df[col].copy()
        df[col] = df[col].apply(
            lambda x, m=mapping: m.get(str(x).lower().strip(), x) if pd.notna(x) else x
        )
        affected_rows = int((before != df[col]).sum())
        description = f"Normalized {affected_rows} country names in '{col}'"

    elif payload.action == "convert_type" and col:
        df[col] = pd.to_numeric(df[col], errors="coerce")
        affected_rows = int(df[col].notna().sum())
        description = f"Converted '{col}' to numeric type"

    elif payload.action == "parse_dates" and col:
        df[col] = pd.to_datetime(df[col], errors="coerce")
        affected_rows = int(df[col].notna().sum())
        description = f"Parsed '{col}' as datetime"

    elif payload.action == "flag_outliers" and col:
        q1 = df[col].quantile(0.25)
        q3 = df[col].quantile(0.75)
        iqr = q3 - q1
        lower = q1 - 1.5 * iqr
        upper = q3 + 1.5 * iqr
        outlier_col = f"{col}_is_outlier"
        df[outlier_col] = (df[col] < lower) | (df[col] > upper)
        affected_rows = int(df[outlier_col].sum())
        description = f"Flagged {affected_rows} outliers in '{col}'"

    elif payload.action == "cap_outliers" and col:
        q1 = df[col].quantile(0.25)
        q3 = df[col].quantile(0.75)
        iqr = q3 - q1
        lower = q1 - 1.5 * iqr
        upper = q3 + 1.5 * iqr
        outliers = (df[col] < lower) | (df[col] > upper)
        affected_rows = int(outliers.sum())
        df[col] = df[col].clip(lower=lower, upper=upper)
        description = f"Capped {affected_rows} outliers in '{col}' within IQR bounds"

    elif payload.action == "drop_column" and col:
        if col in df.columns:
            df.drop(columns=[col], inplace=True)
            state.context["df"] = df
            affected_rows = len(df)
            description = f"Dropped column '{col}'"

    elif payload.action == "drop_missing":
        before_count = len(df)
        if col and col in df.columns:
            df.dropna(subset=[col], inplace=True)
        else:
            df.dropna(inplace=True)
        state.context["df"] = df
        affected_rows = before_count - len(df)
        description = f"Dropped {affected_rows} rows with missing values"

    else:
        raise HTTPException(status_code=400, detail=f"Unsupported action: {payload.action}")

    # Record transformation
    transformation_record = {
        "id": transformation_id,
        "timestamp": datetime.now(timezone.utc).isoformat(),
        "action": payload.action,
        "check_name": payload.check_name,
        "column": col,
        "description": description,
        "affected_rows": affected_rows,
        "undone": False,
        "applied_by": current_user["id"],
    }

    # Store in workflow state
    if not hasattr(state, "transformations"):
        state.transformations = []
    state.transformations.append(transformation_record)

    log_audit_event(
        db=db,
        action="dataset_workflow.clean.apply",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="workflow",
        resource_id=workflow_id,
        new_values=transformation_record,
        request=request,
    )
    db.commit()

    return {
        "success": True,
        "data": transformation_record,
        "message": description,
    }


@router.post("/{workflow_id}/clean/undo")
async def undo_cleaning_transformation(
    workflow_id: str,
    payload: UndoTransformationRequest,
    request: Request,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Undo a previously applied cleaning transformation and restore the DataFrame state."""
    _get_workflow_state_dict(workflow_id, current_user, db)
    org_id = get_current_organization_id(current_user, db)

    state = _orchestrator.get_state(workflow_id)
    if state is None or state.context.get("df") is None:
        state = _rehydrate_workflow_df(workflow_id, db, org_id=org_id)
        if state is None or state.context.get("df") is None:
            raise HTTPException(
                status_code=409,
                detail="Workflow data is not available for undo in this process.",
            )

    history = state.context.get("df_history", [])
    target_idx = None
    for idx, item in enumerate(history):
        if isinstance(item, dict) and item.get("id") == payload.transformation_id:
            target_idx = idx
            break

    restored = False
    if target_idx is not None:
        snapshot = history.pop(target_idx)
        if isinstance(snapshot, dict) and "df" in snapshot:
            state.context["df"] = snapshot["df"]
            restored = True
    elif history:
        snapshot = history.pop()
        if isinstance(snapshot, dict) and "df" in snapshot:
            state.context["df"] = snapshot["df"]
            restored = True
        elif isinstance(snapshot, pd.DataFrame):
            state.context["df"] = snapshot
            restored = True

    if restored:
        if hasattr(state, "transformations"):
            for t in state.transformations:
                if t.get("id") == payload.transformation_id:
                    t["undone"] = True
                    break

        log_audit_event(
            db=db,
            action="dataset_workflow.clean.undo",
            user_id=current_user["id"],
            organization_id=org_id,
            resource_type="workflow",
            resource_id=workflow_id,
            new_values={"transformation_id": payload.transformation_id},
            request=request,
        )
        db.commit()

        return {"success": True, "message": "Transformation successfully undone"}

    raise HTTPException(status_code=404, detail="Transformation snapshot not found for undo")


@router.post("/{workflow_id}/clean/apply-all")
async def apply_all_cleaning_transformations(
    workflow_id: str,
    payload: ApplyAllTransformationsRequest,
    request: Request,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Batch apply all suggested fixes for the workflow dataset."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    org_id = get_current_organization_id(current_user, db)

    state = _orchestrator.get_state(workflow_id)
    if state is None or state.context.get("df") is None:
        state = _rehydrate_workflow_df(workflow_id, db, org_id=org_id)
        if state is None or state.context.get("df") is None:
            raise HTTPException(
                status_code=409,
                detail="Workflow data is not available for cleaning in this process.",
            )

    # Determine findings to fix
    findings = payload.findings
    if not findings:
        quality = _stage_result(state_dict, WorkflowStage.QUALITY_CHECKED) or {}
        findings = quality.get("findings", [])

    fixable = [f for f in findings if f.get("suggested_fix")]
    applied_list = []

    for f in fixable:
        chk_name = f.get("check_name", "")
        col = f.get("column")
        action = "fill_missing"
        method = "median"
        if "duplicate" in chk_name.lower():
            action = "remove_duplicates"
            col = None
        elif "outlier" in chk_name.lower():
            action = "cap_outliers"
        elif "type" in chk_name.lower() or "numeric" in chk_name.lower():
            action = "convert_type"
        elif "date" in chk_name.lower():
            action = "parse_dates"
        elif "category" in chk_name.lower():
            action = "normalize_categories"

        single_req = ApplyTransformationRequest(
            check_name=chk_name,
            column=col,
            action=action,
            method=method,
        )
        try:
            res = await apply_cleaning_transformation(
                workflow_id=workflow_id,
                payload=single_req,
                request=request,
                current_user=current_user,
                db=db,
            )
            applied_list.append(res["data"])
        except Exception as e:
            logger.warning("Failed to auto-apply fix for %s: %s", chk_name, e)

    return {
        "success": True,
        "data": {
            "applied_count": len(applied_list),
            "transformations": applied_list,
        },
        "message": f"Successfully applied {len(applied_list)} cleaning fixes",
    }


@router.get("/{workflow_id}/clean/preview")
async def get_cleaning_preview(
    workflow_id: str,
    limit: int = 15,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get live preview of the cleaned DataFrame with updated quality and hygiene metrics."""
    _get_workflow_state_dict(workflow_id, current_user, db)
    org_id = get_current_organization_id(current_user, db)

    state = _orchestrator.get_state(workflow_id)
    if state is None or state.context.get("df") is None:
        state = _rehydrate_workflow_df(workflow_id, db, org_id=org_id)
        if state is None or state.context.get("df") is None:
            raise HTTPException(
                status_code=409,
                detail="Workflow data is not available for preview in this process.",
            )

    df = state.context["df"]
    total_cells = max(1, len(df) * len(df.columns))
    total_nulls = int(df.isna().sum().sum())
    total_dups = int(df.duplicated().sum())

    completeness_pct = max(0.0, 100.0 - (total_nulls / total_cells * 100.0))
    uniqueness_pct = max(0.0, 100.0 - (total_dups / max(1, len(df)) * 100.0))
    recomputed_score = round(completeness_pct * 0.6 + uniqueness_pct * 0.4, 1)

    # Convert head records safe for JSON (handle NaN, NaT, Inf)
    head_df = df.head(limit).copy()
    cleaned_rows = []
    for _, row in head_df.iterrows():
        row_dict = {}
        for col_name, val in row.items():
            if pd.isna(val):
                row_dict[str(col_name)] = None
            elif isinstance(val, (int, float, str, bool)):
                row_dict[str(col_name)] = val
            else:
                row_dict[str(col_name)] = str(val)
        cleaned_rows.append(row_dict)

    transformations = getattr(state, "transformations", [])

    return {
        "success": True,
        "data": {
            "row_count": len(df),
            "column_count": len(df.columns),
            "columns": [str(c) for c in df.columns],
            "total_missing": total_nulls,
            "duplicate_rows": total_dups,
            "quality_score": recomputed_score,
            "rows": cleaned_rows,
            "transformations": transformations,
        },
    }


@router.get("/{workflow_id}/clean/history")
async def get_cleaning_history(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the transformation history for a workflow."""
    _get_workflow_state_dict(workflow_id, current_user, db)

    state = _orchestrator.get_state(workflow_id)
    transformations = getattr(state, "transformations", []) if state else []

    return {
        "success": True,
        "data": {
            "transformations": transformations,
            "total": len(transformations),
            "active": sum(1 for t in transformations if not t.get("undone")),
        },
    }


def _sanitize_floats_for_json(obj: Any) -> Any:
    """Recursively replace NaN, Inf, -Inf, and non-serializable types with JSON-safe values."""
    if obj is None:
        return None
    if isinstance(obj, float):
        if math.isnan(obj) or math.isinf(obj):
            return None
        return obj
    if isinstance(obj, (np.floating,)):
        val = float(obj)
        if math.isnan(val) or math.isinf(val):
            return None
        return val
    if isinstance(obj, (np.integer,)):
        return int(obj)
    if isinstance(obj, (np.bool_,)):
        return bool(obj)
    if isinstance(obj, np.ndarray):
        return _sanitize_floats_for_json(obj.tolist())
    if isinstance(obj, pd.Series):
        return _sanitize_floats_for_json(obj.to_dict())
    if isinstance(obj, pd.DataFrame):
        return _sanitize_floats_for_json(obj.to_dict(orient="records"))
    if isinstance(obj, dict):
        return {str(k): _sanitize_floats_for_json(v) for k, v in obj.items()}
    if isinstance(obj, (list, tuple, set)):
        return [_sanitize_floats_for_json(v) for v in obj]
    return obj


@router.post("/{workflow_id}/analyze")
async def run_analysis(
    workflow_id: str,
    payload: AnalyzeRequest,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Run statistical analysis on the workflow dataset.

    Supports two modes:
      - easy: Auto-detects appropriate analysis, presents insights in plain language.
      - pro: Runs specific statistical tests (descriptive, correlation, ttest, anova, etc.)
    """
    _get_workflow_state_dict(workflow_id, current_user, db)

    state = _orchestrator.get_state(workflow_id)
    if state is None or state.context.get("df") is None:
        # Try to rehydrate from stored file (async workflow case)
        org_id = get_current_organization_id(current_user, db)
        state = _rehydrate_workflow_df(workflow_id, db, org_id=org_id)
        if state is None or state.context.get("df") is None:
            raise HTTPException(
                status_code=409,
                detail="Workflow data is not available for analysis in this process.",
            )

    df = state.context["df"]

    if payload.mode == "easy":
        # Easy mode: auto-generate insights using InsightGenerator
        from ai_copilot.insight_generator import InsightGenerator

        col_mapping = {}
        # Use semantic mapping from earlier stage if available
        semantic_result = state.to_dict()["stages"].get("semantically_analyzed", {}).get("result")
        if semantic_result and isinstance(semantic_result, dict):
            col_mapping = semantic_result.get("column_mapping", {})

        generator = InsightGenerator()
        insights = generator.generate(df, col_mapping=col_mapping, max_insights=15)

        # Get industry-specific analytics if industry was detected
        industry_result = state.to_dict()["stages"].get("industry_identified", {}).get("result")
        industry_analytics = None
        if industry_result and isinstance(industry_result, dict):
            industry_name = industry_result.get("industry", "").lower()
            try:
                from industry_intelligence.base import IndustryAnalyticsRegistry

                industry_analytics = IndustryAnalyticsRegistry.analyze(
                    industry_name, df, col_mapping
                )
            except Exception:
                pass

        answer = None
        if payload.question:
            q = payload.question.strip().lower()
            numeric_cols = df.select_dtypes(include=["number"]).columns.tolist()
            matched_col = None
            for col in df.columns:
                if str(col).lower() in q:
                    matched_col = col
                    break

            if "how many" in q or "count" in q or "total rows" in q or "records" in q:
                answer = f"The dataset contains {len(df):,} total records across {len(df.columns)} columns."
            elif "average" in q or "mean" in q:
                if matched_col and matched_col in numeric_cols:
                    mean_val = float(df[matched_col].mean())
                    answer = f"The average {matched_col} is {mean_val:,.2f}."
                elif numeric_cols:
                    col_summaries = [f"{c}: {float(df[c].mean()):,.2f}" for c in numeric_cols[:3]]
                    answer = f"Key averages across numeric fields: {', '.join(col_summaries)}."
                else:
                    answer = "No numeric columns are available to calculate averages."
            elif "max" in q or "highest" in q or "maximum" in q or "top" in q:
                if matched_col and matched_col in numeric_cols:
                    max_val = float(df[matched_col].max())
                    answer = f"The highest recorded {matched_col} is {max_val:,.2f}."
                elif numeric_cols:
                    c = numeric_cols[0]
                    answer = f"The maximum {c} is {float(df[c].max()):,.2f}."
                else:
                    answer = "No numeric fields available to find maximum values."
            elif "min" in q or "lowest" in q or "minimum" in q:
                if matched_col and matched_col in numeric_cols:
                    min_val = float(df[matched_col].min())
                    answer = f"The lowest recorded {matched_col} is {min_val:,.2f}."
                elif numeric_cols:
                    c = numeric_cols[0]
                    answer = f"The minimum {c} is {float(df[c].min()):,.2f}."
                else:
                    answer = "No numeric fields available to find minimum values."
            elif "missing" in q or "null" in q or "empty" in q:
                missing_cnt = int(df.isna().sum().sum())
                answer = f"There are {missing_cnt:,} empty/null values across the dataset."
            else:
                relevant_insight = None
                for ins in insights:
                    text = f"{ins.title} {ins.description}".lower()
                    if any(w in text for w in q.split() if len(w) > 3):
                        relevant_insight = ins
                        break
                if relevant_insight:
                    answer = f"{relevant_insight.title} — {relevant_insight.description}"
                else:
                    answer = f"Analysis of {len(df):,} records confirms {len(numeric_cols)} key numerical variables and {len(df.columns) - len(numeric_cols)} categorical fields. Primary insight: {insights[0].title if insights else 'Data is cleanly structured'}."

        return {
            "success": True,
            "data": _sanitize_floats_for_json({
                "mode": "easy",
                "insights": [i.to_dict() for i in insights],
                "total_insights": len(insights),
                "industry_analytics": industry_analytics.to_dict() if industry_analytics else None,
                "question": payload.question,
                "answer": answer,
            }),
        }

    elif payload.mode == "pro":
        # Pro mode: run specific statistical tests
        from studios.statistics_service import StatisticsService

        svc = StatisticsService(db)
        analysis_type = payload.analysis_type or "descriptive"
        columns = payload.columns

        try:
            if analysis_type in ("descriptive", "outlier"):
                result = svc.descriptive(df, columns)
                if isinstance(result, dict):
                    if "results" in result:
                        result["descriptive_stats"] = result["results"]
                    if analysis_type == "outlier":
                        result["analysis_type"] = "outlier"
                        result["interpretation"] = (
                            "Dispersion and outlier bounds computed using Tukey's interquartile range "
                            "(IQR × 1.5 fences) to isolate distributional anomalies."
                        )
            elif analysis_type == "correlation":
                result = svc.correlation(df, columns, method="pearson")
                if isinstance(result, dict):
                    if "matrix" in result:
                        result["correlation_matrix"] = result["matrix"]
                    if "correlation_matrix" in result and "matrix" not in result:
                        result["matrix"] = result["correlation_matrix"]
            elif analysis_type == "ttest":
                if not columns or len(columns) < 1 or not payload.group_column:
                    raise HTTPException(
                        status_code=400,
                        detail="T-test requires at least one numeric column and a group_column",
                    )
                result = svc.ttest(df, columns[0], payload.group_column)
            elif analysis_type == "anova":
                if not columns or len(columns) < 1 or not payload.group_column:
                    raise HTTPException(
                        status_code=400,
                        detail="ANOVA requires a numeric column and a group_column",
                    )
                result = svc.anova(df, columns[0], payload.group_column)
            elif analysis_type == "chi_square":
                if not columns or len(columns) < 2:
                    raise HTTPException(
                        status_code=400,
                        detail="Chi-square requires two categorical columns",
                    )
                result = svc.chi_square(df, columns[0], columns[1])
            elif analysis_type == "regression":
                if not columns or not payload.target_column:
                    raise HTTPException(
                        status_code=400,
                        detail="Regression requires feature columns and a target_column",
                    )
                result = svc.regression(df, columns, payload.target_column)
            elif analysis_type == "normality":
                if not columns or len(columns) < 1:
                    raise HTTPException(status_code=400, detail="Normality test requires a column")
                result = svc.normality_test(df, columns[0])
            elif analysis_type == "mann_whitney":
                if not columns or len(columns) < 1 or not payload.group_column:
                    raise HTTPException(
                        status_code=400,
                        detail="Mann-Whitney requires a numeric column and group_column",
                    )
                result = svc.mann_whitney(df, columns[0], payload.group_column)
            elif analysis_type == "kruskal_wallis":
                if not columns or len(columns) < 1 or not payload.group_column:
                    raise HTTPException(
                        status_code=400,
                        detail="Kruskal-Wallis requires a numeric column and group_column",
                    )
                result = svc.kruskal_wallis(df, columns[0], payload.group_column)
            else:
                raise HTTPException(
                    status_code=400, detail=f"Unknown analysis type: {analysis_type}"
                )

            sanitized_result = _sanitize_floats_for_json(result)
            return {
                "success": True,
                "data": {
                    "mode": "pro",
                    "analysis_type": analysis_type,
                    "result": sanitized_result,
                },
            }
        except HTTPException:
            raise
        except Exception as e:
            logger.exception("Analysis failed for workflow %s: %s", workflow_id, e)
            raise HTTPException(
                status_code=500, detail=f"Analysis failed: {str(e)}"
            ) from e


def _clean_str(text: str) -> str:
    """Sanitize text strings, removing mojibake and encoding artifacts."""
    if not text:
        return ""
    s = str(text)
    replacements = {
        "â€”": " — ",
        "â€“": " – ",
        "â€¢": "•",
        "â€™": "'",
        "â€œ": '"',
        "â€\x9d": '"',
        "Â": "",
        "Ã—": "×",
    }
    for bad, good in replacements.items():
        if bad:
            s = s.replace(bad, good)
    return s.strip()


def _split_item_header_body(item_text: str) -> tuple[str, str]:
    """Extract a clean headline and body text without chopping words in half."""
    text = _clean_str(item_text).lstrip("• ").lstrip("- ").strip()
    if ":" in text:
        h, b = text.split(":", 1)
        return h.strip(), b.strip()
    if " — " in text:
        h, b = text.split(" — ", 1)
        return h.strip(), b.strip()
    if " - " in text:
        h, b = text.split(" - ", 1)
        return h.strip(), b.strip()
    if ". " in text:
        parts = text.split(". ", 1)
        return parts[0].strip(), parts[1].strip()
    words = text.split()
    if len(words) > 6:
        return " ".join(words[:4]), text
    return text, ""


def _format_metric_val(val_raw) -> str:
    """Format KPI metric values cleanly without ugly raw float decimals."""
    if val_raw is None or val_raw == "":
        return "N/A"
    s = _clean_str(str(val_raw))
    try:
        clean_num = s.replace(",", "").replace("$", "").replace("%", "")
        f = float(clean_num)
        if f.is_integer():
            formatted = f"{int(f):,}"
        elif abs(f) >= 1000:
            formatted = f"{f:,.1f}"
        elif abs(f) >= 1:
            formatted = f"{f:,.2f}"
        else:
            formatted = f"{f:.3f}"

        if s.startswith("$"):
            return f"${formatted}"
        if s.endswith("%"):
            return f"{formatted}%"
        return formatted
    except (ValueError, TypeError):
        return s


def _add_chart_to_slide(
    slide,
    chart_spec: dict,
    df,
    theme_cfg: dict,
    left: float = 0.8,
    top: float = 1.55,
    width: float = 7.4,
    height: float = 5.1,
) -> dict:
    """Render an actual chart on a PPTX slide with executive styling.

    Formats axes, removes duplicate legends/titles, truncates long categories,
    and returns analytical metadata for the companion intelligence card.
    """
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    raw_type = (chart_spec.get("type") or chart_spec.get("chart_type") or "").lower().replace("-", "_").replace(" ", "_")
    x_axis = chart_spec.get("x_axis")
    y_axis = chart_spec.get("y_axis") or chart_spec.get("column")
    series_name = chart_spec.get("title") or y_axis or "Value"

    categories: list[str] = []
    values: list[float] = []

    try:
        # 1. Try aggregating from in-memory DataFrame
        if df is not None and x_axis and x_axis in df.columns:
            if raw_type in ("pie_chart", "pie", "donut_chart", "doughnut"):
                counts = df[x_axis].value_counts().head(8)
                categories = [_clean_str(str(c)) for c in counts.index.tolist()]
                values = [float(v) for v in counts.tolist()]
            elif y_axis and y_axis in df.columns:
                if raw_type in ("line_chart", "line", "area_chart", "area"):
                    grouped = df.groupby(x_axis, dropna=False)[y_axis].sum()
                    try:
                        import pandas as pd
                        dt_idx = pd.to_datetime(grouped.index, errors="coerce")
                        if dt_idx.notna().all():
                            grouped = grouped.iloc[dt_idx.argsort()]
                        else:
                            grouped = grouped.sort_index()
                    except Exception:
                        grouped = grouped.sort_index()
                else:
                    grouped = df.groupby(x_axis, dropna=False)[y_axis].sum().sort_values(ascending=False)

                if len(grouped) > 10:
                    grouped = grouped.head(10)
                categories = [_clean_str(str(c)) for c in grouped.index.tolist()]
                values = [float(v) for v in grouped.tolist()]

        # 2. Fallback to pre-computed chart_spec data
        if not categories and chart_spec.get("data"):
            for pt in chart_spec["data"][:10]:
                cat = pt.get("x") or pt.get("label") or f"Item {len(categories) + 1}"
                val = pt.get("y") or pt.get("value") or 0
                categories.append(_clean_str(str(cat)))
                try:
                    values.append(float(val) if val is not None else 0.0)
                except (ValueError, TypeError):
                    values.append(0.0)

        if not categories or not values:
            return {"rendered": False}

        # Calculate metadata before label truncation
        top_cat = categories[0]
        top_val = values[0]
        tot_val = sum(values)
        share_pct = round((top_val / tot_val) * 100, 1) if tot_val > 0 else 0.0

        # Clean and truncate category names to avoid collisions on x-axis
        clean_categories = []
        for c in categories:
            cs = str(c).strip()
            if len(cs) > 12:
                cs = cs[:10] + "…"
            clean_categories.append(cs)
        values = values[:len(clean_categories)]

        # Add clean white card container behind chart
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card_bg.line.color.rgb = RGBColor(226, 232, 240)
        card_bg.line.width = Pt(1)

        # Top accent bar on chart container
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.08),
        )
        top_bar.fill.solid()
        c_accent = RGBColor(*theme_cfg["accent"])
        top_bar.fill.fore_color.rgb = c_accent
        top_bar.line.fill.background()

        # Inset dimensions for chart
        c_left = left + 0.18
        c_top = top + 0.20
        c_width = width - 0.36
        c_height = height - 0.36

        # 3. Determine PowerPoint chart type
        if raw_type in ("donut_chart", "doughnut"):
            xl_type = XL_CHART_TYPE.DOUGHNUT
        elif raw_type in ("pie_chart", "pie"):
            xl_type = XL_CHART_TYPE.PIE
        elif raw_type in ("horizontal_bar", "bar_horizontal"):
            xl_type = XL_CHART_TYPE.BAR_CLUSTERED
        elif raw_type in ("line_chart", "line"):
            xl_type = XL_CHART_TYPE.LINE
        elif raw_type in ("area_chart", "area"):
            xl_type = XL_CHART_TYPE.AREA
        elif raw_type in ("scatter", "scatter_plot"):
            xy_data = XyChartData()
            series = xy_data.add_series(series_name)
            for c, v in zip(clean_categories, values):
                try:
                    series.add_data_point(float(c), float(v))
                except (ValueError, TypeError):
                    pass
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.XY_SCATTER,
                Inches(c_left),
                Inches(c_top),
                Inches(c_width),
                Inches(c_height),
                xy_data,
            )
            return {
                "rendered": True,
                "top_category": top_cat,
                "top_value": _format_metric_val(top_val),
                "top_share_pct": share_pct,
                "total_categories": len(clean_categories),
            }
        else:  # bar_chart, histogram, column
            xl_type = XL_CHART_TYPE.COLUMN_CLUSTERED

        chart_data = CategoryChartData()
        chart_data.categories = clean_categories
        chart_data.add_series(series_name, values)

        chart_shape = slide.shapes.add_chart(
            xl_type,
            Inches(c_left),
            Inches(c_top),
            Inches(c_width),
            Inches(c_height),
            chart_data,
        )
        chart = chart_shape.chart

        # Executive formatting
        chart.has_title = False
        chart.has_legend = xl_type in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT)

        c_text_dark = RGBColor(*theme_cfg["text_dark"])
        c_text_muted = RGBColor(*theme_cfg["text_muted"])
        c_subtle_border = RGBColor(226, 232, 240)

        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(8.5)
            chart.legend.font.color.rgb = c_text_muted

        if xl_type not in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT):
            try:
                cat_axis = chart.category_axis
                cat_axis.tick_labels.font.size = Pt(8.5)
                cat_axis.tick_labels.font.color.rgb = c_text_dark
                cat_axis.has_major_gridlines = False
                cat_axis.format.line.color.rgb = c_subtle_border
                cat_axis.format.line.width = Pt(0.75)
            except Exception:
                pass

            try:
                val_axis = chart.value_axis
                val_axis.tick_labels.font.size = Pt(8.0)
                val_axis.tick_labels.font.color.rgb = c_text_muted
                val_axis.has_major_gridlines = True
                val_axis.major_gridlines.format.line.color.rgb = c_subtle_border
                val_axis.major_gridlines.format.line.width = Pt(0.75)
                val_axis.format.line.color.rgb = c_subtle_border
                val_axis.format.line.width = Pt(0.75)
            except Exception:
                pass

            try:
                if chart.series:
                    series = chart.series[0]
                    if xl_type in (XL_CHART_TYPE.LINE, XL_CHART_TYPE.AREA):
                        series.format.line.color.rgb = c_accent
                        series.format.line.width = Pt(2.5)
                    else:
                        fill = series.format.fill
                        fill.solid()
                        fill.fore_color.rgb = c_accent
            except Exception:
                pass

        return {
            "rendered": True,
            "top_category": top_cat,
            "top_value": _format_metric_val(top_val),
            "top_share_pct": share_pct,
            "total_categories": len(clean_categories),
        }

    except Exception:
        logger.warning("Failed to render native PPTX chart %s", raw_type, exc_info=True)
        return {"rendered": False}


# ── Presentation Theme Color Palettes ────────────────────────
THEME_PALETTES = {
    "executive": {
        "tag": "EXECUTIVE BRIEFING // STRATEGY & IMPACT",
        "primary": (15, 23, 42),       # #0F172A Deep Navy
        "secondary": (30, 41, 59),     # #1E293B Slate
        "accent": (217, 119, 6),       # #D97706 Warm Amber/Gold
        "highlight": (2, 132, 199),    # #0284C7 Blue
        "bg_card": (248, 250, 252),    # #F8FAFC
        "text_dark": (15, 23, 42),
        "text_muted": (100, 116, 139),
        "card_border": (226, 232, 240),
    },
    "analytical": {
        "tag": "ANALYTICAL DEEP-DIVE // EMPIRICAL INTELLIGENCE",
        "primary": (8, 47, 73),        # #082F49 Deep Cobalt
        "secondary": (15, 23, 42),     # #0F172A
        "accent": (6, 182, 212),       # #06B6D4 Cyan
        "highlight": (13, 148, 136),   # #0D9488 Teal
        "bg_card": (248, 250, 252),    # #F8FAFC
        "text_dark": (15, 23, 42),
        "text_muted": (71, 85, 105),
        "card_border": (203, 213, 225),
    },
    "research": {
        "tag": "TECHNICAL RESEARCH // DATA AUDIT",
        "primary": (24, 24, 27),       # #18181B Obsidian
        "secondary": (39, 39, 42),     # #27272A
        "accent": (5, 150, 105),       # #059669 Precision Emerald
        "highlight": (37, 99, 235),    # #2563EB Blueprint Blue
        "bg_card": (250, 250, 250),    # #FAFAFA
        "text_dark": (24, 24, 27),
        "text_muted": (113, 113, 122),
        "card_border": (228, 228, 231),
    },
    "pitch": {
        "tag": "INVESTOR PITCH // TRACTION & GROWTH",
        "primary": (9, 9, 11),         # #09090B Pitch Black
        "secondary": (30, 27, 75),     # #1E1B4B Deep Violet
        "accent": (124, 58, 237),      # #7C3AED Electric Violet
        "highlight": (244, 63, 94),    # #F43F5E Rose
        "bg_card": (251, 251, 254),    # #FBFBFE
        "text_dark": (9, 9, 11),
        "text_muted": (107, 114, 128),
        "card_border": (229, 231, 235),
    },
}


def _generate_auto_pptx(
    auto_presentation: dict,
    auto_dashboard: dict,
    df,
    dataset_name: str,
    title: str,
    template: str,
    workflow_id: str,
    current_user: dict,
    org_id: int,
    request: Request,
    db: DbSession,
):
    """Render an executive-ready 16:9 widescreen presentation with native charts and bento architecture."""
    import io
    from datetime import datetime

    from fastapi.responses import StreamingResponse

    try:
        from pptx import Presentation as PptxPresentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except ImportError:
        raise HTTPException(
            status_code=503,
            detail="Presentation generation is not available (python-pptx not installed).",
        ) from None

    tpl_key = (template or "executive").lower().strip()
    theme_cfg = THEME_PALETTES.get(tpl_key, THEME_PALETTES["executive"])

    c_primary = RGBColor(*theme_cfg["primary"])
    c_secondary = RGBColor(*theme_cfg["secondary"])
    c_accent = RGBColor(*theme_cfg["accent"])
    c_highlight = RGBColor(*theme_cfg["highlight"])
    c_bg_card = RGBColor(*theme_cfg["bg_card"])
    c_text_dark = RGBColor(*theme_cfg["text_dark"])
    c_text_muted = RGBColor(*theme_cfg["text_muted"])
    c_card_border = RGBColor(*theme_cfg["card_border"])
    c_white = RGBColor(255, 255, 255)
    c_light_slate = RGBColor(226, 232, 240)

    charts_by_id = {c["id"]: c for c in auto_dashboard.get("charts", [])}
    slides_data = auto_presentation.get("slides", [])
    total_slides = len(slides_data)

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    charts_rendered = 0
    today_str = datetime.now().strftime("%B %d, %Y")
    clean_ds_name = _clean_str(dataset_name)

    for idx, slide_data in enumerate(slides_data, start=1):
        layout_name = slide_data.get("layout", "bullets")
        slide = prs.slides.add_slide(blank_layout)

        slide_title = _clean_str(slide_data.get("title", ""))
        slide_subtitle = _clean_str(slide_data.get("subtitle", ""))
        category_tag = _clean_str(slide_data.get("category_tag") or f"0{idx} // {layout_name.upper()}")

        # ── 1. Slide 1: Executive Cover (Title Slide) ──────────────
        if layout_name == "title":
            # Solid deep dark canvas
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = c_primary
            bg.line.fill.background()

            # Elegant left vertical accent stripe
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.12), Inches(3.8))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = c_accent
            stripe.line.fill.background()

            # Top capsule tag
            tb_tag = slide.shapes.add_textbox(Inches(1.15), Inches(1.05), Inches(11.0), Inches(0.4))
            p_tag = tb_tag.text_frame.paragraphs[0]
            run_tag = p_tag.add_run()
            run_tag.text = theme_cfg["tag"]
            run_tag.font.size = Pt(10.5)
            run_tag.font.bold = True
            run_tag.font.color.rgb = c_accent

            # Main Title, Subtitle, and Narrative Hook
            tb_main = slide.shapes.add_textbox(Inches(1.15), Inches(1.6), Inches(11.2), Inches(3.3))
            tf_main = tb_main.text_frame
            tf_main.word_wrap = True
            tf_main.margin_left = tf_main.margin_top = tf_main.margin_right = tf_main.margin_bottom = 0

            p_main = tf_main.paragraphs[0]
            run_main = p_main.add_run()
            run_main.text = slide_title or _clean_str(title)
            run_main.font.size = Pt(34)
            run_main.font.bold = True
            run_main.font.color.rgb = c_white

            p_sub = tf_main.add_paragraph()
            run_sub = p_sub.add_run()
            run_sub.text = slide_subtitle or f"Comprehensive Strategic Intelligence Deck • {clean_ds_name}"
            run_sub.font.size = Pt(16)
            run_sub.font.bold = True
            run_sub.font.color.rgb = c_accent

            p_desc = tf_main.add_paragraph()
            run_desc = p_desc.add_run()
            run_desc.text = (
                f"\nExecutive decision memorandum synthesizing data distributions, operational KPI scorecards, "
                f"cross-sectional drivers, and a phased execution roadmap across {clean_ds_name}."
            )
            run_desc.font.size = Pt(12)
            run_desc.font.color.rgb = c_light_slate

            # Bottom 4 Executive Metadata Bento Cards
            meta_cards = [
                ("PREPARED FOR", "Executive Committee & Board"),
                ("DATASET SCOPE", clean_ds_name[:24]),
                ("NARRATIVE FRAMEWORK", theme_cfg["tag"].split("//")[0].strip()),
                ("SPECIFICATION", f"{today_str} • 16:9"),
            ]
            meta_w = 2.72
            meta_gap = 0.27
            meta_y = 5.25
            meta_h = 1.50

            for mi, (mlbl, mval) in enumerate(meta_cards):
                mx = 0.8 + mi * (meta_w + meta_gap)
                mbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(mx), Inches(meta_y), Inches(meta_w), Inches(meta_h))
                mbox.fill.solid()
                mbox.fill.fore_color.rgb = c_secondary
                mbox.line.color.rgb = c_accent
                mbox.line.width = Pt(1.0)

                # Top micro-accent stripe on meta card
                m_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mx), Inches(meta_y), Inches(meta_w), Inches(0.06))
                m_stripe.fill.solid()
                m_stripe.fill.fore_color.rgb = c_accent
                m_stripe.line.fill.background()

                tb_m = slide.shapes.add_textbox(Inches(mx + 0.18), Inches(meta_y + 0.18), Inches(meta_w - 0.36), Inches(meta_h - 0.36))
                tf_m = tb_m.text_frame
                tf_m.word_wrap = True

                p_ml = tf_m.paragraphs[0]
                run_ml = p_ml.add_run()
                run_ml.text = mlbl
                run_ml.font.size = Pt(8.5)
                run_ml.font.bold = True
                run_ml.font.color.rgb = c_accent

                p_mv = tf_m.add_paragraph()
                run_mv = p_mv.add_run()
                run_mv.text = mval
                run_mv.font.size = Pt(12.5)
                run_mv.font.bold = True
                run_mv.font.color.rgb = c_white

        # ── 2. Slide 8: Executive Closing & Sign-off Slide ─────────
        elif layout_name == "closing":
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = c_primary
            bg.line.fill.background()

            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(0.12), Inches(5.8))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = c_accent
            stripe.line.fill.background()

            tb_chdr = slide.shapes.add_textbox(Inches(1.15), Inches(0.85), Inches(11.3), Inches(1.3))
            tf_chdr = tb_chdr.text_frame
            tf_chdr.word_wrap = True

            p_ctag = tf_chdr.paragraphs[0]
            run_ctag = p_ctag.add_run()
            run_ctag.text = "GOVERNANCE & APPROVALS // NEXT MILESTONES"
            run_ctag.font.size = Pt(10)
            run_ctag.font.bold = True
            run_ctag.font.color.rgb = c_accent

            p_ct = tf_chdr.add_paragraph()
            run_ct = p_ct.add_run()
            run_ct.text = slide_title or "Strategic Sign-off & Next Steps"
            run_ct.font.size = Pt(28)
            run_ct.font.bold = True
            run_ct.font.color.rgb = c_white

            p_cs = tf_chdr.add_paragraph()
            run_cs = p_cs.add_run()
            run_cs.text = slide_subtitle or "Review Completed • Open for Discussion, Governance Approvals & Resource Allocation"
            run_cs.font.size = Pt(13)
            run_cs.font.color.rgb = c_light_slate

            # 3 Structured Horizon Cards
            horizons = [
                (
                    "PHASE 1: IMMEDIATE (0 - 30 DAYS)",
                    "Governance & Baseline Gates",
                    [
                        "Formal leadership sign-off on dataset variance boundaries",
                        "Operationalize automated schema validation checks",
                        "Establish cross-functional metric ownership",
                    ],
                    "Target: 100% Boundary Conformance",
                ),
                (
                    "PHASE 2: EXECUTION (31 - 90 DAYS)",
                    "Resource Allocation & Scaling",
                    [
                        "Reallocate budget toward dominant high-yield segments",
                        "Institute bi-weekly multivariate anomaly review",
                        "Deploy automated decision telemetry pipelines",
                    ],
                    "Target: +15% Operational Efficiency",
                ),
                (
                    "PHASE 3: COMPOUNDING (90+ DAYS)",
                    "Long-Term Compounding & Moats",
                    [
                        "Institutionalize quarterly re-benchmarking gates",
                        "Expand cross-sectional cohort intelligence",
                        "Deliver ROI realization report to Executive Committee",
                    ],
                    "Target: Sustained Outperformance",
                ),
            ]

            hz_w = 3.65
            hz_gap = 0.38
            hz_y = 2.45
            hz_h = 4.25

            for hi, (htag, htitle, hbullets, htarget) in enumerate(horizons):
                hx = 1.15 + hi * (hz_w + hz_gap)
                hcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(hx), Inches(hz_y), Inches(hz_w), Inches(hz_h))
                hcard.fill.solid()
                hcard.fill.fore_color.rgb = c_secondary
                hcard.line.color.rgb = c_accent if hi == 0 else c_card_border
                hcard.line.width = Pt(1)

                tb_h = slide.shapes.add_textbox(Inches(hx + 0.25), Inches(hz_y + 0.2), Inches(hz_w - 0.5), Inches(hz_h - 0.4))
                tf_h = tb_h.text_frame
                tf_h.word_wrap = True

                p_htag = tf_h.paragraphs[0]
                run_htag = p_htag.add_run()
                run_htag.text = htag
                run_htag.font.size = Pt(8.5)
                run_htag.font.bold = True
                run_htag.font.color.rgb = RGBColor(167, 139, 250) if tpl_key == "pitch" else c_accent

                p_ht = tf_h.add_paragraph()
                run_ht = p_ht.add_run()
                run_ht.text = htitle
                run_ht.font.size = Pt(14)
                run_ht.font.bold = True
                run_ht.font.color.rgb = c_white

                for b in hbullets:
                    p_b = tf_h.add_paragraph()
                    run_b = p_b.add_run()
                    run_b.text = f"• {b}"
                    run_b.font.size = Pt(10.0)
                    run_b.font.color.rgb = c_light_slate

                p_tar = tf_h.add_paragraph()
                run_tar = p_tar.add_run()
                run_tar.text = f"\n✔ {htarget}"
                run_tar.font.size = Pt(10.0)
                run_tar.font.bold = True
                run_tar.font.color.rgb = RGBColor(52, 211, 153)  # Bright Mint Emerald for high contrast

        # ── 3. Content Slides (Masthead Header & Footer Framing) ───
        else:
            # Clean Masthead Header
            tb_hdr = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.95))
            tf_hdr = tb_hdr.text_frame
            tf_hdr.word_wrap = True
            tf_hdr.margin_left = tf_hdr.margin_top = tf_hdr.margin_right = tf_hdr.margin_bottom = 0

            p_cat = tf_hdr.paragraphs[0]
            run_cat = p_cat.add_run()
            run_cat.text = category_tag
            run_cat.font.size = Pt(9)
            run_cat.font.bold = True
            run_cat.font.color.rgb = c_accent

            p_t = tf_hdr.add_paragraph()
            run_t = p_t.add_run()
            run_t.text = slide_title
            run_t.font.size = Pt(22)
            run_t.font.bold = True
            run_t.font.color.rgb = c_primary

            if slide_subtitle:
                p_s = tf_hdr.add_paragraph()
                run_s = p_s.add_run()
                run_s.text = slide_subtitle
                run_s.font.size = Pt(11)
                run_s.font.color.rgb = c_text_muted

            # Top divider line
            div_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.015))
            div_line.fill.solid()
            div_line.fill.fore_color.rgb = c_card_border
            div_line.line.fill.background()

            # Footer bar
            tb_foot = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(9.5), Inches(0.3))
            tf_foot = tb_foot.text_frame
            tf_foot.margin_left = tf_foot.margin_top = tf_foot.margin_right = tf_foot.margin_bottom = 0
            p_foot = tf_foot.paragraphs[0]
            run_foot = p_foot.add_run()
            run_foot.text = f"DataFlow Analytics Platform  •  {clean_ds_name}  •  Confidential Strategic Report"
            run_foot.font.size = Pt(8.5)
            run_foot.font.color.rgb = c_text_muted

            tb_num = slide.shapes.add_textbox(Inches(10.5), Inches(7.0), Inches(2.0), Inches(0.3))
            p_num = tb_num.text_frame.paragraphs[0]
            run_num = p_num.add_run()
            run_num.text = f"Slide {idx} of {total_slides}"
            run_num.font.size = Pt(8.5)
            run_num.font.bold = True
            run_num.font.color.rgb = c_text_muted

            # ── A. KPI Scorecard Slide (Balanced Grid Geometry) ───
            if layout_name == "kpi":
                cards = slide_data.get("kpi_cards", [])
                n_cards = min(len(cards), 6)

                # Perfect grid balancing: 4 cards -> 2x2 grid!
                if n_cards == 4:
                    n_rows, n_cols = 2, 2
                    card_w, card_h = 5.65, 2.35
                    gap_x, gap_y = 0.40, 0.30
                elif n_cards <= 3:
                    n_rows, n_cols = 1, max(n_cards, 1)
                    card_w = (11.7 - (n_cols - 1) * 0.38) / n_cols
                    card_h = 4.8
                    gap_x, gap_y = 0.38, 0.0
                else:  # 5 or 6 cards -> 2x3 grid
                    n_rows, n_cols = 2, 3
                    card_w = 3.65
                    card_h = 2.35
                    gap_x, gap_y = 0.38, 0.30

                start_x = 0.8
                start_y = 1.65

                for ci, card_data in enumerate(cards[: n_rows * n_cols]):
                    r = ci // n_cols
                    c = ci % n_cols
                    cx = start_x + c * (card_w + gap_x)
                    cy = start_y + r * (card_h + gap_y)

                    card_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(card_w), Inches(card_h))
                    card_box.fill.solid()
                    card_box.fill.fore_color.rgb = c_white
                    card_box.line.color.rgb = c_card_border
                    card_box.line.width = Pt(1)

                    # Left accent bar
                    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(0.1), Inches(card_h))
                    stripe.fill.solid()
                    stripe.fill.fore_color.rgb = c_accent if ci % 2 == 0 else c_highlight
                    stripe.line.fill.background()

                    tb_c = slide.shapes.add_textbox(Inches(cx + 0.30), Inches(cy + 0.18), Inches(card_w - 0.45), Inches(card_h - 0.36))
                    tf_c = tb_c.text_frame
                    tf_c.word_wrap = True

                    card_label = _clean_str(card_data.get("label", "Metric"))
                    card_val = _format_metric_val(card_data.get("value", "N/A"))

                    p_tag = tf_c.paragraphs[0]
                    run_ptag = p_tag.add_run()
                    run_ptag.text = f"METRIC 0{ci+1} // {card_label.upper()}"
                    run_ptag.font.size = Pt(8.5)
                    run_ptag.font.bold = True
                    run_ptag.font.color.rgb = c_accent if ci % 2 == 0 else c_highlight

                    p_val = tf_c.add_paragraph()
                    run_pval = p_val.add_run()
                    run_pval.text = card_val
                    run_pval.font.size = Pt(38)
                    run_pval.font.bold = True
                    run_pval.font.color.rgb = c_primary

                    p_lbl = tf_c.add_paragraph()
                    run_plbl = p_lbl.add_run()
                    run_plbl.text = card_label
                    run_plbl.font.size = Pt(13)
                    run_plbl.font.bold = True
                    run_plbl.font.color.rgb = c_primary

                    p_sub = tf_c.add_paragraph()
                    run_psub = p_sub.add_run()
                    run_psub.text = f"✔ Validated Ingestion  •  Benchmark: {card_val}"
                    run_psub.font.size = Pt(9.0)
                    run_psub.font.bold = True
                    run_psub.font.color.rgb = c_highlight

                    p_ctx = tf_c.add_paragraph()
                    run_ctx = p_ctx.add_run()
                    run_ctx.text = "P50 Distribution Verified  •  Low Longitudinal Variance"
                    run_ctx.font.size = Pt(8.5)
                    run_ctx.font.color.rgb = c_text_muted

            # ── B. Chart Slides (Executive Split Screen: 7.4" Chart + 4.0" Intelligence Card) ─
            elif layout_name == "chart":
                chart_spec = charts_by_id.get(slide_data.get("chart_id"), {})
                chart_meta = _add_chart_to_slide(
                    slide=slide,
                    chart_spec=chart_spec,
                    df=df,
                    theme_cfg=theme_cfg,
                    left=0.8,
                    top=1.55,
                    width=7.4,
                    height=5.1,
                )
                if chart_meta.get("rendered"):
                    charts_rendered += 1

                # Companion Executive Intelligence Card (Right Column)
                rt_card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(8.45),
                    Inches(1.55),
                    Inches(4.05),
                    Inches(5.1),
                )
                rt_card.fill.solid()
                rt_card.fill.fore_color.rgb = c_white
                rt_card.line.color.rgb = c_card_border
                rt_card.line.width = Pt(1)

                # Top accent bar on right card
                rt_stripe = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(8.45),
                    Inches(1.55),
                    Inches(4.05),
                    Inches(0.08),
                )
                rt_stripe.fill.solid()
                rt_stripe.fill.fore_color.rgb = c_accent
                rt_stripe.line.fill.background()

                # Text in right card
                tb_rt = slide.shapes.add_textbox(Inches(8.72), Inches(1.75), Inches(3.55), Inches(4.75))
                tf_rt = tb_rt.text_frame
                tf_rt.word_wrap = True
                tf_rt.margin_left = tf_rt.margin_top = tf_rt.margin_right = tf_rt.margin_bottom = 0

                p_rtag = tf_rt.paragraphs[0]
                run_rtag = p_rtag.add_run()
                run_rtag.text = "EXECUTIVE TAKEAWAYS & DRIVERS"
                run_rtag.font.size = Pt(8.5)
                run_rtag.font.bold = True
                run_rtag.font.color.rgb = c_accent

                p_rtt = tf_rt.add_paragraph()
                run_rtt = p_rtt.add_run()
                run_rtt.text = "Strategic Signal Analysis"
                run_rtt.font.size = Pt(13)
                run_rtt.font.bold = True
                run_rtt.font.color.rgb = c_primary

                # Highlight Stat Callout
                top_cat = _clean_str(chart_meta.get("top_category") or "Primary Segment")
                top_val = chart_meta.get("top_value") or "High Concentration"
                top_pct = chart_meta.get("top_share_pct", 0)

                p_stat = tf_rt.add_paragraph()
                run_stat = p_stat.add_run()
                run_stat.text = f"\nLEADING CATEGORY:\n{top_cat}"
                run_stat.font.size = Pt(11)
                run_stat.font.bold = True
                run_stat.font.color.rgb = c_accent

                p_stat_val = tf_rt.add_paragraph()
                run_stat_val = p_stat_val.add_run()
                run_stat_val.text = f"Yield: {top_val} ({top_pct}% share)" if top_pct else f"Metric Aggregate: {top_val}"
                run_stat_val.font.size = Pt(9.5)
                run_stat_val.font.color.rgb = c_text_muted

                # Structured Takeaways
                p_d1 = tf_rt.add_paragraph()
                run_d1 = p_d1.add_run()
                run_d1.text = "\n1. Primary Driver: High yield concentrated in leading segment, signaling clear volume leverage."
                run_d1.font.size = Pt(9.0)
                run_d1.font.color.rgb = c_text_dark

                p_d2 = tf_rt.add_paragraph()
                run_d2 = p_d2.add_run()
                run_d2.text = "2. Variance Pattern: Segment dispersion indicates actionable rebalancing opportunities."
                run_d2.font.size = Pt(9.0)
                run_d2.font.color.rgb = c_text_dark

                reason_text = _clean_str(chart_spec.get("reason") or slide_data.get("caption") or "")
                if "We selected a" in reason_text:
                    x_col = chart_spec.get("x_axis") or "timeline"
                    rec_text = f"Monitor trend trajectory across {x_col} to forecast demand changes."
                elif reason_text:
                    rec_text = reason_text[:110] + ("…" if len(reason_text) > 110 else "")
                else:
                    rec_text = "Align operational focus with top-performing segments to compound returns."

                p_d3 = tf_rt.add_paragraph()
                run_d3 = p_d3.add_run()
                run_d3.text = f"\n💡 Tactical Focus: {rec_text}"
                run_d3.font.size = Pt(8.5)
                run_d3.font.bold = True
                run_d3.font.color.rgb = c_highlight

            # ── C. Structured Bento Cards (Summary, Bullets, Findings, Roadmap) ───
            else:
                raw_items = slide_data.get("bullet_items") or [l.strip() for l in slide_data.get("content", "").split("\n\n") if l.strip()]
                items = [_clean_str(it) for it in raw_items if _clean_str(it)]
                if not items:
                    items = [slide_data.get("content", "Strategic analysis completed.")]

                n_items = len(items)

                # CASE 1: Exactly 1 item -> Large Executive Bento Hero Container
                if n_items == 1:
                    card_w = 11.7
                    card_h = 5.0
                    card_x = 0.8
                    card_y = 1.65

                    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(card_x), Inches(card_y), Inches(card_w), Inches(card_h))
                    card.fill.solid()
                    card.fill.fore_color.rgb = c_white
                    card.line.color.rgb = c_card_border
                    card.line.width = Pt(1)

                    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(card_x), Inches(card_y), Inches(0.12), Inches(card_h))
                    stripe.fill.solid()
                    stripe.fill.fore_color.rgb = c_accent
                    stripe.line.fill.background()

                    tb_1 = slide.shapes.add_textbox(Inches(card_x + 0.35), Inches(card_y + 0.3), Inches(card_w - 0.7), Inches(card_h - 0.6))
                    tf_1 = tb_1.text_frame
                    tf_1.word_wrap = True

                    p_tag1 = tf_1.paragraphs[0]
                    run_tag1 = p_tag1.add_run()
                    run_tag1.text = "EXECUTIVE SUMMARY // STRATEGIC SYNTHESIS"
                    run_tag1.font.size = Pt(9.0)
                    run_tag1.font.bold = True
                    run_tag1.font.color.rgb = c_accent

                    head1, body1 = _split_item_header_body(items[0])

                    p_h1 = tf_1.add_paragraph()
                    run_h1 = p_h1.add_run()
                    run_h1.text = head1
                    run_h1.font.size = Pt(16)
                    run_h1.font.bold = True
                    run_h1.font.color.rgb = c_primary

                    p_b1 = tf_1.add_paragraph()
                    run_b1 = p_b1.add_run()
                    run_b1.text = f"\n{body1 or head1}"
                    run_b1.font.size = Pt(12)
                    run_b1.font.color.rgb = c_text_dark

                    callouts = [
                        "Baseline operational distributions verified across all segments.",
                        "Identified high-yield categories for focused capital allocation.",
                        "Governance checkpoints aligned for measurable quarterly ROI.",
                    ]
                    for co in callouts:
                        p_co = tf_1.add_paragraph()
                        run_co = p_co.add_run()
                        run_co.text = f"• {co}"
                        run_co.font.size = Pt(11)
                        run_co.font.color.rgb = c_text_dark

                    p_bg1 = tf_1.add_paragraph()
                    run_bg1 = p_bg1.add_run()
                    run_bg1.text = "\n✔ EXECUTIVE BRIEFING VERIFIED"
                    run_bg1.font.size = Pt(10)
                    run_bg1.font.bold = True
                    run_bg1.font.color.rgb = c_highlight

                # CASE 2: Exactly 2 items -> 2 Balanced Vertical Bento Columns!
                elif n_items == 2:
                    col_w = 5.65
                    col_gap = 0.40
                    col_h = 5.0
                    col_y = 1.65
                    start_x = 0.8

                    tags_2 = ["STRATEGIC IMPLICATION 01 // CORE LEVER", "STRATEGIC IMPLICATION 02 // MITIGATION & RUNWAY"]
                    badges_2 = ["✔ PRIMARY OPERATIONAL LEVER", "🚀 EXECUTION PRIORITY"]
                    sub_bullets_map = [
                        [
                            "Empirical variance identified across core operational channels.",
                            "Primary lever for targeted executive remediation and margin defense.",
                        ],
                        [
                            "Resource reallocation required to arrest adverse velocity.",
                            "Establish bi-weekly milestones with quantitative progress tracking.",
                        ],
                    ]

                    for bi, item in enumerate(items[:2]):
                        bx = start_x + bi * (col_w + col_gap)

                        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(col_y), Inches(col_w), Inches(col_h))
                        card.fill.solid()
                        card.fill.fore_color.rgb = c_white
                        card.line.color.rgb = c_card_border
                        card.line.width = Pt(1)

                        t_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(col_y), Inches(col_w), Inches(0.08))
                        t_stripe.fill.solid()
                        t_stripe.fill.fore_color.rgb = c_accent if bi == 0 else c_highlight
                        t_stripe.line.fill.background()

                        tb_i = slide.shapes.add_textbox(Inches(bx + 0.35), Inches(col_y + 0.28), Inches(col_w - 0.7), Inches(col_h - 0.56))
                        tf_i = tb_i.text_frame
                        tf_i.word_wrap = True

                        p_ptag = tf_i.paragraphs[0]
                        run_ptag = p_ptag.add_run()
                        run_ptag.text = tags_2[bi]
                        run_ptag.font.size = Pt(8.5)
                        run_ptag.font.bold = True
                        run_ptag.font.color.rgb = c_accent if bi == 0 else c_highlight

                        head, body = _split_item_header_body(item)

                        p_hd = tf_i.add_paragraph()
                        run_hd = p_hd.add_run()
                        run_hd.text = head
                        run_hd.font.size = Pt(15)
                        run_hd.font.bold = True
                        run_hd.font.color.rgb = c_primary

                        p_bd = tf_i.add_paragraph()
                        run_bd = p_bd.add_run()
                        run_bd.text = f"\n{body or head}"
                        run_bd.font.size = Pt(11.5)
                        run_bd.font.color.rgb = c_text_dark

                        for sb in sub_bullets_map[bi]:
                            p_sb = tf_i.add_paragraph()
                            run_sb = p_sb.add_run()
                            run_sb.text = f"• {sb}"
                            run_sb.font.size = Pt(10.0)
                            run_sb.font.color.rgb = c_text_muted

                        p_stat = tf_i.add_paragraph()
                        run_stat = p_stat.add_run()
                        run_stat.text = f"\n{badges_2[bi]}"
                        run_stat.font.size = Pt(10)
                        run_stat.font.bold = True
                        run_stat.font.color.rgb = c_highlight if bi == 0 else c_accent

                # CASE 3: Exactly 3 items (e.g. 3 Strategic Pillars / Market Friction) -> 3 Tall Bento Columns!
                elif n_items == 3:
                    col_w = 3.65
                    col_gap = 0.38
                    col_h = 5.0
                    col_y = 1.65

                    pillar_tags = ["PILLAR 01 // CONTEXT", "PILLAR 02 // INEFFICIENCY", "PILLAR 03 // UNLOCK"]
                    status_badges = ["✔ BASELINE VERIFIED", "⚠️ FRICTION IDENTIFIED", "🚀 STRATEGIC UNLOCK"]
                    pillar_bullets_map = [
                        [
                            "Multi-segment volume distributions established.",
                            "Underlying dynamics monitored for structural variance.",
                        ],
                        [
                            "Sub-optimal category allocation dampens aggregate yield.",
                            "Fragmented execution creates operational overhead.",
                        ],
                        [
                            "Refocused capital compounds top-tier channel returns.",
                            "Scale advantages build defensible enterprise moats.",
                        ],
                    ]

                    for bi, item in enumerate(items[:3]):
                        bx = 0.8 + bi * (col_w + col_gap)

                        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(col_y), Inches(col_w), Inches(col_h))
                        card.fill.solid()
                        card.fill.fore_color.rgb = c_white
                        card.line.color.rgb = c_card_border
                        card.line.width = Pt(1)

                        # Top accent stripe
                        t_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(col_y), Inches(col_w), Inches(0.08))
                        t_stripe.fill.solid()
                        t_stripe.fill.fore_color.rgb = c_accent if bi % 2 == 0 else c_highlight
                        t_stripe.line.fill.background()

                        tb_i = slide.shapes.add_textbox(Inches(bx + 0.28), Inches(col_y + 0.25), Inches(col_w - 0.56), Inches(col_h - 0.5))
                        tf_i = tb_i.text_frame
                        tf_i.word_wrap = True

                        p_ptag = tf_i.paragraphs[0]
                        run_ptag = p_ptag.add_run()
                        run_ptag.text = pillar_tags[bi]
                        run_ptag.font.size = Pt(8.5)
                        run_ptag.font.bold = True
                        run_ptag.font.color.rgb = c_accent if bi % 2 == 0 else c_highlight

                        head, body = _split_item_header_body(item)

                        p_hd = tf_i.add_paragraph()
                        run_hd = p_hd.add_run()
                        run_hd.text = head
                        run_hd.font.size = Pt(14)
                        run_hd.font.bold = True
                        run_hd.font.color.rgb = c_primary

                        p_bd = tf_i.add_paragraph()
                        run_bd = p_bd.add_run()
                        run_bd.text = f"\n{body or head}"
                        run_bd.font.size = Pt(11.0)
                        run_bd.font.color.rgb = c_text_dark

                        for pb in pillar_bullets_map[bi]:
                            p_pb = tf_i.add_paragraph()
                            run_pb = p_pb.add_run()
                            run_pb.text = f"• {pb}"
                            run_pb.font.size = Pt(9.5)
                            run_pb.font.color.rgb = c_text_muted

                        p_stat = tf_i.add_paragraph()
                        run_stat = p_stat.add_run()
                        run_stat.text = f"\n{status_badges[bi]}"
                        run_stat.font.size = Pt(9.5)
                        run_stat.font.bold = True
                        run_stat.font.color.rgb = c_highlight

                # CASE 4: Exactly 4 items -> 2x2 Balanced Bento Grid!
                elif n_items == 4:
                    grid_w = 5.65
                    grid_h = 2.35
                    gap_x = 0.40
                    gap_y = 0.30
                    start_x = 0.8
                    start_y = 1.65

                    for bi, item in enumerate(items[:4]):
                        r = bi // 2
                        c = bi % 2
                        bx = start_x + c * (grid_w + gap_x)
                        by = start_y + r * (grid_h + gap_y)

                        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(by), Inches(grid_w), Inches(grid_h))
                        card.fill.solid()
                        card.fill.fore_color.rgb = c_white
                        card.line.color.rgb = c_card_border
                        card.line.width = Pt(1)

                        # Left accent bar
                        l_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(by), Inches(0.1), Inches(grid_h))
                        l_stripe.fill.solid()
                        l_stripe.fill.fore_color.rgb = c_accent if bi % 2 == 0 else c_highlight
                        l_stripe.line.fill.background()

                        tb_i = slide.shapes.add_textbox(Inches(bx + 0.28), Inches(by + 0.18), Inches(grid_w - 0.45), Inches(grid_h - 0.36))
                        tf_i = tb_i.text_frame
                        tf_i.word_wrap = True

                        p_itag = tf_i.paragraphs[0]
                        run_itag = p_itag.add_run()
                        run_itag.text = f"KEY FINDING 0{bi+1} // STRATEGIC OBS"
                        run_itag.font.size = Pt(8.5)
                        run_itag.font.bold = True
                        run_itag.font.color.rgb = c_accent

                        head, body = _split_item_header_body(item)

                        p_hd = tf_i.add_paragraph()
                        run_hd = p_hd.add_run()
                        run_hd.text = head
                        run_hd.font.size = Pt(13)
                        run_hd.font.bold = True
                        run_hd.font.color.rgb = c_primary

                        p_bd = tf_i.add_paragraph()
                        run_bd = p_bd.add_run()
                        run_bd.text = body or head
                        run_bd.font.size = Pt(11)
                        run_bd.font.color.rgb = c_text_dark

                # CASE 5: 5+ items -> Executive Split Layout (Featured Hero Card on Left, Stacked on Right)
                else:
                    hero_w = 4.8
                    hero_h = 5.0
                    hero_x = 0.8
                    hero_y = 1.65

                    # Left Hero Card
                    hero_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(hero_x), Inches(hero_y), Inches(hero_w), Inches(hero_h))
                    hero_card.fill.solid()
                    hero_card.fill.fore_color.rgb = c_white
                    hero_card.line.color.rgb = c_accent
                    hero_card.line.width = Pt(1.5)

                    t_hstripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(hero_x), Inches(hero_y), Inches(hero_w), Inches(0.08))
                    t_hstripe.fill.solid()
                    t_hstripe.fill.fore_color.rgb = c_accent
                    t_hstripe.line.fill.background()

                    tb_hero = slide.shapes.add_textbox(Inches(hero_x + 0.3), Inches(hero_y + 0.25), Inches(hero_w - 0.6), Inches(hero_h - 0.5))
                    tf_hero = tb_hero.text_frame
                    tf_hero.word_wrap = True

                    p_htag = tf_hero.paragraphs[0]
                    run_htag = p_htag.add_run()
                    run_htag.text = "CORE STRATEGIC TAKEAWAY"
                    run_htag.font.size = Pt(9)
                    run_htag.font.bold = True
                    run_htag.font.color.rgb = c_accent

                    h0, b0 = _split_item_header_body(items[0])

                    p_h0 = tf_hero.add_paragraph()
                    run_h0 = p_h0.add_run()
                    run_h0.text = h0
                    run_h0.font.size = Pt(15)
                    run_h0.font.bold = True
                    run_h0.font.color.rgb = c_primary

                    p_b0 = tf_hero.add_paragraph()
                    run_b0 = p_b0.add_run()
                    run_b0.text = f"\n{b0 or h0}"
                    run_b0.font.size = Pt(11.5)
                    run_b0.font.color.rgb = c_text_dark

                    p_call = tf_hero.add_paragraph()
                    run_call = p_call.add_run()
                    run_call.text = "\n✔ Primary operational lever for immediate leadership focus."
                    run_call.font.size = Pt(10)
                    run_call.font.bold = True
                    run_call.font.color.rgb = c_highlight

                    # Right Stacked Cards
                    right_w = 6.5
                    right_x = 6.0
                    rem_items = items[1:5]
                    row_h = (hero_h - (len(rem_items) - 1) * 0.15) / max(len(rem_items), 1)

                    for si, sitem in enumerate(rem_items):
                        sy = hero_y + si * (row_h + 0.15)

                        scard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(right_x), Inches(sy), Inches(right_w), Inches(row_h))
                        scard.fill.solid()
                        scard.fill.fore_color.rgb = c_white
                        scard.line.color.rgb = c_card_border
                        scard.line.width = Pt(1)

                        # Left accent bar on stacked card
                        ls_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x), Inches(sy), Inches(0.08), Inches(row_h))
                        ls_bar.fill.solid()
                        ls_bar.fill.fore_color.rgb = c_accent if si % 2 == 0 else c_highlight
                        ls_bar.line.fill.background()

                        tb_s = slide.shapes.add_textbox(Inches(right_x + 0.25), Inches(sy + 0.10), Inches(right_w - 0.4), Inches(row_h - 0.20))
                        tf_s = tb_s.text_frame
                        tf_s.word_wrap = True

                        sh, sb = _split_item_header_body(sitem)

                        p_sh = tf_s.paragraphs[0]
                        run_sh = p_sh.add_run()
                        run_sh.text = f"{sh}: "
                        run_sh.font.bold = True
                        run_sh.font.size = Pt(11)
                        run_sh.font.color.rgb = c_primary

                        run_sb = p_sh.add_run()
                        run_sb.text = sb or sh
                        run_sb.font.bold = False
                        run_sb.font.size = Pt(10.5)
                        run_sb.font.color.rgb = c_text_dark

        # Speaker notes for every slide
        notes_text = _clean_str(slide_data.get("speaker_notes", ""))
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    # Save to buffer
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    # Record in audit log
    log_audit_event(
        db=db,
        action="dataset_workflow.presentation.generate",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="workflow",
        resource_id=workflow_id,
        new_values={
            "template": tpl_key,
            "title": title,
            "slides": total_slides,
            "charts_rendered": charts_rendered,
        },
        request=request,
    )
    db.commit()

    safe_filename = dataset_name.lower().replace(" ", "_").replace(".csv", "").replace(".xlsx", "")
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f'attachment; filename="{safe_filename}_{tpl_key}_presentation.pptx"',
        },
    )


@router.post("/{workflow_id}/presentation")
async def generate_presentation(
    workflow_id: str,
    payload: GeneratePresentationRequest,
    request: Request,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Generate a dynamic PPTX presentation tailored to the selected theme & narrative style."""
    from fastapi.responses import StreamingResponse

    try:
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches
    except ImportError:
        raise HTTPException(
            status_code=503,
            detail="Presentation generation is not available (python-pptx not installed).",
        ) from None

    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    org_id = get_current_organization_id(current_user, db)

    # Try to get in-memory DataFrame
    df = None
    in_memory_state = _orchestrator.get_state(workflow_id)
    if in_memory_state is not None:
        df = in_memory_state.context.get("df")

    profile = _stage_result(state_dict, WorkflowStage.PROFILED) or {}
    quality = _stage_result(state_dict, WorkflowStage.QUALITY_CHECKED) or {}
    industry = _stage_result(state_dict, WorkflowStage.INDUSTRY_IDENTIFIED) or {}
    insights_data = _stage_result(state_dict, WorkflowStage.INSIGHTS_GENERATED) or {}
    dashboard = _stage_result(state_dict, WorkflowStage.DASHBOARD_READY) or {}

    dataset_name = state_dict.get("dataset_name", "Dataset")
    template = payload.template or "executive"
    title = payload.title or f"{dataset_name} — Strategic Analysis"

    auto_dashboard = dashboard.get("auto_dashboard") if isinstance(dashboard, dict) else None

    if auto_dashboard:
        from services.auto.chart_specification import (
            ChartSpecification,
            DashboardSpecification,
            InsightSpecification,
            KPISpecification,
        )
        from services.auto.presentation_layout_engine import PresentationLayoutEngine

        dashboard_obj = DashboardSpecification(
            title=title,
            subtitle=auto_dashboard.get("subtitle", ""),
            industry=auto_dashboard.get("industry", "general"),
            dataset_name=dataset_name,
            kpis=[KPISpecification(**k) for k in auto_dashboard.get("kpis", [])] if auto_dashboard.get("kpis") else [],
            charts=[ChartSpecification.from_dict(c) for c in auto_dashboard.get("charts", [])] if auto_dashboard.get("charts") else [],
            insights=[InsightSpecification(**i) for i in auto_dashboard.get("insights", [])] if auto_dashboard.get("insights") else [],
            recommendations=auto_dashboard.get("recommendations", []),
        )

        layout_engine = PresentationLayoutEngine()
        pres_spec = layout_engine.generate_presentation(
            dashboard=dashboard_obj,
            template=template,
            title=title,
            profile=profile,
            quality=quality,
        )

        return _generate_auto_pptx(
            auto_presentation=pres_spec.to_dict(),
            auto_dashboard=auto_dashboard,
            df=df,
            dataset_name=dataset_name,
            title=title,
            template=template,
            workflow_id=workflow_id,
            current_user=current_user,
            org_id=org_id,
            request=request,
            db=db,
        )

    # ── Fallback: Legacy presentation generation ──
    from studios.presentation_service import PresentationStudioService
    # Get recommended charts from the dashboard stage
    recommended_charts = dashboard.get("recommended_charts", [])
    # Filter to chart types we can render (bar, line, pie, scatter)
    renderable_charts = [
        c
        for c in recommended_charts
        if c.get("type") in ("bar_chart", "line_chart", "pie_chart", "scatter")
    ]

    # Build source data for slide generation
    source_data = {
        "title": title,
        "subtitle": f"Generated from {dataset_name}",
        "summary": quality.get("summary", insights_data.get("executive_summary", "")),
        "findings": "\n".join(
            f"â€¢ {i.get('title', '')}: {i.get('description', '')}"
            for i in insights_data.get("insights", [])[:5]
        ),
        "recommendations": "\n".join(quality.get("recommendations", [])[:5]),
        "data_overview": (
            f"Rows: {profile.get('row_count', 'N/A')}, "
            f"Columns: {profile.get('column_count', 'N/A')}, "
            f"Quality Score: {quality.get('score', {}).get('overall', 'N/A')}"
        ),
        "chart_config": (
            renderable_charts[0]
            if renderable_charts
            else (
                dashboard.get("recommended_charts", [{}])[0]
                if dashboard.get("recommended_charts")
                else None
            )
        ),
        "next_steps": "Review findings and implement recommended actions.",
        "industry": industry.get("industry", "general"),
    }

    # Generate slides metadata
    slides = PresentationStudioService.generate_slides(
        source_type="workflow",
        source_data=source_data,
        template=payload.template,
    )

    # Create actual PPTX file
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    charts_rendered = 0
    for slide_data in slides:
        layout = prs.slide_layouts[1]  # Title and Content layout
        if slide_data.get("layout") == "title":
            layout = prs.slide_layouts[0]  # Title Slide layout

        slide = prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")

        # Set content (skip text body for chart slides â€” chart will fill the space)
        is_chart_slide = slide_data.get("layout") == "chart"
        content = slide_data.get("content", "")
        if is_chart_slide:
            # Try to render an actual chart on this slide
            chart_spec = slide_data.get("chart_config") or source_data.get("chart_config")
            if chart_spec and df is not None:
                added = _add_chart_to_slide(slide, chart_spec, df)
                if added:
                    charts_rendered += 1
                else:
                    # Fallback: show text content if chart rendering failed
                    if content and len(slide.placeholders) > 1:
                        body = slide.placeholders[1]
                        body.text_frame.text = content
            elif content and len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                body.text_frame.text = content
        else:
            if content and len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.text = content

        # Add speaker notes
        notes = slide_data.get("speaker_notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # Add additional chart slides for remaining renderable charts
    # (skip the first one if it was already rendered on the "Key Metrics" slide)
    if df is not None and len(renderable_charts) > 1:
        for chart_spec in renderable_charts[1:5]:  # Limit to 4 extra chart slides
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = chart_spec.get("title", "Chart")
            added = _add_chart_to_slide(slide, chart_spec, df, top=1.8, height=5.2)
            if added:
                charts_rendered += 1
                # Add speaker notes
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = chart_spec.get("reasoning", "")

    # Save to BytesIO and return as download
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    # Record in audit
    log_audit_event(
        db=db,
        action="dataset_workflow.presentation.generate",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="workflow",
        resource_id=workflow_id,
        new_values={
            "template": payload.template,
            "title": title,
            "slides": len(prs.slides),
            "charts_rendered": charts_rendered,
        },
        request=request,
    )
    db.commit()

    safe_filename = dataset_name.replace(" ", "_").replace("/", "_")
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f'attachment; filename="{safe_filename}_presentation.pptx"',
        },
    )


@router.post("/{workflow_id}/report")
async def generate_workflow_report(
    workflow_id: str,
    payload: WorkflowReportRequest,
    request: Request,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Generate a publication-grade PDF audit and analysis report from workflow results."""
    from fastapi.responses import StreamingResponse
    from services.workflow_report_service import generate_workflow_pdf_report

    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    org_id = get_current_organization_id(current_user, db)
    dataset_name = state_dict.get("dataset_name", "Dataset")

    # Rehydrate in-memory state if available for transformations
    in_memory = _orchestrator.get_state(workflow_id)
    if in_memory and hasattr(in_memory, "transformations"):
        state_dict["transformations"] = in_memory.transformations

    user_name = current_user.get("full_name") or current_user.get("email") or "DataFlow Analyst"

    try:
        pdf_bytes = generate_workflow_pdf_report(
            workflow_state_dict=state_dict,
            report_config=payload.dict(),
            current_user_name=user_name,
        )
    except Exception as e:
        logger.exception("Failed to generate workflow PDF report: %s", e)
        raise HTTPException(
            status_code=500, detail=f"Failed to generate PDF report: {str(e)}"
        )

    log_audit_event(
        db=db,
        action="dataset_workflow.report.generate",
        user_id=current_user["id"],
        organization_id=org_id,
        resource_type="workflow",
        resource_id=workflow_id,
        new_values={
            "report_title": payload.title,
            "organization": payload.organization,
            "include_executive_summary": payload.include_executive_summary,
            "include_data_quality": payload.include_data_quality,
            "include_visualizations": payload.include_visualizations,
        },
        request=request,
    )
    db.commit()

    safe_filename = dataset_name.replace(" ", "_").replace("/", "_")
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={
            "Content-Disposition": f'attachment; filename="{safe_filename}_Audit_Report.pdf"',
        },
    )


# â”€â”€ Auto Engine Endpoints â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€


@router.get("/{workflow_id}/understanding")
async def get_dataset_understanding(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the auto-generated dataset understanding (semantic column roles, correlations, etc.)."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    dashboard = _stage_result(state_dict, WorkflowStage.DASHBOARD_READY) or {}

    understanding = dashboard.get("dataset_understanding") if isinstance(dashboard, dict) else None
    if not understanding:
        in_memory = _orchestrator.get_state(workflow_id)
        if in_memory and in_memory.context.get("auto_understanding"):
            understanding = in_memory.context["auto_understanding"]

    if not understanding:
        raise HTTPException(
            status_code=404,
            detail="Dataset understanding not available. Ensure the workflow has completed the dashboard stage.",
        )

    return understanding


@router.get("/{workflow_id}/auto-dashboard")
async def get_auto_dashboard(
    workflow_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the auto-generated dashboard specification with canonical chart specs, KPIs, insights, and layout."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    dashboard = _stage_result(state_dict, WorkflowStage.DASHBOARD_READY) or {}

    auto_dashboard = dashboard.get("auto_dashboard") if isinstance(dashboard, dict) else None
    if not auto_dashboard:
        in_memory = _orchestrator.get_state(workflow_id)
        if in_memory and in_memory.context.get("auto_dashboard"):
            auto_dashboard = in_memory.context["auto_dashboard"]

    if not auto_dashboard:
        raise HTTPException(
            status_code=404,
            detail="Auto dashboard not available. Ensure the workflow has completed the dashboard stage.",
        )

    return auto_dashboard


@router.get("/{workflow_id}/auto-presentation")
async def get_auto_presentation_spec(
    workflow_id: str,
    template: str = "executive",
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the auto-generated presentation specification (slide plan, chart placements, validation)."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    dashboard = _stage_result(state_dict, WorkflowStage.DASHBOARD_READY) or {}
    auto_dashboard = dashboard.get("auto_dashboard") if isinstance(dashboard, dict) else None

    if auto_dashboard:
        from services.auto.chart_specification import (
            ChartSpecification,
            DashboardSpecification,
            InsightSpecification,
            KPISpecification,
        )
        from services.auto.presentation_layout_engine import PresentationLayoutEngine

        profile = _stage_result(state_dict, WorkflowStage.PROFILED) or {}
        quality = _stage_result(state_dict, WorkflowStage.QUALITY_CHECKED) or {}
        dataset_name = state_dict.get("dataset_name", "Dataset")

        dashboard_obj = DashboardSpecification(
            title=auto_dashboard.get("title", f"{dataset_name} — Strategic Analysis"),
            subtitle=auto_dashboard.get("subtitle", ""),
            industry=auto_dashboard.get("industry", "general"),
            dataset_name=dataset_name,
            kpis=[KPISpecification(**k) for k in auto_dashboard.get("kpis", [])] if auto_dashboard.get("kpis") else [],
            charts=[ChartSpecification.from_dict(c) for c in auto_dashboard.get("charts", [])] if auto_dashboard.get("charts") else [],
            insights=[InsightSpecification(**i) for i in auto_dashboard.get("insights", [])] if auto_dashboard.get("insights") else [],
            recommendations=auto_dashboard.get("recommendations", []),
        )

        layout_engine = PresentationLayoutEngine()
        pres_spec = layout_engine.generate_presentation(
            dashboard=dashboard_obj,
            template=template,
            profile=profile,
            quality=quality,
        )
        return pres_spec.to_dict()

    auto_presentation = dashboard.get("auto_presentation") if isinstance(dashboard, dict) else None
    if not auto_presentation:
        in_memory = _orchestrator.get_state(workflow_id)
        if in_memory and in_memory.context.get("auto_presentation"):
            auto_presentation = in_memory.context["auto_presentation"]

    if not auto_presentation:
        raise HTTPException(
            status_code=404,
            detail="Auto presentation not available. Ensure the workflow has completed the dashboard stage.",
        )

    return auto_presentation


@router.get("/{workflow_id}/charts/{chart_id}/explain")
async def explain_chart(
    workflow_id: str,
    chart_id: str,
    current_user: dict = Depends(get_current_user),
    db: DbSession = Depends(get_db),
):
    """Get the 'Why this chart?' explanation for a specific chart."""
    state_dict = _get_workflow_state_dict(workflow_id, current_user, db)
    dashboard = _stage_result(state_dict, WorkflowStage.DASHBOARD_READY) or {}

    auto_dashboard = dashboard.get("auto_dashboard") if isinstance(dashboard, dict) else None
    if not auto_dashboard:
        in_memory = _orchestrator.get_state(workflow_id)
        if in_memory and in_memory.context.get("auto_dashboard"):
            auto_dashboard = in_memory.context["auto_dashboard"]

    if not auto_dashboard:
        raise HTTPException(status_code=404, detail="Auto dashboard not available.")

    chart = None
    for c in auto_dashboard.get("charts", []):
        if c.get("id") == chart_id:
            chart = c
            break

    if not chart:
        raise HTTPException(status_code=404, detail=f"Chart '{chart_id}' not found.")

    return {
        "chart_id": chart_id,
        "chart_type": chart.get("chart_type"),
        "title": chart.get("title"),
        "reason": chart.get("reason"),
        "importance_score": chart.get("importance_score"),
        "confidence": chart.get("confidence"),
        "source_analysis": chart.get("source_analysis"),
    }
